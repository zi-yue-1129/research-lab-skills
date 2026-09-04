"""Tests that a standalone textbox is sized by measurement, not by guess.

`_estimate_text_width` multiplied the longest line's character count by
0.65 em. Every glyph is not 0.65 em wide, so a box could come out narrower
than the text it holds; with wrapping off, LibreOffice then laid the run out
against a box that disagreed with the SVG, which is what the ragged bullet
list in the final verification was.
"""
import shutil
import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parents[2]))

import pytest
from lxml import etree
from pptx import Presentation
from pptx.util import Emu

from fonts import text_width
from svg_to_pptx.converter import CoordSystem
from svg_to_pptx.style_parser import compute_style
from svg_to_pptx.text_converter import _PAD_SVG, add_textbox

CS = CoordSystem(svg_w=1200.0, svg_h=675.0)
_FAMILY = "DejaVu Sans"

_NEEDS_FONTCONFIG = pytest.mark.skipif(
    shutil.which("fc-match") is None,
    # fontconfig is an optional external binary; there is no font file to
    # measure a string with, and no family to resolve, without it.
    reason="fontconfig (fc-match) is not installed",
)


def _blank_slide():
    """Return an empty 16:9 slide and its presentation."""
    prs = Presentation()
    prs.slide_width = Emu(12_192_000)
    prs.slide_height = Emu(6_858_000)
    return prs.slides.add_slide(prs.slide_layouts[6]), prs


def _box(markup: str):
    """Convert one `<text>` element to a standalone textbox.

    Args:
        markup: A complete `<text>` element.

    Returns:
        The created textbox shape.
    """
    slide, _ = _blank_slide()
    elem = etree.fromstring(markup)
    return add_textbox(slide, elem, compute_style(elem, {}), CS)


@_NEEDS_FONTCONFIG
def test_the_width_comes_from_the_glyphs_not_the_character_count() -> None:
    """Six wide glyphs make a wider box than six narrow ones.

    At 0.65 em per character the two are identical, which is the whole defect:
    the box was the same width whatever the text said.
    """
    narrow = _box('<text x="100" y="100" font-size="21" '
                  'font-family="DejaVu Sans">llllll</text>')
    wide = _box('<text x="100" y="100" font-size="21" '
                'font-family="DejaVu Sans">WWWWWW</text>')
    assert wide.width > narrow.width


@_NEEDS_FONTCONFIG
def test_the_width_is_the_measured_advance_plus_padding() -> None:
    """The box is exactly the measured run plus the documented padding."""
    box = _box('<text x="100" y="100" font-size="21" '
               'font-family="DejaVu Sans">One token set now describes '
               'each deck.</text>')
    measured = text_width("One token set now describes each deck.",
                          _FAMILY, 21.0, 400)
    assert box.width == Emu(CS.x(measured + _PAD_SVG))


@_NEEDS_FONTCONFIG
def test_a_bold_run_is_measured_on_the_bold_face() -> None:
    """A bold face is materially wider than its regular sibling.

    Measuring a heading with the regular face is how a title silently
    overflows its own box.
    """
    regular = _box('<text x="100" y="100" font-size="32" font-weight="400" '
                   'font-family="DejaVu Sans">Before and after</text>')
    bold = _box('<text x="100" y="100" font-size="32" font-weight="700" '
                'font-family="DejaVu Sans">Before and after</text>')
    assert bold.width > regular.width


@_NEEDS_FONTCONFIG
def test_the_widest_line_sizes_the_box() -> None:
    """A multi-line label is as wide as its widest line, not its first."""
    box = _box('<text x="100" y="100" font-size="21" '
               'font-family="DejaVu Sans">'
               '<tspan x="100" dy="0">short</tspan>'
               '<tspan x="100" dy="30">a considerably longer second line</tspan>'
               '</text>')
    measured = text_width("a considerably longer second line",
                          _FAMILY, 21.0, 400)
    assert box.width == Emu(CS.x(measured + _PAD_SVG))


@_NEEDS_FONTCONFIG
def test_a_right_anchored_box_still_ends_at_its_svg_x() -> None:
    """Sizing by measurement must not move the anchor.

    A deck footer is anchored at its right edge; the box grows leftward.
    """
    box = _box('<text x="1152" y="636" font-size="12" text-anchor="end" '
               'font-family="DejaVu Sans">report-slides verification</text>')
    assert box.left + box.width == Emu(CS.x(1152.0))


def test_an_undeclared_family_still_produces_a_box() -> None:
    """An SVG that names no font is converted rather than refused.

    Nothing states which face such a document would be rendered in, so there
    is nothing to measure; the box falls back to the character estimate.
    """
    box = _box('<text x="100" y="100" font-size="21">no family here</text>')
    assert box.width > 0
