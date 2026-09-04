import pytest
from pptx import Presentation
from pptx.util import Emu

from svg_to_pptx.converter import CoordSystem

def test_coordsystem_from_viewbox():
    cs = CoordSystem.from_viewbox("0 0 1200 675")
    assert cs.svg_w == 1200.0
    assert cs.svg_h == 675.0

def test_coordsystem_x_scaling():
    cs = CoordSystem.from_viewbox("0 0 1200 675")
    assert cs.x(600) == 6_096_000

def test_coordsystem_y_scaling():
    cs = CoordSystem.from_viewbox("0 0 1200 675")
    assert cs.y(337.5) == 3_429_000

def test_coordsystem_missing_viewbox():
    cs = CoordSystem.from_viewbox("")
    assert cs.svg_w == 1200.0
    assert cs.svg_h == 675.0

def test_coordsystem_malformed_viewbox():
    cs = CoordSystem.from_viewbox("0 0 abc 675")
    assert cs.svg_w == 1200.0
    assert cs.svg_h == 675.0


from svg_to_pptx.style_parser import (
    parse_inline_style, resolve_color, compute_style, apply_fill, apply_stroke
)
from pptx.dml.color import RGBColor
from lxml import etree


def test_parse_inline_style():
    r = parse_inline_style("fill:#3b82f6;stroke:#1d4ed8;stroke-width:2")
    assert r == {"fill": "#3b82f6", "stroke": "#1d4ed8", "stroke-width": "2"}


def test_parse_inline_style_with_spaces():
    r = parse_inline_style(" fill : #fff ; stroke: none ")
    assert r["fill"] == "#fff"
    assert r["stroke"] == "none"


def test_resolve_color_hex6():
    assert resolve_color("#3b82f6") == RGBColor(0x3b, 0x82, 0xf6)


def test_resolve_color_hex3():
    assert resolve_color("#f0f") == RGBColor(0xff, 0x00, 0xff)


def test_resolve_color_named_white():
    assert resolve_color("white") == RGBColor(0xff, 0xff, 0xff)


def test_resolve_color_named_blue():
    assert resolve_color("blue") == RGBColor(0x00, 0x00, 0xff)


def test_resolve_color_none():
    assert resolve_color("none") is None


def test_resolve_color_invalid():
    assert resolve_color("notacolor") is None


def test_compute_style_attr_fallback():
    elem = etree.fromstring('<rect fill="#ff0000"/>')
    result = compute_style(elem, {})
    assert result["fill"] == "#ff0000"


def test_compute_style_inherited():
    elem = etree.fromstring('<rect/>')
    result = compute_style(elem, {"fill": "#00ff00"})
    assert result["fill"] == "#00ff00"


def test_compute_style_inline_overrides_attr():
    elem = etree.fromstring('<rect fill="#ff0000" style="fill:#00ff00"/>')
    result = compute_style(elem, {})
    assert result["fill"] == "#00ff00"


def test_compute_style_inline_overrides_inherited():
    elem = etree.fromstring('<rect style="stroke:none"/>')
    result = compute_style(elem, {"stroke": "#000000"})
    assert result["stroke"] == "none"


from svg_to_pptx.style_parser import (
    parse_transform, apply_transform_to_pos, apply_gradient_fill
)
from pptx import Presentation
from pptx.util import Emu


def test_parse_transform_translate():
    tx, ty, rot, sx, sy = parse_transform("translate(100,50)")
    assert tx == pytest.approx(100.0)
    assert ty == pytest.approx(50.0)
    assert rot == pytest.approx(0.0)


def test_parse_transform_rotate():
    tx, ty, rot, sx, sy = parse_transform("rotate(45)")
    assert rot == pytest.approx(45.0)


def test_parse_transform_scale():
    tx, ty, rot, sx, sy = parse_transform("scale(2,3)")
    assert sx == pytest.approx(2.0)
    assert sy == pytest.approx(3.0)


def test_parse_transform_matrix():
    tx, ty, rot, sx, sy = parse_transform("matrix(1,0,0,1,50,30)")
    assert tx == pytest.approx(50.0)
    assert ty == pytest.approx(30.0)


def test_apply_transform_to_pos():
    x, y, w, h = apply_transform_to_pos(100, 200, 300, 400, "translate(50,25)")
    assert x == pytest.approx(150.0)
    assert y == pytest.approx(225.0)
    assert w == pytest.approx(300.0)


def test_apply_gradient_fill_linear():
    prs = Presentation()
    prs.slide_width = Emu(12_192_000)
    prs.slide_height = Emu(6_858_000)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shape = slide.shapes.add_shape(1, Emu(0), Emu(0), Emu(1000), Emu(500))
    stops = [("0", "#ff0000"), ("100%", "#0000ff")]
    apply_gradient_fill(shape, stops, angle_deg=0.0)
    A = "http://schemas.openxmlformats.org/drawingml/2006/main"
    from pptx.oxml.ns import qn
    spPr = shape._element.find(qn("p:spPr"))
    assert spPr.find(f"{{{A}}}gradFill") is not None


_A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"


def _blank_slide():
    """Return a blank 16:9 slide and its presentation.

    Returns:
        A `(slide, presentation)` pair on the blank layout.
    """
    prs = Presentation()
    prs.slide_width = Emu(12_192_000)
    prs.slide_height = Emu(6_858_000)
    return prs.slides.add_slide(prs.slide_layouts[6]), prs


def test_dash_style_for_token_patterns():
    """The token dash patterns map to the intended OOXML presets."""
    from svg_to_pptx.style_parser import dash_style_for

    assert dash_style_for("", 2) is None
    assert dash_style_for("none", 2) is None
    assert dash_style_for("2 4", 2) == "sysDot"
    assert dash_style_for("8 4", 2) == "dash"
    assert dash_style_for("40 8", 2) == "lgDash"
    assert dash_style_for("8 4 2 4", 2) == "dashDot"
    assert dash_style_for("20 4 2 4 2 4", 2) == "lgDashDotDot"


def test_dash_style_for_accepts_comma_separated_values():
    """SVG permits commas between dash values, as the timeline renderer uses."""
    from svg_to_pptx.style_parser import dash_style_for

    assert dash_style_for("3,2", 1.5) == "sysDot"


def test_dash_style_for_rejects_malformed_array():
    """A non-numeric dasharray is an error, not a silent solid line."""
    from svg_to_pptx.style_parser import dash_style_for

    with pytest.raises(ValueError, match="dasharray"):
        dash_style_for("8 wide", 2)


def test_apply_dash_sets_prst_dash_in_schema_order():
    """apply_dash writes a:prstDash before any line-end child.

    The two helpers are called here in the reverse of the production order to
    prove the ordering comes from `ensure_ln_child`'s schema-position insert
    and not from the sequence of calls.
    """
    from svg_to_pptx.style_parser import apply_dash, apply_line_ends

    slide, _ = _blank_slide()
    conn = slide.shapes.add_connector(1, Emu(0), Emu(0), Emu(100000), Emu(0))
    style = {"stroke": "#475569", "stroke-width": "2",
             "stroke-dasharray": "8 4", "marker-end": "url(#a)"}
    apply_line_ends(conn, style)
    apply_dash(conn, style)
    ln = conn._element.find(f".//{_A}ln")
    tags = [child.tag.split("}")[-1] for child in ln]
    assert "prstDash" in tags
    assert tags.index("prstDash") < tags.index("tailEnd")
    assert ln.find(f"{_A}prstDash").get("val") == "dash"


def test_apply_dash_leaves_solid_lines_alone():
    """No dasharray means no prstDash element at all."""
    from svg_to_pptx.style_parser import apply_dash

    slide, _ = _blank_slide()
    conn = slide.shapes.add_connector(1, Emu(0), Emu(0), Emu(100000), Emu(0))
    apply_dash(conn, {"stroke": "#475569", "stroke-width": "2"})
    ln = conn._element.find(f".//{_A}ln")
    assert ln is None or ln.find(f"{_A}prstDash") is None


def test_apply_alpha_sets_fill_transparency():
    """fill-opacity becomes an a:alpha child of the fill colour."""
    from svg_to_pptx.style_parser import apply_alpha, apply_fill

    slide, _ = _blank_slide()
    shape = slide.shapes.add_shape(1, Emu(0), Emu(0), Emu(100000), Emu(100000))
    style = {"fill": "#3b82f6", "fill-opacity": "0.25"}
    apply_fill(shape, style["fill"])
    apply_alpha(shape, style)
    alpha = shape._element.find(f".//{_A}solidFill/{_A}srgbClr/{_A}alpha")
    assert alpha is not None
    assert alpha.get("val") == "25000"


def test_apply_alpha_multiplies_opacity_and_fill_opacity():
    """opacity and fill-opacity compose multiplicatively, as SVG specifies."""
    from svg_to_pptx.style_parser import apply_alpha, apply_fill

    slide, _ = _blank_slide()
    shape = slide.shapes.add_shape(1, Emu(0), Emu(0), Emu(100000), Emu(100000))
    apply_fill(shape, "#3b82f6")
    apply_alpha(
        shape, {"fill": "#3b82f6", "opacity": "0.5", "fill-opacity": "0.5"})
    alpha = shape._element.find(f".//{_A}solidFill/{_A}srgbClr/{_A}alpha")
    assert alpha.get("val") == "25000"


def test_apply_alpha_is_a_no_op_at_full_opacity():
    """An opaque element gains no alpha element."""
    from svg_to_pptx.style_parser import apply_alpha, apply_fill

    slide, _ = _blank_slide()
    shape = slide.shapes.add_shape(1, Emu(0), Emu(0), Emu(100000), Emu(100000))
    apply_fill(shape, "#3b82f6")
    apply_alpha(shape, {"fill": "#3b82f6"})
    assert shape._element.find(f".//{_A}solidFill/{_A}srgbClr/{_A}alpha") is None


def test_apply_alpha_is_idempotent():
    """Applying twice leaves one alpha element, not two."""
    from svg_to_pptx.style_parser import apply_alpha, apply_fill

    slide, _ = _blank_slide()
    shape = slide.shapes.add_shape(1, Emu(0), Emu(0), Emu(100000), Emu(100000))
    apply_fill(shape, "#3b82f6")
    style = {"fill": "#3b82f6", "fill-opacity": "0.4"}
    apply_alpha(shape, style)
    apply_alpha(shape, style)
    color = shape._element.find(f".//{_A}solidFill/{_A}srgbClr")
    assert len(color.findall(f"{_A}alpha")) == 1
    assert color.find(f"{_A}alpha").get("val") == "40000"


def test_apply_alpha_sets_stroke_transparency_independently():
    """stroke-opacity affects the line colour, not the fill."""
    from svg_to_pptx.style_parser import apply_alpha, apply_fill, apply_stroke

    slide, _ = _blank_slide()
    shape = slide.shapes.add_shape(1, Emu(0), Emu(0), Emu(100000), Emu(100000))
    apply_fill(shape, "#3b82f6")
    apply_stroke(shape, {"stroke": "#1d4ed8", "stroke-width": "2"})
    apply_alpha(shape, {"fill": "#3b82f6", "stroke-opacity": "0.5"})
    line_alpha = shape._element.find(
        f".//{_A}ln/{_A}solidFill/{_A}srgbClr/{_A}alpha")
    assert line_alpha is not None
    assert line_alpha.get("val") == "50000"


def test_apply_alpha_rejects_out_of_range_opacity():
    """An opacity outside 0..1 is an error rather than being clamped silently."""
    from svg_to_pptx.style_parser import apply_alpha, apply_fill

    slide, _ = _blank_slide()
    shape = slide.shapes.add_shape(1, Emu(0), Emu(0), Emu(100000), Emu(100000))
    apply_fill(shape, "#3b82f6")
    with pytest.raises(ValueError, match="fill-opacity"):
        apply_alpha(shape, {"fill": "#3b82f6", "fill-opacity": "1.5"})


def test_apply_alpha_rejects_a_malformed_opacity():
    """A non-numeric opacity is an error, not silently ignored."""
    from svg_to_pptx.style_parser import apply_alpha, apply_fill

    slide, _ = _blank_slide()
    shape = slide.shapes.add_shape(1, Emu(0), Emu(0), Emu(100000), Emu(100000))
    apply_fill(shape, "#3b82f6")
    with pytest.raises(ValueError, match="opacity"):
        apply_alpha(shape, {"fill": "#3b82f6", "opacity": "half"})


def test_a_dashed_rect_exports_dashed():
    """The dispatcher wires apply_dash, not just the helper in isolation.

    The helpers existing but never being called is exactly the defect this
    task removes: `stroke-dasharray` was collected into every style dict and
    no code path applied it.
    """
    from lxml import etree as _etree

    from svg_to_pptx.converter import CoordSystem as _CS
    from svg_to_pptx.shapes import dispatch_shape
    from svg_to_pptx.style_parser import compute_style

    slide, _ = _blank_slide()
    elem = _etree.fromstring(
        '<rect x="0" y="0" width="200" height="100" fill="#fff" '
        'stroke="#475569" stroke-width="2" stroke-dasharray="8 4"/>')
    dispatch_shape(slide, elem, compute_style(elem, {}),
                   _CS(svg_w=1200.0, svg_h=675.0), None)
    prst = slide.shapes[0]._element.find(f".//{_A}ln/{_A}prstDash")
    assert prst is not None
    assert prst.get("val") == "dash"


def test_a_translucent_circle_exports_translucent():
    """The dispatcher wires apply_alpha for ovals as well as rects."""
    from lxml import etree as _etree

    from svg_to_pptx.converter import CoordSystem as _CS
    from svg_to_pptx.shapes import dispatch_shape
    from svg_to_pptx.style_parser import compute_style

    slide, _ = _blank_slide()
    elem = _etree.fromstring(
        '<circle cx="100" cy="100" r="40" fill="#10b981" fill-opacity="0.3"/>')
    dispatch_shape(slide, elem, compute_style(elem, {}),
                   _CS(svg_w=1200.0, svg_h=675.0), None)
    alpha = slide.shapes[0]._element.find(
        f".//{_A}solidFill/{_A}srgbClr/{_A}alpha")
    assert alpha is not None
    assert alpha.get("val") == "30000"


def test_a_dashed_connector_exports_dashed():
    """The connector path wires apply_dash too.

    The timeline renderer draws every event stem with `stroke-dasharray="3,2"`,
    so this is the pattern the deterministic renderer actually emits.
    """
    from lxml import etree as _etree

    from svg_to_pptx.connector import dispatch_connector
    from svg_to_pptx.converter import CoordSystem as _CS
    from svg_to_pptx.style_parser import compute_style

    slide, _ = _blank_slide()
    elem = _etree.fromstring(
        '<line x1="0" y1="0" x2="0" y2="24" stroke="#1e3a5f" '
        'stroke-width="1.5" stroke-dasharray="3,2"/>')
    conns = dispatch_connector(slide, elem, compute_style(elem, {}),
                               _CS(svg_w=1200.0, svg_h=675.0))
    prst = conns[0]._element.find(f".//{_A}ln/{_A}prstDash")
    assert prst is not None
    assert prst.get("val") == "sysDot"
