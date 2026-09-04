import pytest
from lxml import etree
from pptx import Presentation
from pptx.util import Emu
from pptx.oxml.ns import qn

from svg_to_pptx.converter import CoordSystem
from svg_to_pptx.connector import dispatch_connector, build_anchor_map, bind_connector_end, _add_line
from svg_to_pptx.style_parser import compute_style

CS = CoordSystem(svg_w=1200.0, svg_h=675.0)


def _slide():
    prs = Presentation()
    prs.slide_width = Emu(12_192_000)
    prs.slide_height = Emu(6_858_000)
    return prs.slides.add_slide(prs.slide_layouts[6])


def test_line_creates_connector():
    slide = _slide()
    elem = etree.fromstring('<line x1="100" y1="100" x2="400" y2="100" stroke="#000"/>')
    dispatch_connector(slide, elem, compute_style(elem, {}), CS)
    assert len(slide.shapes) == 1


def test_polyline_creates_n_minus_1_connectors():
    slide = _slide()
    elem = etree.fromstring('<polyline points="0,0 100,0 200,100" stroke="#000"/>')
    dispatch_connector(slide, elem, compute_style(elem, {}), CS)
    assert len(slide.shapes) == 2


def test_polygon_closes_path():
    slide = _slide()
    elem = etree.fromstring('<polygon points="0,0 100,0 50,100" stroke="#000"/>')
    dispatch_connector(slide, elem, compute_style(elem, {}), CS)
    assert len(slide.shapes) == 3


def test_build_anchor_map_rect():
    anchors = build_anchor_map([(None, 100.0, 50.0, 200.0, 100.0, 42)])
    assert 42 in anchors
    assert len(anchors[42]) == 8


def test_build_anchor_map_nearest():
    anchors = build_anchor_map([(None, 100.0, 50.0, 200.0, 100.0, 42)])
    top_center = next(a for a in anchors[42] if a[2] == 1)
    assert top_center[0] == 200.0
    assert top_center[1] == 50.0


def test_bind_connector_injects_xml():
    slide = _slide()
    conn = _add_line(slide, 0, 0, 100, 0, {})
    shape = slide.shapes.add_shape(1, Emu(0), Emu(0), Emu(CS.x(200)), Emu(CS.y(100)))
    bind_connector_end(conn, True, shape.shape_id, 3)
    cxnSp = conn._element
    nvCxnSpPr = cxnSp.find(qn('p:nvCxnSpPr'))
    cNvCxnSpPr = nvCxnSpPr.find(qn('p:cNvCxnSpPr'))
    stCxn = cNvCxnSpPr.find(qn('a:stCxn'))
    assert stCxn is not None
    assert stCxn.get('id') == str(shape.shape_id)
    assert stCxn.get('idx') == '3'


_A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"


def _line_ends(conn):
    """Return the (headEnd, tailEnd) elements of a connector's line properties.

    Args:
        conn: A PPTX connector shape.

    Returns:
        A pair of elements, each None when that end carries no arrowhead.
    """
    ln = conn._element.find(f".//{_A}ln")
    if ln is None:
        return None, None
    return ln.find(f"{_A}headEnd"), ln.find(f"{_A}tailEnd")


def _dispatch(markup: str):
    """Dispatch one SVG connector element and return its connectors.

    Args:
        markup: A single SVG element as a string.

    Returns:
        The list of connectors `dispatch_connector` created.
    """
    slide = _slide()
    elem = etree.fromstring(markup)
    return dispatch_connector(slide, elem, compute_style(elem, {}), CS)


def test_plain_line_has_no_arrowheads():
    """A line without marker attributes exports without arrowheads."""
    conns = _dispatch(
        '<line x1="0" y1="0" x2="100" y2="0" stroke="#475569" stroke-width="2"/>')
    head, tail = _line_ends(conns[0])
    assert head is None
    assert tail is None


def test_marker_end_produces_a_tail_arrowhead():
    """marker-end becomes an OOXML tailEnd triangle."""
    conns = _dispatch(
        '<line x1="0" y1="0" x2="100" y2="0" stroke="#475569" '
        'stroke-width="2" marker-end="url(#arrow)"/>')
    head, tail = _line_ends(conns[0])
    assert head is None
    assert tail is not None
    assert tail.get("type") == "triangle"


def test_marker_start_produces_a_head_arrowhead():
    """marker-start becomes an OOXML headEnd."""
    conns = _dispatch(
        '<line x1="0" y1="0" x2="100" y2="0" stroke="#475569" '
        'stroke-width="2" marker-start="url(#arrow)"/>')
    head, tail = _line_ends(conns[0])
    assert head is not None
    assert head.get("type") == "triangle"
    assert tail is None


def test_marker_none_is_not_an_arrowhead():
    """An explicit marker-end="none" adds nothing."""
    conns = _dispatch(
        '<line x1="0" y1="0" x2="100" y2="0" stroke="#475569" '
        'stroke-width="2" marker-end="none"/>')
    _, tail = _line_ends(conns[0])
    assert tail is None


def test_explicit_arrowhead_type_is_honoured():
    """data-pptx-arrowhead on the element overrides the triangle default."""
    conns = _dispatch(
        '<line x1="0" y1="0" x2="100" y2="0" stroke="#475569" stroke-width="2" '
        'marker-end="url(#a)" data-pptx-arrowhead="stealth"/>')
    _, tail = _line_ends(conns[0])
    assert tail.get("type") == "stealth"


def test_unknown_arrowhead_type_raises():
    """An unrecognised arrowhead name is an error, not a silent default."""
    with pytest.raises(ValueError, match="harpoon"):
        _dispatch(
            '<line x1="0" y1="0" x2="100" y2="0" stroke="#475569" '
            'stroke-width="2" marker-end="url(#a)" '
            'data-pptx-arrowhead="harpoon"/>')


def test_unknown_arrowhead_size_raises():
    """An unrecognised arrowhead size is an error too."""
    with pytest.raises(ValueError, match="enormous"):
        _dispatch(
            '<line x1="0" y1="0" x2="100" y2="0" stroke="#475569" '
            'stroke-width="2" marker-end="url(#a)" '
            'data-pptx-arrowhead-size="enormous"/>')


def test_a_polyline_carries_arrowheads_only_at_its_ends():
    """An elbow connector has two ends, not one per bend.

    `dispatch_connector` turns a polyline into one connector per segment. If the
    arrowhead is applied per segment, a four-point elbow -- the standard shape
    for routing around a node in an architecture diagram -- exports with three
    arrowheads pointing into the middle of itself.
    """
    conns = _dispatch(
        '<polyline points="0,0 100,0 100,80 200,80" stroke="#475569" '
        'stroke-width="2" marker-end="url(#a)"/>')
    assert len(conns) == 3
    tails = [_line_ends(conn)[1] for conn in conns]
    assert [tail is not None for tail in tails] == [False, False, True]


def test_a_polyline_head_is_on_its_first_segment_only():
    """marker-start lands on the first segment, not on every one."""
    conns = _dispatch(
        '<polyline points="0,0 100,0 100,80 200,80" stroke="#475569" '
        'stroke-width="2" marker-start="url(#a)"/>')
    heads = [_line_ends(conn)[0] for conn in conns]
    assert [head is not None for head in heads] == [True, False, False]


def test_a_closed_polygon_has_no_arrowheads():
    """A polygon has no ends, so a marker on it declares nothing to apply."""
    conns = _dispatch(
        '<polygon points="0,0 100,0 50,80" stroke="#475569" '
        'stroke-width="2" marker-end="url(#a)"/>')
    assert all(_line_ends(conn) == (None, None) for conn in conns)


def test_the_two_ends_can_differ():
    """A connector may start plain and end in a stealth arrow."""
    conns = _dispatch(
        '<line x1="0" y1="0" x2="100" y2="0" stroke="#475569" stroke-width="2" '
        'marker-start="url(#a)" marker-end="url(#a)" '
        'data-pptx-arrowhead-start="oval" data-pptx-arrowhead-end="stealth"/>')
    head, tail = _line_ends(conns[0])
    assert head.get("type") == "oval"
    assert tail.get("type") == "stealth"


def test_ensure_ln_child_keeps_schema_order():
    """Line-property children are inserted in DrawingML's required order.

    CT_LineProperties fixes the sequence fill, dash, join, headEnd, tailEnd.
    Appending blindly yields a file PowerPoint often tolerates but LibreOffice
    -- the renderer this skill's own review gate uses -- may reject.
    """
    from svg_to_pptx.style_parser import ensure_ln_child

    slide = _slide()
    conn = slide.shapes.add_connector(1, Emu(0), Emu(0), Emu(100000), Emu(0))
    ln = conn.line._get_or_add_ln()
    # Deliberately request them out of order.
    ensure_ln_child(ln, "a:tailEnd")
    ensure_ln_child(ln, "a:prstDash")
    ensure_ln_child(ln, "a:headEnd")
    ensure_ln_child(ln, "a:solidFill")
    tags = [child.tag.split("}")[-1] for child in ln]
    assert tags == ["solidFill", "prstDash", "headEnd", "tailEnd"]


def test_ensure_ln_child_is_idempotent():
    """Asking twice returns the same element rather than duplicating it."""
    from svg_to_pptx.style_parser import ensure_ln_child

    slide = _slide()
    conn = slide.shapes.add_connector(1, Emu(0), Emu(0), Emu(100000), Emu(0))
    ln = conn.line._get_or_add_ln()
    first = ensure_ln_child(ln, "a:tailEnd")
    second = ensure_ln_child(ln, "a:tailEnd")
    assert first is second
    assert len(ln.findall(f"{_A}tailEnd")) == 1


def test_ensure_ln_child_rejects_an_unknown_tag():
    """A tag outside the line-property schema is an error."""
    from svg_to_pptx.style_parser import ensure_ln_child

    slide = _slide()
    conn = slide.shapes.add_connector(1, Emu(0), Emu(0), Emu(100000), Emu(0))
    with pytest.raises(ValueError, match="a:nonsense"):
        ensure_ln_child(conn.line._get_or_add_ln(), "a:nonsense")
