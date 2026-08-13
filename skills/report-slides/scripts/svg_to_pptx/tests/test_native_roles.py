import json
from pathlib import Path

import pytest
from pptx import Presentation

from svg_to_pptx.converter import NativeObjectMarkerError, SvgConverter


def _write_svg(tmp_path: Path, body: str) -> Path:
    svg_path = tmp_path / "slide01_table.svg"
    svg_path.write_text(
        f'<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 1200 675">'
        f'{body}</svg>',
        encoding="utf-8",
    )
    return svg_path


def test_table_role_creates_native_table_and_skips_children(tmp_path: Path):
    (tmp_path / "table_data.json").write_text(json.dumps({
        "columns": ["Model", "QWK"],
        "rows": [["GPT-4", "0.671"]],
    }), encoding="utf-8")
    svg_path = _write_svg(tmp_path, (
        '<g data-pptx-role="table" data-pptx-source="table_data.json" '
        'data-pptx-bbox="60,75,1080,100">'
        '<rect x="60" y="75" width="1080" height="50" fill="#1e3a5f"/>'
        '<text x="600" y="100">Model</text>'
        '</g>'
    ))
    prs = Presentation()
    layout = prs.slide_layouts[6]
    conv = SvgConverter(str(svg_path))
    slide = conv.convert(prs, layout)

    table_frames = [s for s in slide.shapes if getattr(s, "has_table", False)]
    assert len(table_frames) == 1
    assert table_frames[0].table.cell(0, 0).text_frame.text == "Model"
    # The hand-drawn preview rect/text inside the marker must not also be
    # converted into a plain autoshape/textbox.
    assert len(slide.shapes) == 1


def test_chart_role_creates_native_chart(tmp_path: Path):
    (tmp_path / "chart_data.json").write_text(json.dumps({
        "type": "bar_chart",
        "categories": ["Q1", "Q2"],
        "series": [{"label": "EN", "color": "#1e3a5f", "values": [0.8, 0.9]}],
        "y_max": 1.0,
    }), encoding="utf-8")
    svg_path = _write_svg(tmp_path, (
        '<g data-pptx-role="chart" data-pptx-source="chart_data.json" '
        'data-pptx-bbox="130,100,970,420">'
        '<rect x="130" y="300" width="100" height="120" fill="#1e3a5f"/>'
        '</g>'
    ))
    prs = Presentation()
    layout = prs.slide_layouts[6]
    conv = SvgConverter(str(svg_path))
    slide = conv.convert(prs, layout)

    chart_frames = [s for s in slide.shapes if getattr(s, "has_chart", False)]
    assert len(chart_frames) == 1


def test_source_fragment_resolves_slide_data_json_by_index(tmp_path: Path):
    (tmp_path / "slide_data.json").write_text(json.dumps({
        "meta": {},
        "slides": [
            {"index": 0, "type": "bullet_list"},
            {"index": 4, "type": "table", "columns": ["A"], "rows": [["1"]]},
        ],
    }), encoding="utf-8")
    svg_path = _write_svg(tmp_path, (
        '<g data-pptx-role="table" data-pptx-source="slide_data.json#4" '
        'data-pptx-bbox="60,75,300,100">'
        '<rect x="60" y="75" width="300" height="50" fill="#1e3a5f"/>'
        '</g>'
    ))
    prs = Presentation()
    layout = prs.slide_layouts[6]
    conv = SvgConverter(str(svg_path))
    slide = conv.convert(prs, layout)

    table_frames = [s for s in slide.shapes if getattr(s, "has_table", False)]
    # cell(0, 0) is the header row (column name "A"); cell(1, 0) is the data
    # row that proves the correct slide entry (index 4, not index 0) was
    # resolved from slide_data.json.
    assert table_frames[0].table.cell(1, 0).text_frame.text == "1"


def test_group_role_wraps_child_shapes_in_native_group(tmp_path: Path):
    # The text sits below the rect's bbox (rect spans y=100..180) rather than
    # inside it: the pre-existing _compute_text_attachments pass embeds any
    # text whose baseline falls inside a sibling shape's bbox as that shape's
    # own label (a single autoshape, by design, for labeled-box content), so
    # an overlapping text here would collapse to 1 collected shape and never
    # exercise the 2-shapes-grouped path this test targets.
    svg_path = _write_svg(tmp_path, (
        '<g data-pptx-role="group" data-node-id="node1">'
        '<rect x="100" y="100" width="200" height="80" fill="#f8fafc"/>'
        '<text x="200" y="200">Encoder</text>'
        '</g>'
    ))
    prs = Presentation()
    layout = prs.slide_layouts[6]
    conv = SvgConverter(str(svg_path))
    slide = conv.convert(prs, layout)

    from pptx.enum.shapes import MSO_SHAPE_TYPE
    groups = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.GROUP]
    assert len(groups) == 1
    assert len(groups[0].shapes) == 2


def test_missing_source_and_bbox_raises_native_object_marker_error(tmp_path: Path):
    svg_path = _write_svg(tmp_path, (
        '<g data-pptx-role="table">'
        '<rect x="60" y="75" width="300" height="50" fill="#1e3a5f"/>'
        '</g>'
    ))
    prs = Presentation()
    layout = prs.slide_layouts[6]
    conv = SvgConverter(str(svg_path))

    with pytest.raises(NativeObjectMarkerError):
        conv.convert(prs, layout)


def test_unrecognized_chart_type_raises_native_object_marker_error(tmp_path: Path):
    # data.get("type") that isn't "bar_chart"/"line_chart"/"pie_chart" must
    # raise instead of silently defaulting to a bar chart -- a wrong chart
    # built with no error would be an unnoticed data-integrity failure.
    (tmp_path / "chart_data.json").write_text(json.dumps({
        "type": "scatter_chart",
        "categories": ["Q1", "Q2"],
        "series": [{"label": "EN", "color": "#1e3a5f", "values": [0.8, 0.9]}],
        "y_max": 1.0,
    }), encoding="utf-8")
    svg_path = _write_svg(tmp_path, (
        '<g data-pptx-role="chart" data-pptx-source="chart_data.json" '
        'data-pptx-bbox="130,100,970,420">'
        '<rect x="130" y="300" width="100" height="120" fill="#1e3a5f"/>'
        '</g>'
    ))
    prs = Presentation()
    layout = prs.slide_layouts[6]
    conv = SvgConverter(str(svg_path))

    with pytest.raises(NativeObjectMarkerError):
        conv.convert(prs, layout)


# ── Malformed-marker failure modes (all must be hard blockers) ───────────────
#
# _dispatch_children's broad per-child guard swallows any plain exception, so
# every failure mode in the native-marker code path must surface as
# NativeObjectMarkerError -- otherwise the table/chart is silently deleted
# from the deck while every gate stays green.

def _convert(svg_path: Path):
    prs = Presentation()
    layout = prs.slide_layouts[6]
    conv = SvgConverter(str(svg_path))
    return conv.convert(prs, layout)


def _table_marker(source: str, bbox: str = "60,75,300,100") -> str:
    return (
        f'<g data-pptx-role="table" data-pptx-source="{source}" '
        f'data-pptx-bbox="{bbox}">'
        '<rect x="60" y="75" width="300" height="50" fill="#1e3a5f"/>'
        '</g>'
    )


def _chart_marker(source: str, bbox: str = "130,100,970,420") -> str:
    return (
        f'<g data-pptx-role="chart" data-pptx-source="{source}" '
        f'data-pptx-bbox="{bbox}">'
        '<rect x="130" y="300" width="100" height="120" fill="#1e3a5f"/>'
        '</g>'
    )


def _write_json(directory: Path, name: str, payload) -> None:
    (directory / name).write_text(json.dumps(payload), encoding="utf-8")


VALID_TABLE = {"columns": ["Model", "QWK"], "rows": [["GPT-4", "0.671"]]}


@pytest.mark.parametrize("bbox", ["60,75,300", "60,75,300,100,120", "", ","])
def test_malformed_bbox_arity_raises_native_object_marker_error(
        tmp_path: Path, bbox: str):
    _write_json(tmp_path, "table_data.json", VALID_TABLE)
    svg_path = _write_svg(tmp_path, _table_marker("table_data.json", bbox))
    with pytest.raises(NativeObjectMarkerError):
        _convert(svg_path)


@pytest.mark.parametrize("bbox", ["a,b,c,d", "60,75,300,wide"])
def test_non_numeric_bbox_raises_native_object_marker_error(
        tmp_path: Path, bbox: str):
    _write_json(tmp_path, "table_data.json", VALID_TABLE)
    svg_path = _write_svg(tmp_path, _table_marker("table_data.json", bbox))
    with pytest.raises(NativeObjectMarkerError):
        _convert(svg_path)


def test_missing_sidecar_file_raises_native_object_marker_error(tmp_path: Path):
    # A misspelled sidecar filename is the realistic Path C authoring slip.
    svg_path = _write_svg(tmp_path, _table_marker("tabel_data.json"))
    with pytest.raises(NativeObjectMarkerError):
        _convert(svg_path)


def test_malformed_sidecar_json_raises_native_object_marker_error(tmp_path: Path):
    (tmp_path / "table_data.json").write_text("{not json,", encoding="utf-8")
    svg_path = _write_svg(tmp_path, _table_marker("table_data.json"))
    with pytest.raises(NativeObjectMarkerError):
        _convert(svg_path)


@pytest.mark.parametrize("payload", [
    {"rows": [["GPT-4", "0.671"]]},
    {"columns": ["Model", "QWK"]},
])
def test_table_sidecar_missing_required_keys_raises(tmp_path: Path, payload):
    _write_json(tmp_path, "table_data.json", payload)
    svg_path = _write_svg(tmp_path, _table_marker("table_data.json"))
    with pytest.raises(NativeObjectMarkerError):
        _convert(svg_path)


@pytest.mark.parametrize("payload", [
    {"type": "bar_chart", "series": [{"label": "EN", "values": [0.8]}]},
    {"type": "bar_chart", "categories": ["Q1"]},
    {"type": "pie_chart", "values": [50, 50]},
])
def test_chart_sidecar_missing_required_keys_raises(tmp_path: Path, payload):
    _write_json(tmp_path, "chart_data.json", payload)
    svg_path = _write_svg(tmp_path, _chart_marker("chart_data.json"))
    with pytest.raises(NativeObjectMarkerError):
        _convert(svg_path)


def test_unresolvable_slide_index_raises_native_object_marker_error(tmp_path: Path):
    _write_json(tmp_path, "slide_data.json", {
        "meta": {}, "slides": [{"index": 0, "type": "bullet_list"}]})
    svg_path = _write_svg(tmp_path, _table_marker("slide_data.json#4"))
    with pytest.raises(NativeObjectMarkerError):
        _convert(svg_path)


def test_non_integer_slide_index_raises_native_object_marker_error(tmp_path: Path):
    _write_json(tmp_path, "slide_data.json", {
        "meta": {}, "slides": [{"index": 0, "type": "table",
                                "columns": ["A"], "rows": [["1"]]}]})
    svg_path = _write_svg(tmp_path, _table_marker("slide_data.json#first"))
    with pytest.raises(NativeObjectMarkerError):
        _convert(svg_path)


def test_native_object_marker_error_chains_the_root_cause(tmp_path: Path):
    """The original exception must stay visible in the traceback."""
    svg_path = _write_svg(tmp_path, _table_marker("missing.json"))
    with pytest.raises(NativeObjectMarkerError) as excinfo:
        _convert(svg_path)
    assert isinstance(excinfo.value.__cause__, FileNotFoundError)


# ── M1: sidecar path containment ─────────────────────────────────────────────

@pytest.mark.parametrize("source", ["../outside.json", "sub/../../outside.json"])
def test_sidecar_escaping_the_svg_directory_raises(tmp_path: Path, source: str):
    svg_dir = tmp_path / "slides"
    svg_dir.mkdir()
    _write_json(tmp_path, "outside.json", VALID_TABLE)
    svg_path = _write_svg(svg_dir, _table_marker(source))
    with pytest.raises(NativeObjectMarkerError):
        _convert(svg_path)


def test_absolute_sidecar_path_raises(tmp_path: Path):
    svg_dir = tmp_path / "slides"
    svg_dir.mkdir()
    _write_json(tmp_path, "outside.json", VALID_TABLE)
    svg_path = _write_svg(
        svg_dir, _table_marker(str(tmp_path / "outside.json")))
    with pytest.raises(NativeObjectMarkerError):
        _convert(svg_path)


def test_sidecar_in_a_subdirectory_of_the_svg_directory_is_allowed(tmp_path: Path):
    sub = tmp_path / "data"
    sub.mkdir()
    _write_json(sub, "table_data.json", VALID_TABLE)
    svg_path = _write_svg(tmp_path, _table_marker("data/table_data.json"))
    slide = _convert(svg_path)
    assert any(getattr(s, "has_table", False) for s in slide.shapes)


# ── I1: malformed data-pptx-style ────────────────────────────────────────────

def test_malformed_pptx_style_json_raises_native_object_marker_error(tmp_path: Path):
    _write_json(tmp_path, "table_data.json", VALID_TABLE)
    svg_path = _write_svg(tmp_path, (
        '<g data-pptx-role="table" data-pptx-source="table_data.json" '
        'data-pptx-bbox="60,75,300,100" data-pptx-style="{not json">'
        '<rect x="60" y="75" width="300" height="50" fill="#1e3a5f"/>'
        '</g>'
    ))
    with pytest.raises(NativeObjectMarkerError):
        _convert(svg_path)


def test_absent_pptx_style_still_falls_back_to_defaults(tmp_path: Path):
    _write_json(tmp_path, "table_data.json", VALID_TABLE)
    svg_path = _write_svg(tmp_path, _table_marker("table_data.json"))
    slide = _convert(svg_path)
    table = next(s for s in slide.shapes if getattr(s, "has_table", False)).table
    header_run = table.cell(0, 0).text_frame.paragraphs[0].runs[0]
    assert str(header_run.font.color.rgb) == "FFFFFF"


# ── I2: per-child error tolerance inside a group marker ──────────────────────

def test_group_with_one_malformed_child_still_groups_the_valid_ones(tmp_path: Path):
    """One bad child must not delete the whole node, matching _dispatch_children."""
    svg_path = _write_svg(tmp_path, (
        '<g data-pptx-role="group" data-node-id="node1">'
        '<rect x="100" y="100" width="200" height="80" fill="#f8fafc"/>'
        '<rect x="100" y="oops" width="200" height="80" fill="#f8fafc"/>'
        '<text x="200" y="200">Encoder</text>'
        '</g>'
    ))
    slide = _convert(svg_path)

    from pptx.enum.shapes import MSO_SHAPE_TYPE
    groups = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.GROUP]
    assert len(groups) == 1
    assert len(groups[0].shapes) == 2


def test_group_child_marker_error_still_propagates(tmp_path: Path):
    """NativeObjectMarkerError stays a hard blocker even inside a group."""
    svg_path = _write_svg(tmp_path, (
        '<g data-pptx-role="group" data-node-id="node1">'
        '<rect x="100" y="100" width="200" height="80" fill="#f8fafc"/>'
        '<g data-pptx-role="table"><rect x="1" y="1" width="2" height="2"/></g>'
        '</g>'
    ))
    with pytest.raises(NativeObjectMarkerError):
        _convert(svg_path)
