import shutil
import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parents[2]))

import pytest
from lxml import etree
from pptx import Presentation
from pptx.util import Emu
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN

from svg_to_pptx.converter import CoordSystem
from svg_to_pptx.shapes import dispatch_shape
from svg_to_pptx.style_parser import compute_style

CS = CoordSystem(svg_w=1200.0, svg_h=675.0)


def _blank_slide():
    prs = Presentation()
    prs.slide_width = Emu(12_192_000)
    prs.slide_height = Emu(6_858_000)
    return prs.slides.add_slide(prs.slide_layouts[6]), prs


def test_rect_creates_one_shape():
    slide, _ = _blank_slide()
    elem = etree.fromstring('<rect x="100" y="50" width="200" height="100" fill="#3b82f6"/>')
    style = compute_style(elem, {})
    dispatch_shape(slide, elem, style, CS, None)
    assert len(slide.shapes) == 1


def test_rect_position_and_size():
    slide, _ = _blank_slide()
    elem = etree.fromstring('<rect x="0" y="0" width="1200" height="675" fill="#fff"/>')
    style = compute_style(elem, {})
    dispatch_shape(slide, elem, style, CS, None)
    shape = slide.shapes[0]
    assert shape.left == 0
    assert shape.top == 0
    assert shape.width == 12_192_000
    assert shape.height == 6_858_000


def test_circle_creates_oval():
    slide, _ = _blank_slide()
    elem = etree.fromstring('<circle cx="600" cy="337" r="100" fill="#10b981"/>')
    style = compute_style(elem, {})
    dispatch_shape(slide, elem, style, CS, None)
    assert len(slide.shapes) == 1
    shape = slide.shapes[0]
    expected_w = CS.x(200)
    assert abs(shape.width - expected_w) <= 1


def test_rect_with_label_writes_text_in_shape():
    slide, _ = _blank_slide()
    rect_elem = etree.fromstring('<rect x="40" y="80" width="160" height="70" fill="#3b82f6"/>')
    text_elem = etree.fromstring('<text x="120" y="115" fill="white">Hello</text>')
    style = compute_style(rect_elem, {})
    dispatch_shape(slide, rect_elem, style, CS, text_elem)
    shape = slide.shapes[0]
    assert shape.has_text_frame
    assert shape.text_frame.paragraphs[0].runs[0].text == "Hello"


def test_rect_with_tspan_only_label_still_writes_text_in_shape():
    """Regression: a <text> with coordinates only on its <tspan> (valid SVG,
    but violates this repo's own x/y-on-<text> authoring convention) must
    still resolve its position via _text_xy's fallback and embed correctly
    -- not silently anchor at (0, 0) and miss the shape entirely."""
    slide, _ = _blank_slide()
    rect_elem = etree.fromstring('<rect x="40" y="80" width="160" height="70" fill="#3b82f6"/>')
    text_elem = etree.fromstring(
        '<text fill="white" text-anchor="middle">'
        '<tspan x="120" y="108">Multi-Head</tspan>'
        '<tspan x="120" dy="15">Attention</tspan></text>'
    )
    style = compute_style(rect_elem, {})
    dispatch_shape(slide, rect_elem, style, CS, text_elem)
    shape = slide.shapes[0]
    assert shape.has_text_frame
    paragraphs = shape.text_frame.paragraphs
    assert paragraphs[0].runs[0].text == "Multi-Head"
    assert paragraphs[1].runs[0].text == "Attention"


def test_attached_labels_use_explicit_svg_text_layout() -> None:
    """Verify attached labels use a frame offset and preserve paragraph layout."""
    slide, _ = _blank_slide()
    rect_elem = etree.fromstring(
        '<rect x="40" y="80" width="160" height="70" fill="#3b82f6"/>'
    )
    label_elems = [
        etree.fromstring(
            '<text x="120" y="112" fill="white" text-anchor="middle" '
            'font-size="14">Title</text>'
        ),
        etree.fromstring(
            '<text x="180" y="132" fill="white" text-anchor="end" '
            'font-size="11">Subtitle</text>'
        ),
    ]
    style = compute_style(rect_elem, {})
    dispatch_shape(slide, rect_elem, style, CS, label_elems)

    tf = slide.shapes[0].text_frame
    assert tf.margin_left == Emu(0)
    assert tf.margin_right == Emu(0)
    assert tf.margin_bottom == Emu(0)
    assert tf.vertical_anchor == MSO_ANCHOR.TOP

    paragraphs = tf.paragraphs
    assert [paragraph.alignment for paragraph in paragraphs] == [
        PP_ALIGN.CENTER,
        PP_ALIGN.RIGHT,
    ]
    assert [
        run.text
        for paragraph in paragraphs
        for run in paragraph.runs
    ] == ["Title", "Subtitle"]
    expected_top_margin = CS.y(112 - 80 - 14 * 0.75 + 4.0)
    assert tf.margin_top == Emu(expected_top_margin)
    assert tf.margin_top > Emu(0)
    assert paragraphs[0].space_before == Emu(0)
    assert paragraphs[0].space_after == Emu(0)
    assert paragraphs[0].line_spacing == Emu(CS.y(14 * 1.2))
    assert paragraphs[1].space_before == Emu(CS.y(132 - 112 - 14 * 1.2))
    assert paragraphs[1].space_after == Emu(0)
    assert paragraphs[1].line_spacing == Emu(CS.y(11 * 1.2))


def test_attached_label_top_margin_is_bounded_by_shape_slack() -> None:
    """Verify the attached frame offset cannot consume later-line space."""
    slide, _ = _blank_slide()
    rect_elem = etree.fromstring(
        '<rect x="40" y="80" width="160" height="40" fill="#3b82f6"/>'
    )
    label_elems = [
        etree.fromstring(
            '<text x="120" y="112" fill="white" text-anchor="middle" '
            'font-size="14">Title</text>'
        ),
        etree.fromstring(
            '<text x="180" y="132" fill="white" text-anchor="end" '
            'font-size="11">Subtitle</text>'
        ),
    ]
    style = compute_style(rect_elem, {})
    dispatch_shape(slide, rect_elem, style, CS, label_elems)

    tf = slide.shapes[0].text_frame
    paragraphs = tf.paragraphs
    expected_bounded_margin = CS.y(40.0 - (16.8 + 3.2 + 13.2))
    assert tf.margin_top == Emu(expected_bounded_margin)
    assert tf.margin_top > Emu(0)
    assert tf.margin_top < Emu(CS.y(21.5 + 4.0))
    assert paragraphs[0].space_before == Emu(0)
    assert paragraphs[1].space_before == Emu(CS.y(3.2))
    assert (
        tf.margin_top
        + paragraphs[0].line_spacing
        + paragraphs[1].space_before
        + paragraphs[1].line_spacing
        <= Emu(CS.y(40.0))
    )


def test_ellipse_creates_oval():
    slide, _ = _blank_slide()
    elem = etree.fromstring('<ellipse cx="300" cy="200" rx="80" ry="40" fill="#f00"/>')
    style = compute_style(elem, {})
    dispatch_shape(slide, elem, style, CS, None)
    assert len(slide.shapes) == 1


from svg_to_pptx.converter import SvgConverter
import os, tempfile
from pptx import Presentation
from pptx.util import Emu

_ATTACHMENT_SVG = """<svg viewBox="0 0 1200 675" xmlns="http://www.w3.org/2000/svg">
  <rect x="40" y="80" width="160" height="70" fill="#3b82f6"/>
  <text x="120" y="115" fill="white" text-anchor="middle" font-size="14">Deep Research</text>
  <rect x="300" y="80" width="160" height="70" fill="#10b981"/>
  <text x="380" y="115" fill="white" text-anchor="middle" font-size="14">Academic Paper</text>
</svg>"""


def test_text_attaches_to_enclosing_rect():
    with tempfile.NamedTemporaryFile(suffix=".svg", mode="w", delete=False) as f:
        f.write(_ATTACHMENT_SVG)
        svg_path = f.name
    try:
        prs = Presentation()
        prs.slide_width = Emu(12_192_000)
        prs.slide_height = Emu(6_858_000)
        conv = SvgConverter(svg_path)
        slide = conv.convert(prs, prs.slide_layouts[6])
        assert len(slide.shapes) == 2
        texts = [sh.text_frame.paragraphs[0].runs[0].text
                 for sh in slide.shapes if sh.has_text_frame
                 and sh.text_frame.paragraphs[0].runs]
        assert "Deep Research" in texts
        assert "Academic Paper" in texts
    finally:
        os.unlink(svg_path)


_NEEDS_FONTCONFIG = pytest.mark.skipif(
    shutil.which("fc-match") is None,
    # fontconfig is an optional external binary; family resolution has no
    # meaning without it, so these four tests cannot run on a box that lacks it.
    reason="fontconfig (fc-match) is not installed",
)


def _label_font_name(markup: str) -> object:
    """Write one SVG text element into a shape and return its run font name.

    Args:
        markup: A single SVG `<text>` element as a string.

    Returns:
        The `font.name` of the first run, which is `None` when the converter
        set no explicit family.
    """
    from svg_to_pptx.shapes import _write_label

    slide, _ = _blank_slide()
    elem = etree.fromstring(markup)
    style = compute_style(elem, {})
    shape = slide.shapes.add_textbox(Emu(0), Emu(0), Emu(1000000), Emu(500000))
    _write_label(shape, elem, style, CS, (100, 100, 200, 40))
    return shape.text_frame.paragraphs[0].runs[0].font.name


@_NEEDS_FONTCONFIG
def test_multi_word_font_family_is_not_truncated():
    """A quoted multi-word family is never split into its first word."""
    name = _label_font_name(
        '<text x="100" y="100" font-size="21" '
        "font-family=\"'DejaVu Sans', Arial, sans-serif\">Node</text>"
    )
    assert name == "DejaVu Sans"
    assert name != "DejaVu"


@_NEEDS_FONTCONFIG
def test_uninstalled_first_family_falls_through_to_installed_one():
    """Resolution skips an uninstalled first choice rather than using it."""
    name = _label_font_name(
        '<text x="100" y="100" font-size="21" '
        "font-family=\"Totally Not A Real Font 9x7, 'DejaVu Sans'\">Node</text>"
    )
    assert name == "DejaVu Sans"


@_NEEDS_FONTCONFIG
def test_fully_uninstalled_stack_raises():
    """A stack with no installed family raises rather than substituting."""
    from fonts import FontError

    with pytest.raises(FontError):
        _label_font_name(
            '<text x="100" y="100" font-size="21" '
            'font-family="Nope One 9x7, \'Nope Two 9x7\'">Node</text>'
        )


@_NEEDS_FONTCONFIG
def test_absent_font_family_leaves_the_run_name_unset():
    """No font-family attribute means no explicit run font name."""
    assert _label_font_name(
        '<text x="100" y="100" font-size="21">Node</text>') is None


@_NEEDS_FONTCONFIG
def test_a_bare_generic_family_raises():
    """`sans-serif` alone names no installed face, so the export refuses it.

    A CSS generic is a request for whatever the renderer's default sans is. The
    old code wrote the literal string `sans-serif` as the PowerPoint font name,
    which PowerPoint cannot match and therefore substitutes -- the exact silent
    divergence between preview and deck that this change exists to stop. There
    is no family here to preserve, so the author is told to name one.
    """
    from fonts import FontError

    with pytest.raises(FontError, match="sans-serif"):
        _label_font_name(
            '<text x="100" y="100" font-size="21" '
            'font-family="sans-serif">Node</text>')

_A_NS = "{http://schemas.openxmlformats.org/drawingml/2006/main}"


def _preset_geometry(markup: str):
    """Dispatch one SVG element and return its shape and preset geometry name.

    Args:
        markup: A single SVG element as a string.

    Returns:
        A `(shape, prst)` pair, where `prst` is the DrawingML preset geometry
        name PowerPoint will render.
    """
    slide, _ = _blank_slide()
    elem = etree.fromstring(markup)
    style = compute_style(elem, {})
    dispatch_shape(slide, elem, style, CS, None)
    shape = slide.shapes[0]
    prst = shape._element.find(f".//{_A_NS}prstGeom").get("prst")
    return shape, prst


def test_sharp_rect_stays_a_rectangle():
    """A rect with no rx keeps the plain rectangle geometry."""
    _, prst = _preset_geometry(
        '<rect x="0" y="0" width="200" height="100" fill="#fff"/>')
    assert prst == "rect"


def test_rounded_rect_becomes_round_rect_geometry():
    """A rect with rx exports as roundRect, not a sharp box."""
    _, prst = _preset_geometry(
        '<rect x="0" y="0" width="200" height="100" rx="8" fill="#fff"/>')
    assert prst == "roundRect"


def test_rounded_rect_adjustment_matches_the_svg_radius():
    """The corner adjustment is rx as a fraction of the shorter side."""
    shape, _ = _preset_geometry(
        '<rect x="0" y="0" width="200" height="100" rx="10" fill="#fff"/>')
    # rx=10 against a 100-unit shorter side -> 0.10
    assert shape.adjustments[0] == pytest.approx(0.10, abs=0.005)


def test_rounded_rect_adjustment_is_clamped_to_a_stadium():
    """An rx larger than half the shorter side clamps instead of overflowing."""
    shape, _ = _preset_geometry(
        '<rect x="0" y="0" width="200" height="100" rx="400" fill="#fff"/>')
    assert shape.adjustments[0] == pytest.approx(0.5, abs=0.001)


def test_ry_alone_also_rounds_the_rect():
    """SVG allows ry without rx; both must round the shape."""
    _, prst = _preset_geometry(
        '<rect x="0" y="0" width="200" height="100" ry="8" fill="#fff"/>')
    assert prst == "roundRect"


def test_a_zero_radius_is_not_a_rounded_rect():
    """`rx="0"` is an explicitly sharp corner, not a degenerate roundRect."""
    _, prst = _preset_geometry(
        '<rect x="0" y="0" width="200" height="100" rx="0" fill="#fff"/>')
    assert prst == "rect"


def test_a_non_numeric_radius_raises():
    """A malformed radius must not silently become a sharp corner."""
    with pytest.raises(ValueError, match="rx"):
        _preset_geometry(
            '<rect x="0" y="0" width="200" height="100" rx="8px" fill="#fff"/>')


def test_the_token_card_radius_survives_export():
    """A card drawn at the token radius exports as a rounded card.

    The token contract gives `card` a 12-unit radius and `node` an 8-unit one.
    Every such surface exported as a sharp box, which is the single most
    visible difference between the SVG preview and the delivered deck.
    """
    shape, prst = _preset_geometry(
        '<rect x="48" y="108" width="352" height="144" rx="12" fill="#f8fafc"/>')
    assert prst == "roundRect"
    assert shape.adjustments[0] == pytest.approx(12 / 144, abs=0.005)


def _label_bold(markup: str) -> object:
    """Write one SVG text element into a shape and return its run bold flag.

    Args:
        markup: A single SVG `<text>` element as a string.

    Returns:
        The `font.bold` of the first run, which is `None` when the converter
        set no explicit weight.
    """
    from svg_to_pptx.shapes import _write_label

    slide, _ = _blank_slide()
    elem = etree.fromstring(markup)
    style = compute_style(elem, {})
    shape = slide.shapes.add_textbox(Emu(0), Emu(0), Emu(1000000), Emu(500000))
    _write_label(shape, elem, style, CS, (100, 100, 200, 40))
    return shape.text_frame.paragraphs[0].runs[0].font.bold


def test_semibold_text_exports_as_bold():
    """Weight 600 must reach PPTX as bold.

    `typography.roles.node_label` and `typography.roles.takeaway` are both
    weight 600. DrawingML has no numeric weight, only a bold flag, so a
    semibold role that is not mapped to it exports as regular type: the SVG
    shows an emphasised node label and the deck does not, on every diagram.
    """
    assert _label_bold(
        '<text x="100" y="100" font-size="18" font-weight="600">Node</text>'
    ) is True


def test_bold_text_still_exports_as_bold():
    """Weight 700 keeps working."""
    assert _label_bold(
        '<text x="100" y="100" font-size="32" font-weight="700">Title</text>'
    ) is True


def test_regular_text_is_not_made_bold():
    """Weight 400 must not be emphasised.

    A threshold that swept everything into bold would be as wrong as the
    omission: body copy at weight 400 is most of the text on a deck.
    """
    assert _label_bold(
        '<text x="100" y="100" font-size="21" font-weight="400">Body</text>'
    ) is not True


def test_medium_text_is_not_made_bold():
    """Weight 500 is medium, not semibold, and stays regular.

    CSS font matching treats 600 and above as the bold end of a family; 500 is
    on the regular side, and the token contract never asks for it.
    """
    assert _label_bold(
        '<text x="100" y="100" font-size="21" font-weight="500">Body</text>'
    ) is not True


def _para_indent(paragraph: object) -> int:
    """Return a paragraph's left indent in EMU.

    Args:
        paragraph: A python-pptx paragraph.

    Returns:
        The `a:pPr/@marL` value, or 0 when the paragraph has no indent.
    """
    pPr = paragraph._p.find(
        "{http://schemas.openxmlformats.org/drawingml/2006/main}pPr")
    if pPr is None:
        return 0
    return int(pPr.get("marL") or 0)


def test_start_anchored_label_keeps_its_horizontal_inset() -> None:
    """A label anchored at `start` inside a card keeps its padding.

    Nothing carried the label's x, so a left-aligned label rendered flush
    against the shape's border instead of at the x the SVG gave it: on a card
    with 24 units of padding the body copy sat on the border.
    """
    slide, _ = _blank_slide()
    rect_elem = etree.fromstring(
        '<rect x="612" y="108" width="540" height="531" fill="#f8fafc"/>'
    )
    label = etree.fromstring(
        '<text x="636" y="150" font-size="26" text-anchor="start" '
        'fill="#1e3a5f">After</text>'
    )
    style = compute_style(rect_elem, {})
    dispatch_shape(slide, rect_elem, style, CS, [label])

    frame = slide.shapes[0].text_frame
    assert _para_indent(frame.paragraphs[0]) == CS.x(24.0)
    # The indent belongs to the paragraph; the frame margin stays uniform.
    assert frame.margin_left == Emu(0)


def test_each_label_keeps_its_own_indent() -> None:
    """Labels at different x values indent independently.

    The conclusion layout draws a bullet dot at cx=80 and starts its copy at
    101.88 to clear it, while the panel heading starts at 72. One margin for
    the whole frame can only honour one of those, and honouring the heading's
    put the copy on top of the dot.
    """
    slide, _ = _blank_slide()
    rect_elem = etree.fromstring(
        '<rect x="48" y="108" width="540" height="531" fill="#f8fafc"/>'
    )
    heading = etree.fromstring(
        '<text x="72" y="158" font-size="26" fill="#1e3a5f">Conclusions</text>'
    )
    body = etree.fromstring(
        '<text x="101.88" y="217" font-size="21" text-anchor="start" '
        'fill="#374151">One token set now describes each deck.</text>'
    )
    style = compute_style(rect_elem, {})
    dispatch_shape(slide, rect_elem, style, CS, [heading, body])

    paragraphs = slide.shapes[0].text_frame.paragraphs
    assert _para_indent(paragraphs[0]) == CS.x(24.0)
    assert _para_indent(paragraphs[1]) == CS.x(53.88)


def test_a_tspan_keeps_its_own_indent() -> None:
    """A tspan that restates x indents from that x, not from its parent."""
    slide, _ = _blank_slide()
    rect_elem = etree.fromstring(
        '<rect x="48" y="108" width="540" height="531" fill="#f8fafc"/>'
    )
    label = etree.fromstring(
        '<text x="72" y="199" font-size="21" text-anchor="start" '
        'fill="#374151"><tspan x="72" dy="0">first</tspan>'
        '<tspan x="120" dy="30.4">second</tspan></text>'
    )
    style = compute_style(rect_elem, {})
    dispatch_shape(slide, rect_elem, style, CS, [label])

    paragraphs = slide.shapes[0].text_frame.paragraphs
    assert _para_indent(paragraphs[0]) == CS.x(24.0)
    assert _para_indent(paragraphs[1]) == CS.x(72.0)


def test_centred_label_takes_no_horizontal_inset() -> None:
    """A `middle`-anchored label is centred by alignment, not by an indent.

    Indenting a centred label would shift it off the shape's centre, which is
    where every diagram node label belongs.
    """
    slide, _ = _blank_slide()
    rect_elem = etree.fromstring(
        '<rect x="40" y="80" width="160" height="70" fill="#3b82f6"/>'
    )
    label = etree.fromstring(
        '<text x="120" y="112" font-size="14" text-anchor="middle" '
        'fill="white">Node</text>'
    )
    style = compute_style(rect_elem, {})
    dispatch_shape(slide, rect_elem, style, CS, [label])

    assert _para_indent(slide.shapes[0].text_frame.paragraphs[0]) == 0


def test_a_start_anchored_inset_cannot_exceed_the_shape() -> None:
    """A label anchored outside its own shape does not push text off it.

    `_compute_text_attachments` only attaches a label whose anchor is inside
    the shape, but `dispatch_shape` is also called directly, and an indent
    wider than the shape would leave an empty box.
    """
    slide, _ = _blank_slide()
    rect_elem = etree.fromstring(
        '<rect x="40" y="80" width="100" height="70" fill="#3b82f6"/>'
    )
    label = etree.fromstring(
        '<text x="900" y="112" font-size="14" text-anchor="start" '
        'fill="white">Node</text>'
    )
    style = compute_style(rect_elem, {})
    dispatch_shape(slide, rect_elem, style, CS, [label])

    assert _para_indent(slide.shapes[0].text_frame.paragraphs[0]) == 0


def test_an_unanchored_label_is_left_aligned_like_the_svg() -> None:
    """SVG's default `text-anchor` is `start`, so PPTX must left-align too.

    The default here was `middle`, so a panel heading the SVG draws from its
    x rendered centred in the exported deck. The two renderings of the same
    slide disagreed, which is exactly what the export is supposed to rule out.
    """
    slide, _ = _blank_slide()
    rect_elem = etree.fromstring(
        '<rect x="612" y="108" width="540" height="531" fill="#f8fafc"/>'
    )
    label = etree.fromstring(
        '<text x="636" y="150" font-size="26" fill="#1e3a5f">After</text>'
    )
    style = compute_style(rect_elem, {})
    dispatch_shape(slide, rect_elem, style, CS, [label])

    paragraph = slide.shapes[0].text_frame.paragraphs[0]
    assert paragraph.alignment == PP_ALIGN.LEFT
    assert _para_indent(paragraph) == CS.x(24.0)
