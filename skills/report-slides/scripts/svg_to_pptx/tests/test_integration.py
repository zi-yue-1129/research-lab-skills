"""test_integration.py — end-to-end SVG → PPTX shape verification."""
import io, os, tempfile
import pytest
from pptx import Presentation
from pptx.util import Emu

from svg_to_pptx.converter import CoordSystem, SvgConverter, convert_file

_DIAGRAM_SVG = """\
<svg viewBox="0 0 1200 675" xmlns="http://www.w3.org/2000/svg">
  <rect x="40" y="80" width="160" height="70" fill="#3b82f6" stroke="#1d4ed8" stroke-width="2"/>
  <text x="120" y="112" fill="white" text-anchor="middle" font-size="14" font-weight="bold">Deep Research</text>
  <text x="120" y="132" fill="white" text-anchor="middle" font-size="11">13-agent pipeline</text>
  <circle cx="600" cy="337" r="50" fill="#10b981"/>
  <text x="600" y="341" fill="white" text-anchor="middle" font-size="13">Hub</text>
  <line x1="200" y1="115" x2="550" y2="337" stroke="#6b7280" stroke-width="1.5"/>
</svg>"""

_BEZIER_SVG = """\
<svg viewBox="0 0 400 200" xmlns="http://www.w3.org/2000/svg">
  <path d="M 0 100 C 100 0 300 0 400 100" stroke="#3b82f6" fill="none" stroke-width="3"/>
</svg>"""

_ARC_SVG = """\
<svg viewBox="0 0 400 400" xmlns="http://www.w3.org/2000/svg">
  <path d="M 200 50 A 150 150 0 0 1 350 200" stroke="#f59e0b" fill="none" stroke-width="2"/>
</svg>"""

_CANVAS_BACKGROUND_SVG = """\
<svg viewBox="0 0 1200 675" xmlns="http://www.w3.org/2000/svg">
  <rect x="0" y="0" width="1200.0000001" height="675.0000001" fill="white"/>
  <rect x="40" y="80" width="160" height="70" fill="#3b82f6"/>
  <text x="120" y="115" fill="white" text-anchor="middle" font-size="14">Inside</text>
  <text x="600" y="610" fill="black" text-anchor="start" font-size="18">Outside</text>
</svg>"""


def _make_tmp_svg(content: str) -> str:
    f = tempfile.NamedTemporaryFile(suffix=".svg", mode="w", delete=False)
    f.write(content); f.close()
    return f.name


def _conv(svg_content: str):
    path = _make_tmp_svg(svg_content)
    try:
        prs = Presentation()
        prs.slide_width = Emu(12_192_000)
        prs.slide_height = Emu(6_858_000)
        conv = SvgConverter(path)
        slide = conv.convert(prs, prs.slide_layouts[6])
        return slide, prs
    finally:
        os.unlink(path)


def test_diagram_shape_count():
    slide, _ = _conv(_DIAGRAM_SVG)
    assert len(slide.shapes) == 3


def test_diagram_text_in_rect():
    slide, _ = _conv(_DIAGRAM_SVG)
    rect = slide.shapes[0]
    assert rect.has_text_frame
    texts = [r.text for p in rect.text_frame.paragraphs for r in p.runs]
    assert any("Deep Research" in t for t in texts)


def test_canvas_background_does_not_capture_standalone_text():
    slide, _ = _conv(_CANVAS_BACKGROUND_SVG)

    assert len(slide.shapes) == 3
    assert slide.shapes[0].text == ""
    assert slide.shapes[1].text == "Inside"
    standalone = slide.shapes[2]
    assert standalone.text == "Outside"

    cs = CoordSystem(svg_w=1200.0, svg_h=675.0)
    assert standalone.left == cs.x(600.0)
    assert standalone.top == cs.y(610.0 - 18.0 * 0.75)


def test_diagram_connector_present():
    slide, _ = _conv(_DIAGRAM_SVG)
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    connectors = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.LINE]
    assert len(connectors) >= 1


def test_bezier_path_saves_valid_pptx():
    slide, prs = _conv(_BEZIER_SVG)
    assert len(slide.shapes) == 1
    buf = io.BytesIO()
    prs.save(buf)
    assert len(buf.getvalue()) > 5000


def test_arc_path_saves_valid_pptx():
    slide, prs = _conv(_ARC_SVG)
    assert len(slide.shapes) == 1
    buf = io.BytesIO()
    prs.save(buf)
    assert len(buf.getvalue()) > 5000


def test_convert_file_creates_pptx(tmp_path):
    svg_dir = tmp_path / "slides"
    svg_dir.mkdir()
    (svg_dir / "slide01.svg").write_text(_DIAGRAM_SVG)
    out = str(tmp_path / "deck.pptx")
    convert_file(str(svg_dir), out)
    assert os.path.exists(out)
    prs = Presentation(out)
    assert len(prs.slides) == 1


def test_all_shapes_have_positive_size():
    slide, _ = _conv(_DIAGRAM_SVG)
    for shape in slide.shapes:
        assert shape.width > 0
        assert shape.height > 0


_ROOT_FONT_SVG = """\
<svg viewBox="0 0 1200 675" xmlns="http://www.w3.org/2000/svg" \
font-family="DejaVu Sans">
  <text x="100" y="100" font-size="21" data-style-role="body">Inherited</text>
</svg>"""


def test_root_svg_font_family_reaches_the_run():
    """A font-family declared once on <svg> must reach every text run.

    The root traversal was seeded with an empty inherited style, so every
    presentation attribute on the `<svg>` element itself was discarded. The
    deterministic renderer declares the deck's resolved font exactly once,
    there -- so the generated decks carried no run font name at all and
    PowerPoint fell back to the layout font, diverging from the SVG preview on
    every slide. Fixing the truncation in `_apply_font` alone left this intact,
    because there was no `font-family` in scope for it to resolve.
    """
    path = _make_tmp_svg(_ROOT_FONT_SVG)
    try:
        prs = Presentation()
        prs.slide_width = Emu(12_192_000)
        prs.slide_height = Emu(6_858_000)
        SvgConverter(path).convert(prs, prs.slide_layouts[6])
        names = {
            run.font.name
            for slide in prs.slides
            for shape in slide.shapes if shape.has_text_frame
            for para in shape.text_frame.paragraphs
            for run in para.runs
        }
        assert names == {"DejaVu Sans"}
    finally:
        os.unlink(path)


_A_NS = "{http://schemas.openxmlformats.org/drawingml/2006/main}"

_GRADIENT_SVG = """\
<svg viewBox="0 0 1200 675" xmlns="http://www.w3.org/2000/svg">
  <defs>
    <linearGradient id="fade" x1="0" y1="0" x2="1" y2="0">
      <stop offset="0%" stop-color="#1e3a5f"/>
      <stop offset="100%" stop-color="#0f766e"/>
    </linearGradient>
  </defs>
  <rect x="100" y="100" width="400" height="200" fill="url(#fade)"/>
</svg>"""

_BAD_GEOMETRY_SVG = """\
<svg viewBox="0 0 1200 675" xmlns="http://www.w3.org/2000/svg">
  <rect x="100" y="100" width="not-a-number" height="200" fill="#1e3a5f"/>
</svg>"""


def test_gradient_fill_produces_a_grad_fill_element() -> None:
    """A url(#id) fill resolves to a native OOXML gradient, not an empty fill."""
    slide, _ = _conv(_GRADIENT_SVG)
    shape = slide.shapes[0]
    grad = shape._element.find(f".//{_A_NS}gradFill")
    assert grad is not None
    stops = grad.findall(f"{_A_NS}gsLst/{_A_NS}gs")
    assert len(stops) == 2
    assert stops[0].find(f"{_A_NS}srgbClr").get("val") == "1E3A5F"
    assert stops[1].find(f"{_A_NS}srgbClr").get("val") == "0F766E"


def test_gradient_stops_carry_their_offsets() -> None:
    """Stop offsets survive as OOXML gs positions."""
    slide, _ = _conv(_GRADIENT_SVG)
    grad = slide.shapes[0]._element.find(f".//{_A_NS}gradFill")
    positions = [
        stop.get("pos") for stop in grad.findall(f"{_A_NS}gsLst/{_A_NS}gs")
    ]
    assert positions == ["0", "100000"]


def test_solid_fill_is_unaffected_by_gradient_support() -> None:
    """A plain hex fill still produces a solidFill, not a gradient."""
    slide, _ = _conv(_DIAGRAM_SVG)
    shape = slide.shapes[0]
    assert shape._element.find(f".//{_A_NS}solidFill") is not None
    assert shape._element.find(f".//{_A_NS}gradFill") is None


def test_percentage_gradient_coordinates_are_accepted() -> None:
    """`x2="100%"` is ordinary SVG and must not crash the converter.

    linearGradient coordinates default to objectBoundingBox units, where a
    percentage and a fraction mean the same thing, and most authoring tools
    emit the percentage. `float("100%")` raises `ValueError`, so a naive parse
    would die on markup it should render.
    """
    svg = _GRADIENT_SVG.replace('x2="1"', 'x2="100%"')
    slide, _ = _conv(svg)
    assert slide.shapes[0]._element.find(f".//{_A_NS}gradFill") is not None


def test_a_nonsense_gradient_coordinate_is_named() -> None:
    """An unparsable coordinate says which attribute was wrong."""
    svg = _GRADIENT_SVG.replace('x2="1"', 'x2="halfway"')
    with pytest.raises(ValueError, match="x2"):
        _conv(svg)


def test_an_unknown_paint_server_is_reported() -> None:
    """A fill naming a gradient that does not exist is an error.

    It must also reach the caller: `_dispatch_children` guards each child with
    a broad `except Exception`, so an ordinary error here would be logged at
    verbose level and the shape would export unfilled with nothing to explain
    it.
    """
    svg = _GRADIENT_SVG.replace('fill="url(#fade)"', 'fill="url(#nope)"')
    with pytest.raises(ValueError, match="nope"):
        _conv(svg)


@pytest.mark.parametrize("attribute,value", [
    ("gradientUnits", "userSpaceOnUse"),
    ("spreadMethod", "reflect"),
    ("gradientTransform", "rotate(45)"),
])
def test_an_unsupported_gradient_feature_is_refused(
        attribute: str, value: str) -> None:
    """Features that change the rendering are refused, not dropped.

    Accepting the markup and ignoring the feature exports a deck that looks
    wrong with nothing in the log to explain it -- and the person who sees the
    render is not the person reading the code.
    """
    svg = _GRADIENT_SVG.replace(
        '<linearGradient id="fade"',
        f'<linearGradient id="fade" {attribute}="{value}"')
    with pytest.raises(ValueError, match=attribute):
        _conv(svg)


def test_a_stop_opacity_that_would_be_dropped_is_refused() -> None:
    """DrawingML stops here carry no alpha, so a translucent stop is an error."""
    svg = _GRADIENT_SVG.replace(
        'stop-color="#0f766e"', 'stop-color="#0f766e" stop-opacity="0.4"')
    with pytest.raises(ValueError, match="stop-opacity"):
        _conv(svg)


def test_a_stop_without_a_colour_is_refused() -> None:
    """A stop with no stop-color contributes nothing and is an error."""
    svg = _GRADIENT_SVG.replace('stop-color="#0f766e"', "")
    with pytest.raises(ValueError, match="stop-color"):
        _conv(svg)


def test_malformed_geometry_raises_instead_of_dropping_the_shape() -> None:
    """A bad geometry attribute is reported, not silently swallowed.

    The registry block caught every exception and passed. The shape then never
    entered `_shape_registry`, so `_bind_connectors` could not anchor anything
    to it -- producing exactly the dangling connectors the review gate exists
    to catch, with no diagnostic anywhere.
    """
    with pytest.raises(ValueError, match="not-a-number"):
        _conv(_BAD_GEOMETRY_SVG)


def test_the_whole_deck_route_still_exports(tmp_path) -> None:
    """The gradient path also survives the directory-to-deck entry point.

    `convert_file(slides_dir, out_path)` is the route `SKILL.md` actually
    calls; the in-memory `_conv` helper bypasses `prs.save`, so a shape that
    python-pptx accepts in memory but refuses to serialise would slip through
    every other test in this task.
    """
    slides_dir = tmp_path / "slides"
    slides_dir.mkdir()
    (slides_dir / "slide01.svg").write_text(_GRADIENT_SVG, encoding="utf-8")
    out = tmp_path / "deck.pptx"
    convert_file(str(slides_dir), str(out))
    assert out.exists()
    assert len(Presentation(str(out)).slides) == 1


_CARD_WITH_FOOTER_SVG = """\
<svg viewBox="0 0 1200 675" xmlns="http://www.w3.org/2000/svg">
  <rect x="0" y="0" width="1200" height="675" fill="#ffffff"/>
  <text x="1152" y="636" font-size="12" data-style-role="footnote"
        text-anchor="end" fill="#64748b">deck footer</text>
  <rect x="612" y="108" width="540" height="531" fill="#f8fafc"/>
  <text x="636" y="150" font-size="26" data-style-role="takeaway"
        text-anchor="start" fill="#1e3a5f">After</text>
  <text x="636" y="199" font-size="21" data-style-role="body"
        text-anchor="start" fill="#374151">Every size resolves.</text>
</svg>"""


def _card_shape(slide):
    """Return the card rect, identified by the heading it carries.

    Args:
        slide: The converted slide.

    Returns:
        The shape whose text frame holds the panel heading.
    """
    for shape in slide.shapes:
        if shape.has_text_frame and "After" in shape.text_frame.text:
            return shape
    raise AssertionError("no shape carries the panel heading")


def test_slide_chrome_is_not_absorbed_by_a_card():
    """The deck footer is not any shape's label.

    The footer's anchor falls inside a full-height panel, so the attachment
    pass claimed it. Because labels are collected in document order and the
    footer is emitted before the panels, its baseline then set the shape's
    first-paragraph margin: the panel heading exported roughly 400 units below
    where the SVG drew it, and the footer lost its own position entirely.
    """
    slide, _ = _conv(_CARD_WITH_FOOTER_SVG)
    card = _card_shape(slide)
    assert "deck footer" not in card.text_frame.text
    assert "Every size resolves." in card.text_frame.text


def test_absorbed_chrome_no_longer_drives_the_frame_margin():
    """The card's first-paragraph margin comes from its own heading.

    This is the assertion that pins the visible defect rather than its cause:
    26.5 units is the heading's own offset inside the card, and the footer's
    baseline would have produced roughly 523.
    """
    slide, _ = _conv(_CARD_WITH_FOOTER_SVG)
    card = _card_shape(slide)
    cs = CoordSystem(svg_w=1200.0, svg_h=675.0)
    assert card.text_frame.margin_top == Emu(cs.y(150 - 108 - 26 * 0.75 + 4.0))


def test_slide_chrome_keeps_its_own_position():
    """The footer stays a standalone textbox at its own baseline."""
    slide, _ = _conv(_CARD_WITH_FOOTER_SVG)
    footers = [s for s in slide.shapes
               if s.has_text_frame and s.text_frame.text == "deck footer"]
    assert len(footers) == 1
    cs = CoordSystem(svg_w=1200.0, svg_h=675.0)
    assert footers[0].top == Emu(cs.y(636.0 - 12.0 * 0.75))


_BADGE_SVG = """\
<svg viewBox="0 0 1200 675" xmlns="http://www.w3.org/2000/svg">
  <rect x="0" y="0" width="1200" height="675" fill="#ffffff"/>
  <text x="1152" y="636" font-size="12" data-style-role="footnote"
        text-anchor="end" fill="#64748b">deck footer</text>
  <rect x="612" y="108" width="540" height="531" fill="#f8fafc"/>
  <circle cx="652" cy="210.28" r="10.2" fill="#047857"/>
  <text x="652" y="215.89" font-size="12" font-weight="700"
        data-style-role="footnote" text-anchor="middle" fill="#ffffff">1</text>
  <text x="678.2" y="217" font-size="21" data-style-role="body"
        text-anchor="start" fill="#374151">Bind the linter.</text>
</svg>"""


def test_a_badge_digit_stays_on_its_badge():
    """A chrome-sized role inside a badge is still that badge's label.

    `data-style-role` names a type role, not a place: the numbered-bullet
    badge on the conclusion layout carries `footnote` because it is set in
    12/700 type. Excluding the role outright detached the digit from its
    circle and left it floating as its own textbox.
    """
    slide, _ = _conv(_BADGE_SVG)
    badges = [s for s in slide.shapes
              if s.has_text_frame and s.text_frame.text == "1"]
    assert len(badges) == 1
    cs = CoordSystem(svg_w=1200.0, svg_h=675.0)
    assert badges[0].left == Emu(cs.x(652 - 10.2))
    assert badges[0].width == Emu(cs.x(20.4))
