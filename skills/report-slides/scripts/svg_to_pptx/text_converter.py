"""text_converter.py — standalone TextBox for SVG <text> elements."""
from __future__ import annotations

import re
from typing import Any, Dict, List, Optional, Tuple

from pptx.util import Emu
from pptx.enum.text import PP_ALIGN

from fonts import parse_font_stack, resolve_font_stack, text_width
from .style_parser import compute_style
from .converter import CoordSystem, _local_tag, _text_xy
from .shapes import _apply_font, _ALIGN

_ASCENT_FACTOR = 0.75
# Only reached by an SVG that names no font family anywhere: nothing states
# which face such a document would be rendered in, so there is nothing to
# measure. Every deck this skill produces declares one.
_CHAR_W_FACTOR = 0.65
_MIN_TB_W_SVG = 30.0
_PAD_SVG = 24.0         # horizontal padding in SVG units
_EMPTY_LABEL_CHARS = 6
_WEIGHT_KEYWORDS = {"normal": 400, "bold": 700}


def _font_size_svg(style: Dict) -> float:
    raw = style.get("font-size", "14")
    try:
        return float(re.sub(r"[^0-9.]", "", raw) or "14")
    except ValueError:
        return 14.0


def _tspan_dy(ts: Any) -> float:
    try:
        return float(ts.get("dy", "0") or "0")
    except (ValueError, TypeError):
        return 0.0


def _weight_number(style: Dict) -> int:
    """Return a CSS font-weight as a number.

    Args:
        style: A resolved style dict.

    Returns:
        The numeric weight, defaulting to 400. A bold face is materially wider
        than its regular sibling, so measuring one with the other is how a
        heading silently overflows.
    """
    raw = str(style.get("font-weight", "normal")).strip()
    if raw in _WEIGHT_KEYWORDS:
        return _WEIGHT_KEYWORDS[raw]
    try:
        return int(float(raw))
    except ValueError:
        return 400


def _text_lines(elem: Any, style: Dict) -> List[Tuple[str, Dict]]:
    """Return each rendered line of a `<text>` element with its own style.

    Args:
        elem: The `<text>` element.
        style: Its resolved style, which the tspans inherit.

    Returns:
        One `(text, style)` pair per non-empty line.
    """
    lines: List[Tuple[str, Dict]] = []
    direct = (elem.text or "").strip()
    if direct:
        lines.append((direct, style))
    for ts in elem:
        if _local_tag(ts) != "tspan":
            continue
        ts_text = (ts.text or "").strip()
        if ts_text:
            lines.append((ts_text, compute_style(ts, style)))
    return lines


def _measured_line_width(line: str, style: Dict) -> Optional[float]:
    """Measure one line in the face it will actually be set in.

    Args:
        line: The line's text.
        style: The line's resolved style.

    Returns:
        The advance width in SVG units, or None when the style names no font
        family and there is therefore nothing to measure.
    """
    family_css = style.get("font-family", "")
    if not family_css or not parse_font_stack(family_css):
        return None
    return text_width(line, resolve_font_stack(family_css),
                      _font_size_svg(style), _weight_number(style))


def _text_box_width(elem: Any, style: Dict, fs: float) -> float:
    """Return the textbox width in SVG units, measured from the glyphs.

    The width was the longest line's character count times 0.65 em. No font
    has a uniform advance, so the box could come out narrower than the text in
    it, and with wrapping off the run was then laid out against a box that
    disagreed with the SVG the deck was reviewed in.

    Args:
        elem: The `<text>` element.
        style: Its resolved style.
        fs: Its font size in SVG units, used to size an empty label.

    Returns:
        The width in SVG units, never below `_MIN_TB_W_SVG`.
    """
    widths = []
    for line, line_style in _text_lines(elem, style):
        measured = _measured_line_width(line, line_style)
        if measured is None:
            measured = len(line) * _font_size_svg(line_style) * _CHAR_W_FACTOR
        widths.append(measured)
    if not widths:
        widths.append(_EMPTY_LABEL_CHARS * fs * _CHAR_W_FACTOR)
    return max(max(widths) + _PAD_SVG, _MIN_TB_W_SVG)


def add_textbox(slide: Any, elem: Any, style: Dict, cs: CoordSystem) -> Any:
    tx, ty = _text_xy(elem)

    fs = _font_size_svg(style)
    anchor = style.get("text-anchor", "start")

    # SVG y is text baseline; PPTX textbox is positioned from top
    ty_top = ty - fs * _ASCENT_FACTOR

    tw = _text_box_width(elem, style, fs)
    if anchor == "middle":
        tx_left = tx - tw / 2
    elif anchor == "end":
        tx_left = tx - tw
    else:
        tx_left = tx

    tspans = [t for t in elem if _local_tag(t) == "tspan"]
    dy_values = [d for ts in tspans if (d := _tspan_dy(ts)) > 0]
    n_newlines = len(dy_values)
    dy_rep = dy_values[0] if dy_values else fs * 1.4
    tb_h_svg = fs * _ASCENT_FACTOR + n_newlines * dy_rep + fs * 0.3

    ex = cs.x(tx_left)
    ey = cs.y(ty_top)
    ew = cs.x(tw)
    eh = cs.y(max(tb_h_svg, fs * 1.3))

    tb = slide.shapes.add_textbox(Emu(ex), Emu(ey), Emu(ew), Emu(eh))
    tf = tb.text_frame
    tf.word_wrap = False
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)

    _fill_text_frame(tf, elem, style, cs, anchor)
    return tb


def _fill_text_frame(tf: Any, elem: Any, style: Dict, cs: CoordSystem, anchor: str) -> None:
    tspans = [t for t in elem if _local_tag(t) == "tspan"]
    fs = _font_size_svg(style)
    align = _ALIGN.get(anchor, PP_ALIGN.LEFT)

    current_para = tf.paragraphs[0]
    current_para.alignment = align

    dt = (elem.text or "").strip()
    if dt:
        run = current_para.add_run()
        run.text = dt
        _apply_font(run, style, style)

    for ts in tspans:
        dy = _tspan_dy(ts)
        ts_style = compute_style(ts, style)
        ts_text = (ts.text or "").strip()
        ts_anchor = ts_style.get("text-anchor", anchor)
        ts_align = _ALIGN.get(ts_anchor, PP_ALIGN.LEFT)

        if dy > 0:
            p = tf.add_paragraph()
            p.alignment = ts_align
            space_svg = max(0.0, dy - fs)
            if space_svg > 0:
                p.space_before = Emu(cs.y(space_svg))
            current_para = p
        else:
            current_para.alignment = ts_align

        if ts_text:
            run = current_para.add_run()
            run.text = ts_text
            _apply_font(run, ts_style, style)
