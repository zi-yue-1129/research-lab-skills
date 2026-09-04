"""converter.py — orchestration layer for SVG → native PPTX conversion."""
from __future__ import annotations

from dataclasses import dataclass
import math
import re
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

from lxml import etree
from pptx import Presentation
from pptx.util import Emu

PPTX_W = 12_192_000
PPTX_H = 6_858_000
_SVG_NS = "http://www.w3.org/2000/svg"
_XLINK_NS = "http://www.w3.org/1999/xlink"
_CANVAS_TOLERANCE = 1e-6
_CHART_TYPE_MAP = {"bar_chart": "bar", "line_chart": "line", "pie_chart": "pie"}


class SvgSourceError(ValueError):
    """Raised for an error in the SVG source that must reach the caller.

    `_dispatch_children` guards each child with a broad `except Exception` so
    one malformed element does not lose the rest of the slide. That guard is
    right for a shape python-pptx merely refuses, and wrong for a source error
    the author can fix: a non-numeric geometry attribute or a fill naming a
    gradient that does not exist would be logged at verbose level and the deck
    would export missing a shape, an anchor, or a fill, with nothing to explain
    it. Like `NativeObjectMarkerError`, this propagates.
    """


class NativeObjectMarkerError(ValueError):
    """Raised when a data-pptx-role="table"/"chart" marker is missing or
    malformed. Must propagate as a hard blocker -- never silently caught
    by _dispatch_children's broad per-child exception guard."""


@dataclass
class CoordSystem:
    svg_w: float
    svg_h: float

    def x(self, v: float) -> int:
        return round(float(v) / self.svg_w * PPTX_W)

    def y(self, v: float) -> int:
        return round(float(v) / self.svg_h * PPTX_H)

    @classmethod
    def from_viewbox(cls, viewbox: str) -> "CoordSystem":
        parts = viewbox.strip().split()
        if len(parts) == 4:
            try:
                return cls(svg_w=float(parts[2]), svg_h=float(parts[3]))
            except ValueError:
                pass
        return cls(svg_w=1200.0, svg_h=675.0)


def _local_tag(elem: Any) -> str:
    tag = elem.tag
    if not isinstance(tag, str):
        return ""
    return tag.split("}")[-1] if "}" in tag else tag


# Roles a slide's own furniture is set in. `data-style-role` names a type
# role rather than a place, so these also cover badge digits and pill labels;
# `_CHROME_HOST_HEIGHT_RATIO` is what tells the two apart.
_SLIDE_CHROME_ROLES = frozenset({"deck_title", "slide_title", "footnote"})

# A badge, pill, or node is drawn around its label and is a shape or two
# taller than the type in it. A panel that merely encloses a footer's anchor
# is tens of times taller. Chrome type is attached only to a host of the
# first kind.
_CHROME_HOST_HEIGHT_RATIO = 3.0


def _text_xy(text_elem: Any) -> Tuple[float, float]:
    """Resolve a ``<text>`` element's anchor point, tolerating SVG that
    omits ``x``/``y`` on the tag itself and only sets them on the first
    ``<tspan>``.

    This codebase's own authoring convention requires ``x``/``y`` on the
    ``<text>`` element itself (matching every renderer in this repo, e.g.
    ``generate_slides.py``'s ``tlines()`` -- enforced separately by
    ``validate_native_objects.py``'s SVG-source scan). This fallback exists
    only as a second line of defense so a violation of that convention
    degrades gracefully instead of silently attaching/positioning the text
    at ``(0, 0)``.

    Args:
        text_elem: An lxml ``<text>`` element.

    Returns:
        ``(x, y)`` in the SVG's own user units. ``0.0`` for whichever
        coordinate is present on neither the element nor its first
        ``<tspan>``.
    """
    tx_raw, ty_raw = text_elem.get("x"), text_elem.get("y")
    if tx_raw is None or ty_raw is None:
        for child in text_elem:
            if _local_tag(child) == "tspan":
                if tx_raw is None:
                    tx_raw = child.get("x")
                if ty_raw is None:
                    ty_raw = child.get("y")
                break
    try:
        tx = float(tx_raw) if tx_raw is not None else 0.0
    except ValueError:
        tx = 0.0
    try:
        ty = float(ty_raw) if ty_raw is not None else 0.0
    except ValueError:
        ty = 0.0
    return tx, ty


class SvgConverter:
    """Converts one SVG file into one PPTX slide."""

    def __init__(self, svg_path: str, verbose: bool = False) -> None:
        self.svg_path = svg_path
        self.verbose = verbose
        tree = etree.parse(svg_path)
        self.root = tree.getroot()
        vb = self.root.get("viewBox", "0 0 1200 675")
        self.cs = CoordSystem.from_viewbox(vb)
        self._shape_labels: Dict[int, List[Any]] = {}
        self._text_to_shape: Dict[int, Any] = {}
        self._defs: Dict[str, Any] = {}
        self._gradient_defs: Dict[str, Any] = {}
        self._connector_registry: List = []
        self._shape_registry: List = []
        self._pending_texts: List = []
        self._group_collect: Optional[List[Any]] = None

    def convert(self, prs: Presentation, slide_layout: Any) -> Any:
        self._connector_registry = []
        self._shape_registry = []
        self._shape_labels = {}
        self._text_to_shape = {}
        self._pending_texts = []
        self._group_collect = None
        slide = prs.slides.add_slide(slide_layout)
        self._resolve_defs()
        self._resolve_use_elements()
        self._compute_text_attachments()
        # Seed with the root element's own presentation attributes. An empty
        # dict discarded every style declared on `<svg>` itself, which is where
        # a document states the properties it means to apply throughout -- the
        # deck's font family above all. Per SVG, those attributes are inherited
        # by descendants.
        from .style_parser import compute_style as _compute_root_style
        self._dispatch_children(
            slide, self.root, _compute_root_style(self.root, {}))
        # Add all text boxes after shapes so they appear on top
        from .text_converter import add_textbox
        for text_elem, text_style in self._pending_texts:
            add_textbox(slide, text_elem, text_style, self.cs)
        self._bind_connectors(slide)
        return slide

    def _resolve_defs(self) -> None:
        """Index `<defs>` children by id, parsing gradients as they are found.

        Raises:
            ValueError: If a `<linearGradient>` is malformed or uses an
                unsupported feature. This runs before `_dispatch_children`, so
                the error reaches the caller rather than its per-child guard.
        """
        from .style_parser import parse_linear_gradient
        for elem in self.root.iter():
            if _local_tag(elem) == "defs":
                for child in elem:
                    eid = child.get("id")
                    if eid and _local_tag(child) == "linearGradient":
                        self._gradient_defs[eid] = parse_linear_gradient(child)
                    if eid:
                        self._defs[eid] = child

    def _resolve_use_elements(self) -> None:
        from copy import deepcopy
        for use_elem in list(self.root.iter()):
            if _local_tag(use_elem) != "use":
                continue
            href = use_elem.get("href") or use_elem.get(
                "{http://www.w3.org/1999/xlink}href", "")
            if not href.startswith("#"):
                continue
            ref_id = href[1:]
            ref = self._defs.get(ref_id)
            if ref is None:
                continue
            clone = deepcopy(ref)
            use_x = use_elem.get("x", "0")
            use_y = use_elem.get("y", "0")
            if use_x != "0" or use_y != "0":
                existing = clone.get("transform", "")
                clone.set("transform",
                          f"translate({use_x},{use_y}) {existing}".strip())
            parent = use_elem.getparent()
            if parent is not None:
                idx = list(parent).index(use_elem)
                parent.insert(idx, clone)
                parent.remove(use_elem)

    def _compute_text_attachments(self) -> None:
        """Attach text to the smallest semantic shape containing its baseline.

        A full-canvas rectangle is commonly emitted as a rendering background.
        It remains a native PPTX shape, but is deliberately excluded from this
        semantic attachment pass so labels outside smaller shapes stay editable
        standalone textboxes. Positions are still conservatively based on the
        SVG ``x``/``y`` attributes; transformed text is not inferred here.

        Slide furniture is excluded from it. A deck footer sits inside
        whatever panel reaches the bottom of the canvas, and a slide title
        inside a full-width header band, so a purely geometric rule hands them
        to that shape: the panel then renders the footer as its first line and
        takes its own first-paragraph offset from the footer's baseline, which
        put the panel heading roughly 400 units below where the SVG drew it.
        Role alone does not identify furniture -- a numbered-bullet badge
        carries `footnote` because it is set in 12/700 type -- so the shape
        must also be too tall to be that text's host.
        """
        from .shapes import _font_size_svg, _geometry
        from .style_parser import compute_style
        shape_bboxes: List[Tuple[Any, float, float, float, float]] = []
        for elem in self.root.iter():
            tag = _local_tag(elem)
            # Malformed geometry is not swallowed here: a shape missing from
            # this pass silently loses its label, which looks like an authoring
            # mistake in the SVG and is not one. `_geometry` raises
            # `SvgSourceError`, which the per-child guards re-raise.
            if tag == "rect":
                x = _geometry(elem, "x")
                y = _geometry(elem, "y")
                w = _geometry(elem, "width")
                h = _geometry(elem, "height")
                if self._is_canvas_background(x, y, w, h):
                    continue
                shape_bboxes.append((elem, x, y, w, h))
            elif tag == "circle":
                cx = _geometry(elem, "cx")
                cy = _geometry(elem, "cy")
                r = _geometry(elem, "r")
                shape_bboxes.append((elem, cx - r, cy - r, 2 * r, 2 * r))
            elif tag == "ellipse":
                cx = _geometry(elem, "cx")
                cy = _geometry(elem, "cy")
                rx = _geometry(elem, "rx")
                ry = _geometry(elem, "ry")
                shape_bboxes.append((elem, cx - rx, cy - ry, 2 * rx, 2 * ry))

        for text_elem in self.root.iter():
            if _local_tag(text_elem) != "text":
                continue
            tx, ty = _text_xy(text_elem)
            candidates = []
            for shape_elem, sx, sy, sw, sh in shape_bboxes:
                if sx <= tx <= sx + sw and sy <= ty <= sy + sh:
                    candidates.append((sw * sh, shape_elem, sh))
            if not candidates:
                continue
            candidates.sort(key=lambda c: c[0])
            _, best_shape, host_height = candidates[0]
            if text_elem.get("data-style-role") in _SLIDE_CHROME_ROLES:
                font_size = _font_size_svg(compute_style(text_elem, {}))
                if host_height > font_size * _CHROME_HOST_HEIGHT_RATIO:
                    continue
            self._shape_labels.setdefault(id(best_shape), []).append(text_elem)
            self._text_to_shape[id(text_elem)] = best_shape

    def _is_canvas_background(self, x: float, y: float,
                              width: float, height: float) -> bool:
        """Return whether a rectangle covers the converter's SVG canvas.

        The tolerance absorbs harmless decimal serialization noise while
        keeping rectangles that are materially smaller than the viewBox
        eligible as semantic text containers.
        """
        return all(
            math.isclose(value, expected, rel_tol=1e-9,
                         abs_tol=_CANVAS_TOLERANCE)
            for value, expected in (
                (x, 0.0),
                (y, 0.0),
                (width, self.cs.svg_w),
                (height, self.cs.svg_h),
            )
        )

    def _dispatch_children(self, slide: Any, parent: Any, inherited: Dict) -> None:
        from .style_parser import compute_style
        for elem in parent:
            tag = _local_tag(elem)
            try:
                style = compute_style(elem, inherited)
                self._dispatch_element(slide, elem, style)
            except (NativeObjectMarkerError, SvgSourceError):
                raise
            except Exception as exc:
                if self.verbose:
                    print(f"  [warn] {tag}: {exc}")

    def _resolve_paint_server(self, style: Dict) -> None:
        """Resolve a `fill="url(#id)"` reference into gradient style entries.

        Args:
            style: Computed style mapping, modified in place.

        Raises:
            SvgSourceError: If the fill names a paint server that no `<defs>`
                entry defines.
        """
        fill = style.get("fill", "")
        match = re.match(r"url\(#([^)]+)\)", fill.strip()) if fill else None
        if match is None:
            return
        gradient_id = match.group(1)
        if gradient_id not in self._gradient_defs:
            raise SvgSourceError(
                f"fill references unknown paint server {gradient_id!r}; "
                f"defined gradients: {sorted(self._gradient_defs)}"
            )
        stops, angle = self._gradient_defs[gradient_id]
        style["_gradient_stops"] = stops
        style["_gradient_angle"] = angle

    def _dispatch_element(self, slide: Any, elem: Any, style: Dict) -> None:
        tag = _local_tag(elem)
        self._resolve_paint_server(style)
        if tag in ("rect", "circle", "ellipse", "image"):
            from .shapes import dispatch_shape
            shape = dispatch_shape(
                slide, elem, style, self.cs, self._shape_labels.get(id(elem))
            )
            if shape is not None and self._group_collect is not None:
                self._group_collect.append(shape)
            if shape is not None and tag in ("rect", "circle", "ellipse"):
                # A malformed geometry attribute must not silently drop the
                # shape from the anchor registry: `_bind_connectors` would then
                # leave connectors dangling with no diagnostic, which is
                # precisely what the review gate is meant to catch.
                from .shapes import _geometry
                if tag == "rect":
                    bx = _geometry(elem, "x")
                    by = _geometry(elem, "y")
                    bw = _geometry(elem, "width")
                    bh = _geometry(elem, "height")
                else:
                    cx = _geometry(elem, "cx")
                    cy = _geometry(elem, "cy")
                    if tag == "circle":
                        r = _geometry(elem, "r")
                        bx, by, bw, bh = cx - r, cy - r, 2 * r, 2 * r
                    else:
                        rx = _geometry(elem, "rx")
                        ry = _geometry(elem, "ry")
                        bx, by, bw, bh = cx - rx, cy - ry, 2 * rx, 2 * ry
                self._shape_registry.append(
                    (elem, bx, by, bw, bh, shape.shape_id))
        elif tag == "text":
            if id(elem) not in self._text_to_shape:
                self._pending_texts.append((elem, style))
        elif tag in ("line", "polyline", "polygon"):
            from .connector import dispatch_connector
            conns = dispatch_connector(slide, elem, style, self.cs)
            self._connector_registry.extend(
                [(conn, elem) for conn in conns]
            )
            if self._group_collect is not None:
                self._group_collect.extend(conns)
        elif tag == "path":
            from .path_parser import parse_path
            from .path_to_pptx import add_path_shape
            d = elem.get("d", "")
            if d:
                path_shape = add_path_shape(slide, parse_path(d), self.cs, style)
                if path_shape is not None and self._group_collect is not None:
                    self._group_collect.append(path_shape)
        elif tag == "g":
            role = elem.get("data-pptx-role")
            if role in ("table", "chart"):
                self._dispatch_native_data(slide, elem, role)
            elif role == "group":
                self._dispatch_native_group(slide, elem, style)
            else:
                self._dispatch_children(slide, elem, style)

    def _bind_connectors(self, slide: Any) -> None:
        from .connector import build_anchor_map
        if not self._connector_registry or not self._shape_registry:
            return
        anchor_map = build_anchor_map(self._shape_registry)
        for conn, conn_elem in self._connector_registry:
            tag = _local_tag(conn_elem)
            if tag == "line":
                try:
                    begin_pt = (float(conn_elem.get("x1", 0)), float(conn_elem.get("y1", 0)))
                    end_pt = (float(conn_elem.get("x2", 0)), float(conn_elem.get("y2", 0)))
                    _try_bind(conn, begin_pt, end_pt, anchor_map)
                except Exception:
                    pass

    def _dispatch_native_data(self, slide: Any, elem: Any, role: str) -> None:
        """Materialize a ``data-pptx-role="table"/"chart"`` marker natively.

        Every failure in this path -- missing attributes, a malformed bbox, an
        unreadable or invalid sidecar, an unresolvable ``#<index>`` fragment,
        missing schema keys -- is raised as ``NativeObjectMarkerError``. A
        plain exception would be swallowed by ``_dispatch_children``'s
        per-child guard, silently dropping the table or chart from the deck
        while every gate still reported green.

        Args:
            slide: The python-pptx slide the object is added to.
            elem: The ``<g>`` element carrying the marker attributes.
            role: Either ``"table"`` or ``"chart"``.

        Raises:
            NativeObjectMarkerError: For any malformed or unusable marker.
        """
        source = elem.get("data-pptx-source")
        bbox_raw = elem.get("data-pptx-bbox")
        if not source or not bbox_raw:
            raise NativeObjectMarkerError(
                f'data-pptx-role="{role}" requires both data-pptx-source and '
                f"data-pptx-bbox (svg={self.svg_path})"
            )
        try:
            data = self._load_pptx_source(source)
            bbox = self._resolve_pptx_bbox(bbox_raw, role)
            style = self._pptx_style(elem)
            self._add_native_object(slide, role, source, data, bbox, style)
        except NativeObjectMarkerError:
            raise
        except Exception as exc:  # noqa: BLE001 - re-typed as a hard blocker
            raise NativeObjectMarkerError(
                f'data-pptx-role="{role}" with data-pptx-source={source!r} '
                f"could not be materialized ({type(exc).__name__}: {exc}) "
                f"(svg={self.svg_path})"
            ) from exc

    def _add_native_object(self, slide: Any, role: str, source: str,
                           data: Dict[str, Any], bbox: Tuple[int, int, int, int],
                           style: Dict[str, str]) -> None:
        """Build the native table/chart object described by ``data``."""
        import sys
        from pathlib import Path
        sys.path.insert(0, str(Path(__file__).resolve().parent.parent))
        import pptx_native

        if role == "table":
            self._require_source_keys(data, ("columns", "rows"), role, source)
            pptx_native.add_native_table(
                slide, data["columns"], data["rows"], bbox, style,
                highlight_col=data.get("highlight_col"),
            )
            return

        raw_type = data.get("type")
        if raw_type not in _CHART_TYPE_MAP:
            raise NativeObjectMarkerError(
                f'data-pptx-role="chart" source has unrecognized type '
                f"{raw_type!r} (expected one of {sorted(_CHART_TYPE_MAP)}) "
                f"(svg={self.svg_path})"
            )
        chart_type = _CHART_TYPE_MAP[raw_type]
        if chart_type == "pie":
            self._require_source_keys(data, ("categories", "values"), role, source)
            series = [{"label": data.get("title", ""), "values": data["values"]}]
            pptx_native.add_native_chart(
                slide, "pie", data["categories"], series, bbox, style,
                colors=data.get("colors"),
            )
        else:
            self._require_source_keys(data, ("categories", "series"), role, source)
            pptx_native.add_native_chart(
                slide, chart_type, data["categories"], data["series"], bbox,
                style, y_max=data.get("y_max"),
            )

    def _require_source_keys(self, data: Dict[str, Any], keys: Tuple[str, ...],
                             role: str, source: str) -> None:
        """Assert the resolved sidecar payload carries every required key."""
        if not isinstance(data, dict):
            raise NativeObjectMarkerError(
                f'data-pptx-role="{role}" source {source!r} must resolve to a '
                f"JSON object, got {type(data).__name__} (svg={self.svg_path})"
            )
        missing = [key for key in keys if key not in data]
        if missing:
            raise NativeObjectMarkerError(
                f'data-pptx-role="{role}" source {source!r} is missing required '
                f"key(s) {missing} (svg={self.svg_path})"
            )

    def _resolve_pptx_bbox(self, bbox_raw: str,
                           role: str) -> Tuple[int, int, int, int]:
        """Parse ``data-pptx-bbox`` ("x,y,w,h" in SVG user units) into EMU."""
        parts = [part.strip() for part in bbox_raw.split(",")]
        if len(parts) != 4:
            raise NativeObjectMarkerError(
                f'data-pptx-role="{role}" data-pptx-bbox must be "x,y,w,h", '
                f"got {bbox_raw!r} (svg={self.svg_path})"
            )
        try:
            bx, by, bw, bh = (float(part) for part in parts)
        except ValueError as exc:
            raise NativeObjectMarkerError(
                f'data-pptx-role="{role}" data-pptx-bbox has non-numeric '
                f"values: {bbox_raw!r} (svg={self.svg_path})"
            ) from exc
        return (self.cs.x(bx), self.cs.y(by), self.cs.x(bw), self.cs.y(bh))

    def _dispatch_native_group(self, slide: Any, elem: Any, inherited_style: Dict) -> None:
        """Convert a ``data-pptx-role="group"`` subtree into one native Group.

        Per-child failures are tolerated the same way ``_dispatch_children``
        tolerates them: one malformed child is skipped and the remaining
        children still become a group, rather than the whole node vanishing.
        ``NativeObjectMarkerError`` still propagates as a hard blocker.
        """
        from .style_parser import compute_style
        from .text_converter import add_textbox

        collected: List[Any] = []
        previous_sink = self._group_collect
        self._group_collect = collected
        local_texts: List = []
        try:
            for child in elem:
                tag = _local_tag(child)
                try:
                    style = compute_style(child, inherited_style)
                    if tag == "text" and id(child) not in self._text_to_shape:
                        local_texts.append((child, style))
                        continue
                    self._dispatch_element(slide, child, style)
                except (NativeObjectMarkerError, SvgSourceError):
                    raise
                except Exception as exc:  # noqa: BLE001 - skip one bad child
                    if self.verbose:
                        print(f"  [warn] {tag}: {exc}")
            for text_elem, text_style in local_texts:
                try:
                    box = add_textbox(slide, text_elem, text_style, self.cs)
                except Exception as exc:  # noqa: BLE001 - skip one bad label
                    if self.verbose:
                        print(f"  [warn] text: {exc}")
                    continue
                if box is not None:
                    collected.append(box)
        finally:
            self._group_collect = previous_sink

        if len(collected) >= 2:
            import sys
            from pathlib import Path
            sys.path.insert(0, str(Path(__file__).resolve().parent.parent))
            import pptx_native
            pptx_native.add_native_group(slide, collected)

    def _load_pptx_source(self, source: str) -> Dict[str, Any]:
        """Load the JSON payload referenced by ``data-pptx-source``.

        Accepts either a path relative to the SVG's own directory or that
        path plus a ``#<slide_index>`` fragment selecting one entry of a
        ``slide_data.json`` ``slides`` array.

        Args:
            source: The raw ``data-pptx-source`` attribute value.

        Returns:
            The resolved JSON object.

        Raises:
            NativeObjectMarkerError: If the path escapes the SVG's directory,
                the file cannot be read, the JSON is invalid, or the
                ``#<index>`` fragment does not resolve.
        """
        import json
        from pathlib import Path
        base_dir = Path(self.svg_path).resolve().parent
        filename, separator, index_str = source.partition("#")
        path = self._resolve_source_path(base_dir, filename, source)

        try:
            raw = path.read_text(encoding="utf-8")
        except OSError as exc:
            raise NativeObjectMarkerError(
                f"data-pptx-source {source!r} could not be read: {exc} "
                f"(svg={self.svg_path})"
            ) from exc
        try:
            payload = json.loads(raw)
        except json.JSONDecodeError as exc:
            raise NativeObjectMarkerError(
                f"data-pptx-source {source!r} is not valid JSON: {exc} "
                f"(svg={self.svg_path})"
            ) from exc

        if not separator:
            return payload

        try:
            target_index = int(index_str)
        except ValueError as exc:
            raise NativeObjectMarkerError(
                f"data-pptx-source {source!r} has a non-integer slide index "
                f"{index_str!r} (svg={self.svg_path})"
            ) from exc
        if not isinstance(payload, dict):
            raise NativeObjectMarkerError(
                f"data-pptx-source {source!r} must resolve to a JSON object "
                f"with a 'slides' array (svg={self.svg_path})"
            )
        for entry in payload.get("slides", []):
            if isinstance(entry, dict) and entry.get("index") == target_index:
                return entry
        raise NativeObjectMarkerError(
            f"data-pptx-source {source!r}: no slide with index {target_index} "
            f"in {filename} (svg={self.svg_path})"
        )

    def _resolve_source_path(self, base_dir: Path, filename: str,
                             source: str) -> Path:
        """Resolve a sidecar filename, keeping it inside the SVG's directory.

        A marker may only reference data shipped alongside the deck, so
        parent-directory traversal and absolute paths are rejected rather
        than silently loaded.
        """
        candidate = (base_dir / filename).resolve()
        try:
            candidate.relative_to(base_dir)
        except ValueError as exc:
            raise NativeObjectMarkerError(
                f"data-pptx-source {source!r} resolves to {candidate}, outside "
                f"the SVG's own directory {base_dir} (svg={self.svg_path})"
            ) from exc
        return candidate

    def _pptx_style(self, elem: Any) -> Dict[str, str]:
        """Resolve the color/font style for a native table or chart.

        Prefers an explicit ``data-pptx-style`` JSON attribute (the SVG
        producer's already-resolved style values, so custom project styles
        carry through); falls back to generate_slides.py's built-in defaults
        when the attribute is absent, which is the normal case for
        hand-authored Path C markup.

        Raises:
            NativeObjectMarkerError: If the attribute is present but is not a
                parseable JSON object -- falling back to defaults there would
                silently discard the producer's intended styling.
        """
        import json
        defaults = {
            "accent": "#1e3a5f", "white": "#ffffff", "card": "#f8fafc",
            "bg": "#ffffff", "body": "#374151", "good": "#059669",
            "danger": "#dc2626", "font": "'Helvetica Neue', Arial, sans-serif",
        }
        raw = elem.get("data-pptx-style")
        if raw is None:
            return defaults
        try:
            overrides = json.loads(raw)
        except (json.JSONDecodeError, TypeError) as exc:
            raise NativeObjectMarkerError(
                f"data-pptx-style is not valid JSON: {raw!r} "
                f"(svg={self.svg_path})"
            ) from exc
        if not isinstance(overrides, dict):
            raise NativeObjectMarkerError(
                f"data-pptx-style must be a JSON object, got "
                f"{type(overrides).__name__} (svg={self.svg_path})"
            )
        merged = dict(defaults)
        merged.update({k: v for k, v in overrides.items() if v})
        return merged


def _try_bind(conn: Any, begin_pt: tuple, end_pt: tuple, anchor_map: Dict) -> None:
    import math
    from .connector import THRESHOLD, bind_connector_end
    for is_begin, pt in [(True, begin_pt), (False, end_pt)]:
        best_dist = float("inf")
        best_sp_id = None
        best_idx = None
        for sp_id, anchors in anchor_map.items():
            for ax, ay, aidx in anchors:
                d = math.hypot(pt[0] - ax, pt[1] - ay)
                if d < best_dist:
                    best_dist = d
                    best_sp_id = sp_id
                    best_idx = aidx
        if best_dist <= THRESHOLD and best_sp_id is not None:
            bind_connector_end(conn, is_begin, best_sp_id, best_idx)


def convert_file(slides_dir: str, out_path: str, verbose: bool = False) -> None:
    prs = Presentation()
    prs.slide_width = Emu(PPTX_W)
    prs.slide_height = Emu(PPTX_H)
    layout = prs.slide_layouts[6]

    svg_files = sorted(Path(slides_dir).glob("slide*.svg"))
    if not svg_files:
        raise ValueError(f"No slide*.svg files found in {slides_dir}")

    for svg_path in svg_files:
        conv = SvgConverter(str(svg_path), verbose=verbose)
        conv.convert(prs, layout)
        if verbose:
            print(f"  + {svg_path.name}")

    Path(out_path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_path)
    print(f"\n{len(svg_files)} slide(s) → {out_path}")
