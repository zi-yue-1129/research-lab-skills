"""shapes.py — rect, circle, ellipse, image → python-pptx shapes."""
from __future__ import annotations

import re
from typing import Any, Dict, List, Optional, Tuple

from pptx.util import Emu, Pt
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN

from fonts import parse_font_stack, resolve_font_stack
from .style_parser import (apply_alpha, apply_dash, apply_fill,
                           apply_stroke, compute_style, resolve_color)
from .converter import CoordSystem, _local_tag, _text_xy

_ALIGN = {"start": PP_ALIGN.LEFT, "middle": PP_ALIGN.CENTER, "end": PP_ALIGN.RIGHT}
_ATTACHED_LABEL_TOP_PADDING_SVG = 4.0


def dispatch_shape(slide: Any, elem: Any, style: Dict,
                   cs: CoordSystem, label_elem: Optional[Any]) -> Optional[Any]:
    tag = _local_tag(elem)
    if tag == "rect":
        return _add_rect(slide, elem, style, cs, label_elem)
    elif tag in ("circle", "ellipse"):
        return _add_oval(slide, elem, style, cs, label_elem)
    elif tag == "image":
        return _add_image(slide, elem, style, cs)
    return None


_MSO_RECTANGLE = 1
_MSO_ROUNDED_RECTANGLE = 5
_MAX_CORNER_ADJUSTMENT = 0.5


def _corner_radius(elem: Any) -> float:
    """Return the SVG corner radius of a rect element.

    SVG permits `rx` alone, `ry` alone, or both; a single value mirrors to the
    other axis. PowerPoint's roundRect has one symmetric radius, so the larger
    of the two is used.

    Args:
        elem: The `<rect>` element.

    Returns:
        The corner radius in SVG units; 0 when the rect is sharp.

    Raises:
        ValueError: If a radius attribute is present but not a number.
    """
    values = []
    for attr in ("rx", "ry"):
        raw = elem.get(attr)
        if raw is None:
            continue
        try:
            values.append(abs(float(raw)))
        except ValueError:
            # A malformed radius must not silently become a sharp corner.
            raise ValueError(
                f"rect has non-numeric {attr}={raw!r}; fix the SVG source"
            )
    return max(values) if values else 0.0


def _add_rect(slide: Any, elem: Any, style: Dict,
              cs: CoordSystem, label_elem: Optional[Any]) -> Any:
    """Add one SVG rect to the slide as a native PPTX shape.

    A rect carrying `rx` or `ry` becomes a roundRect whose corner adjustment
    reproduces the authored radius, so token surface radii survive export. The
    former unconditional `add_shape(1, ...)` never read either attribute, so
    every rounded card in the preview arrived as a sharp box.

    Args:
        slide: Target PPTX slide.
        elem: The `<rect>` element.
        style: Computed style mapping for the element.
        cs: Coordinate system mapping SVG units to EMU.
        label_elem: Optional `<text>` element to write into the shape.

    Returns:
        The created PPTX shape.
    """
    svg_x = float(elem.get("x", 0))
    svg_y = float(elem.get("y", 0))
    svg_w = float(elem.get("width", 0))
    svg_h = float(elem.get("height", 0))
    x = cs.x(svg_x)
    y = cs.y(svg_y)
    w = max(1, cs.x(svg_w))
    h = max(1, cs.y(svg_h))

    radius = _corner_radius(elem)
    shorter = min(svg_w, svg_h)
    if radius > 0 and shorter > 0:
        shape = slide.shapes.add_shape(
            _MSO_ROUNDED_RECTANGLE, Emu(x), Emu(y), Emu(w), Emu(h))
        # The adjustment is the radius as a fraction of the shorter side,
        # clamped at a stadium so an oversized radius cannot overflow.
        shape.adjustments[0] = min(radius / shorter, _MAX_CORNER_ADJUSTMENT)
    else:
        shape = slide.shapes.add_shape(
            _MSO_RECTANGLE, Emu(x), Emu(y), Emu(w), Emu(h))

    apply_fill(shape, style.get("fill", "black"))
    apply_stroke(shape, style)
    apply_dash(shape, style)
    apply_alpha(shape, style)
    if label_elem is not None:
        _write_label(shape, label_elem, style, cs,
                     (svg_x, svg_y, svg_w, svg_h))
    return shape


def _add_oval(slide: Any, elem: Any, style: Dict,
              cs: CoordSystem, label_elem: Optional[Any]) -> Any:
    tag = _local_tag(elem)
    if tag == "circle":
        cx = float(elem.get("cx", 0))
        cy = float(elem.get("cy", 0))
        r = float(elem.get("r", 0))
        x, y, w, h = cx - r, cy - r, 2 * r, 2 * r
    else:
        cx = float(elem.get("cx", 0))
        cy = float(elem.get("cy", 0))
        rx = float(elem.get("rx", 0))
        ry = float(elem.get("ry", 0))
        x, y, w, h = cx - rx, cy - ry, 2 * rx, 2 * ry
    ex = cs.x(x)
    ey = cs.y(y)
    ew = max(1, cs.x(w))
    eh = max(1, cs.y(h))
    shape = slide.shapes.add_shape(9, Emu(ex), Emu(ey), Emu(ew), Emu(eh))
    apply_fill(shape, style.get("fill", "black"))
    apply_stroke(shape, style)
    apply_dash(shape, style)
    apply_alpha(shape, style)
    if label_elem is not None:
        _write_label(shape, label_elem, style, cs, (x, y, w, h))
    return shape


def _add_image(slide: Any, elem: Any, style: Dict, cs: CoordSystem) -> Optional[Any]:
    href = elem.get("href") or elem.get(
        "{http://www.w3.org/1999/xlink}href", "")
    if not href or href.startswith("data:"):
        return None
    x = cs.x(float(elem.get("x", 0)))
    y = cs.y(float(elem.get("y", 0)))
    w = max(1, cs.x(float(elem.get("width", 100))))
    h = max(1, cs.y(float(elem.get("height", 100))))
    try:
        return slide.shapes.add_picture(href, Emu(x), Emu(y), Emu(w), Emu(h))
    except Exception:
        return None


def _write_label(shape: Any, text_elem: Any, parent_style: Dict,
                 cs: CoordSystem, shape_bbox: Tuple[float, float, float, float]) -> None:
    """Write attached SVG labels into a deterministic native shape text frame.

    PowerPoint does not expose SVG baseline coordinates directly. The
    conservative mapping below uses a top-anchored text frame, explicit line
    heights, and paragraph spacing derived from each SVG baseline. The first
    paragraph's vertical offset is moved into a positive frame margin because
    LibreOffice can collapse first-paragraph spacing; that margin is bounded
    by the height remaining after later paragraphs. This preserves relative
    vertical placement and text-anchor alignment without turning labels into
    non-editable overlays; labels that overlap in SVG remain separate
    paragraphs and cannot be made pixel-perfect in every PowerPoint-compatible
    renderer.
    """
    tf = shape.text_frame
    tf.word_wrap = False
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    # Accept a single text element or a list
    elems = text_elem if isinstance(text_elem, list) else [text_elem]
    all_lines: List[Tuple[str, Dict, float]] = []
    for te in elems:
        all_lines.extend(_collect_text_lines(te, parent_style))

    shape_top = shape_bbox[1]
    base_space_befores: List[float] = []
    line_heights: List[float] = []
    previous_baseline: Optional[float] = None
    previous_line_height = 0.0
    content_height = 0.0
    for _, line_style, baseline_y in all_lines:
        font_size = _font_size_svg(line_style)
        line_height = max(font_size * 1.2, 1.0)
        if previous_baseline is None:
            space_before = max(
                0.0, baseline_y - shape_top - font_size * 0.75
            )
        else:
            space_before = max(
                0.0, baseline_y - previous_baseline - previous_line_height
            )
        base_space_befores.append(space_before)
        line_heights.append(line_height)
        content_height += space_before + line_height
        previous_baseline = baseline_y
        previous_line_height = line_height

    first_space_before = base_space_befores[0]
    remaining_shape_height = max(
        0.0, shape_bbox[3] - (content_height - first_space_before)
    )
    first_paragraph_margin = min(
        first_space_before + _ATTACHED_LABEL_TOP_PADDING_SVG,
        remaining_shape_height,
    )
    tf.margin_top = Emu(cs.y(first_paragraph_margin))

    for i, (text, line_style, _) in enumerate(all_lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        anchor = line_style.get("text-anchor", parent_style.get("text-anchor", "middle"))
        para.alignment = _ALIGN.get(anchor, PP_ALIGN.CENTER)
        space_before = base_space_befores[i]
        if i == 0:
            space_before = 0.0
        para.space_before = Emu(cs.y(space_before))
        para.space_after = Emu(0)
        para.line_spacing = Emu(cs.y(line_heights[i]))
        run = para.add_run()
        run.text = text
        _apply_font(run, line_style, parent_style)


def _collect_text_lines(text_elem: Any, parent_style: Dict) -> List[Tuple[str, Dict, float]]:
    """Collect attached text content with its SVG baseline coordinates."""
    lines: List[Tuple[str, Dict, float]] = []
    # Compute the text element's own style first so tspan children inherit correctly
    # (not from the enclosing shape's style, which is the parent_style here)
    text_style = compute_style(text_elem, parent_style)
    _, baseline_y = _text_xy(text_elem)
    direct_text = (text_elem.text or "").strip()
    if direct_text:
        lines.append((direct_text, text_style, baseline_y))
    current_baseline = baseline_y
    for tspan in text_elem:
        if _local_tag(tspan) == "tspan":
            ts = compute_style(tspan, text_style)
            if tspan.get("y") is not None:
                current_baseline = float(tspan.get("y"))
            else:
                current_baseline += float(tspan.get("dy", 0) or 0)
            t = (tspan.text or "").strip()
            if t:
                lines.append((t, ts, current_baseline))
    if not lines:
        lines.append(("", text_style, baseline_y))
    return lines


def _font_size_svg(style: Dict) -> float:
    """Return an attached-label font size in SVG user units."""
    size_raw = style.get("font-size", "14")
    try:
        return float(re.sub(r"[^0-9.]", "", str(size_raw)) or "14")
    except ValueError:
        return 14.0


def _apply_font(run: Any, style: Dict, parent_style: Dict) -> None:
    size_raw = style.get("font-size", parent_style.get("font-size", "14"))
    try:
        run.font.size = Pt(float(re.sub(r"[^0-9.]", "", size_raw) or "14"))
    except ValueError:
        run.font.size = Pt(14)
    weight = style.get("font-weight", parent_style.get("font-weight", "normal"))
    if weight in ("bold", "700", "800", "900"):
        run.font.bold = True
    fstyle = style.get("font-style", parent_style.get("font-style", "normal"))
    if fstyle == "italic":
        run.font.italic = True
    color_val = style.get("fill", parent_style.get("fill", "#000000"))
    rgb = resolve_color(color_val)
    if rgb:
        run.font.color.rgb = rgb
    ff = style.get("font-family", parent_style.get("font-family", ""))
    if ff and parse_font_stack(ff):
        # Resolve to a family that is actually installed. Splitting on
        # `[,\s]+` and taking the first token truncated every multi-word
        # family: the default stack `'Helvetica Neue', Arial, sans-serif`
        # became "Helvetica", a face installed on neither Linux nor the token
        # stack, so PowerPoint substituted silently and every text extent in
        # the deck diverged from the SVG preview it was built from.
        run.font.name = resolve_font_stack(ff)
