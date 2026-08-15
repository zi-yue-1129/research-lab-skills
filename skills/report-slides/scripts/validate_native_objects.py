#!/usr/bin/env python3
"""validate_native_objects.py -- safety-net gate for native PPTX table/chart/
group objects.

The primary defense is the agent-authoring rule (SKILL.md Stage 9.1 and the
data_visualization_worker_agent / architecture_diagram_worker_agent
instructions): author tables/charts/semantic node clusters with the
data-pptx-role SVG marker so svg_to_pptx/converter.py materializes real
PPTX objects. This script catches the two ways that rule can be missed:

  --svg-dir DIR   static pre-conversion scan: flags hand-drawn SVG shape
                  patterns (table grids, bar clusters, box+icon+label node
                  clusters) that are NOT wrapped in a data-pptx-role marker.
  --pptx FILE     post-conversion scan: opens the produced .pptx and flags
                  slides with an autoshape cluster matching those same
                  patterns and no accompanying native table/chart/group
                  object -- catching a converter bug that silently fell
                  through to shape-flattening despite a correctly
                  marked-up source SVG.

Both modes normalize geometry into the skill's reference canvas
(1200 x 675 SVG user units) and then apply the same structural detectors, so
a pattern is judged identically before and after conversion.

Design bias: precision over recall. The detectors require genuine structure
(a stacked grid of equally-sized cells; a baseline-aligned run of equal-width
bars with varying heights; a tight cluster of a labeled box plus an icon-like
shape) rather than "several shapes on one slide", because the skill's own
routine layouts -- metric cards, two-column panels, conclusion blocks -- are
multi-shape by construction and must never be blocked. Per the design spec,
false negatives here are an accepted residual risk.
"""
from __future__ import annotations

import argparse
import sys
from dataclasses import dataclass
from pathlib import Path
from typing import Any, Dict, List, Optional, Sequence, Tuple

import yaml
from lxml import etree

DEFAULT_THRESHOLDS_PATH = (
    Path(__file__).resolve().parent.parent / "references" / "native_object_thresholds.yaml"
)

#: Reference canvas (SVG user units) all geometry is normalized into. Matches
#: generate_slides.S["w"]/S["h"] and svg_to_pptx's default viewBox.
REFERENCE_CANVAS_W = 1200.0
REFERENCE_CANVAS_H = 675.0
REFERENCE_CANVAS_AREA = REFERENCE_CANVAS_W * REFERENCE_CANVAS_H

REQUIRED_THRESHOLDS: Tuple[str, ...] = (
    "min_grid_rows",
    "min_grid_cols",
    "max_row_gap_px",
    "min_bar_count",
    "min_cluster_shapes",
    "node_adjacency_px",
    "max_cluster_area_ratio",
    "max_icon_area_ratio",
    "chrome_extent_ratio",
    "position_tolerance_px",
)

_TABLE_REMEDIATION = (
    'Wrap the table in <g data-pptx-role="table" data-pptx-source="..." '
    'data-pptx-bbox="...">.'
)
_CHART_REMEDIATION = (
    'Wrap the chart in <g data-pptx-role="chart" data-pptx-source="..." '
    'data-pptx-bbox="...">.'
)
_GROUP_REMEDIATION = (
    'Wrap the node in <g data-pptx-role="group" data-node-id="...">.'
)
_TEXT_COORDS_REMEDIATION = (
    "Set x and y directly on the <text> element itself (matching every "
    "renderer in this skill, e.g. generate_slides.py's tlines()) -- do not "
    "rely on a <tspan> to carry the only copy of the position."
)


@dataclass(frozen=True)
class Finding:
    slide: str
    kind: str
    message: str


@dataclass(eq=False)
class _Box:
    """An axis-aligned box in reference-canvas units.

    Attributes:
        x: Left edge.
        y: Top edge.
        w: Width.
        h: Height.
        labeled: Whether a text label is tightly bound inside this box.
    """

    x: float
    y: float
    w: float
    h: float
    labeled: bool = False

    @property
    def right(self) -> float:
        """Right edge of the box."""
        return self.x + self.w

    @property
    def bottom(self) -> float:
        """Bottom edge of the box."""
        return self.y + self.h

    @property
    def area(self) -> float:
        """Area of the box."""
        return self.w * self.h

    def contains_point(self, px: float, py: float) -> bool:
        """Return whether the given point lies inside this box."""
        return self.x <= px <= self.right and self.y <= py <= self.bottom


def load_thresholds(path: Path = DEFAULT_THRESHOLDS_PATH) -> Dict[str, float]:
    """Load the configurable detection thresholds for the safety-net scan.

    Args:
        path: Path to the thresholds YAML file. Defaults to
            `references/native_object_thresholds.yaml`.

    Returns:
        A dict with exactly the keys listed in `REQUIRED_THRESHOLDS`.

    Raises:
        ValueError: If a required key is missing from the YAML document.
    """
    doc = yaml.safe_load(path.read_text(encoding="utf-8")) or {}
    thresholds: Dict[str, float] = {}
    for key in REQUIRED_THRESHOLDS:
        if key not in doc:
            raise ValueError(f"{path}: missing required key {key!r}")
        thresholds[key] = doc[key]
    return thresholds


def _local_tag(elem: Any) -> str:
    tag = elem.tag
    return tag.split("}")[-1] if isinstance(tag, str) and "}" in tag else str(tag)


def _has_pptx_role_ancestor(elem: Any) -> bool:
    node = elem.getparent()
    while node is not None:
        if node.get("data-pptx-role") in ("table", "chart", "group"):
            return True
        node = node.getparent()
    return False


def _close(a: float, b: float, tolerance: float) -> bool:
    return abs(a - b) <= tolerance


# ── Structural detectors (shared by the SVG and PPTX modes) ──────────────────

def _is_chrome(box: _Box, thresholds: Dict[str, float]) -> bool:
    """Return whether a box is full-bleed frame chrome rather than content.

    A canvas background fill, a header/footer bar, or a full-width divider
    spans essentially the whole canvas in one axis. Such shapes are
    decorative (`frame()` in generate_slides.py emits two of them on every
    slide) and would otherwise contribute phantom "rows" and "clusters" to
    every detector. Mirrors `converter._is_canvas_background`, generalized
    from "covers the whole canvas" to "spans one full axis".
    """
    ratio = thresholds["chrome_extent_ratio"]
    return (box.w >= ratio * REFERENCE_CANVAS_W
            or box.h >= ratio * REFERENCE_CANVAS_H)


def _content_boxes(boxes: Sequence[_Box],
                   thresholds: Dict[str, float]) -> List[_Box]:
    """Drop frame chrome, keeping only boxes that carry slide content."""
    return [box for box in boxes if not _is_chrome(box, thresholds)]


def _row_bands(boxes: Sequence[_Box], tol: float) -> List[List[_Box]]:
    """Group boxes into bands that share a top edge, ordered top to bottom."""
    bands: List[List[_Box]] = []
    for box in sorted(boxes, key=lambda b: (b.y, b.x)):
        for band in bands:
            if _close(band[0].y, box.y, tol):
                band.append(box)
                break
        else:
            bands.append([box])
    return bands


def _columns_of(band: Sequence[_Box], tol: float) -> Optional[Tuple[float, float, List[float]]]:
    """Find the widest run of equally-sized, non-overlapping boxes in a band.

    A genuine table row is N cells of identical size laid out side by side
    without overlap. A metric-card row (a card rect plus a full-width accent
    stripe sharing the card's left edge) collapses to the card columns only,
    because the stripe has a different height and duplicate x positions are
    merged into a single column.

    Args:
        band: Boxes that share a top edge.
        tol: Position tolerance in reference-canvas units.

    Returns:
        `(width, height, column_x_positions)` for the widest qualifying run,
        or None if no two same-sized boxes sit side by side.
    """
    best: Optional[Tuple[float, float, List[float]]] = None
    for anchor in band:
        same_size = [b for b in band
                     if _close(b.w, anchor.w, tol) and _close(b.h, anchor.h, tol)]
        columns: List[float] = []
        for box in sorted(same_size, key=lambda b: b.x):
            if columns and box.x - columns[-1] < anchor.w - tol:
                continue  # overlapping / duplicate x -> same column
            columns.append(box.x)
        if len(columns) >= 2 and (best is None or len(columns) > len(best[2])):
            best = (anchor.w, anchor.h, columns)
    return best


def _detect_grid(boxes: Sequence[_Box],
                 thresholds: Dict[str, float]) -> Optional[Tuple[int, int]]:
    """Detect a hand-drawn table grid among the given boxes.

    Requires a run of vertically stacked rows that (a) each hold at least
    `min_grid_cols` equally-sized, non-overlapping boxes, (b) share the same
    column positions and cell size, and (c) are separated by at most
    `max_row_gap_px` -- table cells are adjacent, whereas card and panel
    layouts are separated by visible gutters.

    Args:
        boxes: Content boxes (chrome already removed).
        thresholds: Detection thresholds as returned by `load_thresholds`.

    Returns:
        `(row_count, column_count)` of the detected grid, or None.
    """
    tol = thresholds["position_tolerance_px"]
    max_gap = thresholds["max_row_gap_px"]
    min_cols = thresholds["min_grid_cols"]

    rows: List[Tuple[float, float, float, List[float]]] = []
    for band in _row_bands(boxes, tol):
        candidate = _columns_of(band, tol)
        if candidate is None:
            continue
        width, height, columns = candidate
        if len(columns) >= min_cols:
            rows.append((band[0].y, width, height, columns))
    rows.sort(key=lambda row: row[0])

    best_run: List[Tuple[float, float, float, List[float]]] = []
    run: List[Tuple[float, float, float, List[float]]] = []
    for entry in rows:
        if run and _rows_stack(run[-1], entry, tol, max_gap):
            run.append(entry)
        else:
            run = [entry]
        if len(run) > len(best_run):
            best_run = list(run)

    if len(best_run) >= thresholds["min_grid_rows"]:
        return len(best_run), len(best_run[0][3])
    return None


def _rows_stack(previous: Tuple[float, float, float, List[float]],
                current: Tuple[float, float, float, List[float]],
                tol: float, max_gap: float) -> bool:
    """Return whether two grid rows are stacked cells of the same table."""
    prev_y, prev_w, prev_h, prev_cols = previous
    cur_y, cur_w, cur_h, cur_cols = current
    if not (_close(prev_w, cur_w, tol) and _close(prev_h, cur_h, tol)):
        return False
    if len(prev_cols) != len(cur_cols):
        return False
    if not all(_close(a, b, tol) for a, b in zip(prev_cols, cur_cols)):
        return False
    gap = cur_y - (prev_y + prev_h)
    return -tol <= gap <= max_gap


def _detect_bars(boxes: Sequence[_Box],
                 thresholds: Dict[str, float]) -> Optional[Tuple[int, float]]:
    """Detect a hand-drawn bar chart among the given boxes.

    Requires at least `min_bar_count` equal-width, non-overlapping boxes
    sharing a baseline (bottom edge) with at least two distinct heights --
    equal-height siblings are a card/panel row, not a bar chart.

    Args:
        boxes: Content boxes (chrome already removed).
        thresholds: Detection thresholds as returned by `load_thresholds`.

    Returns:
        `(bar_count, baseline_y)` of the detected bar run, or None.
    """
    tol = thresholds["position_tolerance_px"]
    min_bars = thresholds["min_bar_count"]

    baselines: List[List[_Box]] = []
    for box in sorted(boxes, key=lambda b: (b.bottom, b.x)):
        for group in baselines:
            if _close(group[0].bottom, box.bottom, tol):
                group.append(box)
                break
        else:
            baselines.append([box])

    for group in baselines:
        for anchor in group:
            same_width = [b for b in group if _close(b.w, anchor.w, tol)]
            bars: List[_Box] = []
            for box in sorted(same_width, key=lambda b: b.x):
                if bars and box.x - bars[-1].x < anchor.w - tol:
                    continue
                bars.append(box)
            if len(bars) < min_bars:
                continue
            heights = _distinct_values([b.h for b in bars], tol)
            if len(heights) > 1:
                return len(bars), group[0].bottom
    return None


def _distinct_values(values: Sequence[float], tol: float) -> List[float]:
    """Collapse near-equal values into distinct representatives."""
    distinct: List[float] = []
    for value in sorted(values):
        if not distinct or abs(value - distinct[-1]) > tol:
            distinct.append(value)
    return distinct


def _cluster_boxes(boxes: Sequence[_Box], adjacency: float) -> List[List[_Box]]:
    """Group boxes into clusters of shapes that touch or overlap."""
    remaining = list(boxes)
    clusters: List[List[_Box]] = []
    while remaining:
        seed = remaining.pop()
        cluster = [seed]
        changed = True
        while changed:
            changed = False
            for candidate in list(remaining):
                if any(_boxes_adjacent(candidate, member, adjacency)
                       for member in cluster):
                    remaining.remove(candidate)
                    cluster.append(candidate)
                    changed = True
        clusters.append(cluster)
    return clusters


def _boxes_adjacent(a: _Box, b: _Box, adjacency: float) -> bool:
    """Return whether two boxes overlap or sit within `adjacency` of each other."""
    return (a.x - adjacency <= b.right and b.x - adjacency <= a.right
            and a.y - adjacency <= b.bottom and b.y - adjacency <= a.bottom)


def _detect_node_clusters(boxes: Sequence[_Box],
                          thresholds: Dict[str, float]) -> List[List[_Box]]:
    """Detect ungrouped box+icon+label diagram nodes among the given boxes.

    A diagram node is a *tight* cluster: a labeled host shape with at least
    one icon-like companion (small in both axes relative to the host),
    covering only a small fraction of the canvas. This deliberately excludes
    the skill's routine layouts: a metric card's accent stripe spans the
    card's full width (not icon-like), and a two-column/conclusion panel
    covers roughly a third of the canvas (not a node).

    Args:
        boxes: Content boxes (chrome already removed), with `labeled` set.
        thresholds: Detection thresholds as returned by `load_thresholds`.

    Returns:
        The clusters that look like unmarked diagram nodes.
    """
    max_area = thresholds["max_cluster_area_ratio"] * REFERENCE_CANVAS_AREA
    max_icon_ratio = thresholds["max_icon_area_ratio"]
    detected: List[List[_Box]] = []

    for cluster in _cluster_boxes(boxes, thresholds["node_adjacency_px"]):
        if len(cluster) < thresholds["min_cluster_shapes"]:
            continue
        span_w = max(b.right for b in cluster) - min(b.x for b in cluster)
        span_h = max(b.bottom for b in cluster) - min(b.y for b in cluster)
        if span_w * span_h > max_area:
            continue
        labeled = [b for b in cluster if b.labeled]
        if not labeled:
            continue
        host = max(labeled, key=lambda b: b.area)
        if host.area <= 0:
            continue
        icons = [b for b in cluster
                 if b is not host
                 and b.area <= max_icon_ratio * host.area
                 and b.w <= 0.5 * host.w and b.h <= 0.5 * host.h]
        if icons:
            detected.append(cluster)
    return sorted(detected, key=lambda c: (min(b.y for b in c), min(b.x for b in c)))


# ── SVG-source scan ──────────────────────────────────────────────────────────

def _viewbox_scale(root: Any) -> Tuple[float, float]:
    """Return the (x, y) factors mapping SVG units into the reference canvas."""
    parts = (root.get("viewBox") or "").strip().split()
    if len(parts) == 4:
        try:
            width = float(parts[2])
            height = float(parts[3])
            if width > 0 and height > 0:
                return REFERENCE_CANVAS_W / width, REFERENCE_CANVAS_H / height
        except ValueError:
            pass
    return 1.0, 1.0


def _element_box(elem: Any, sx: float, sy: float) -> Optional[_Box]:
    """Return the normalized bounding box of a primitive SVG shape element."""
    tag = _local_tag(elem)
    try:
        if tag == "rect":
            x, y = float(elem.get("x", 0)), float(elem.get("y", 0))
            w, h = float(elem.get("width", 0)), float(elem.get("height", 0))
        elif tag == "circle":
            cx, cy = float(elem.get("cx", 0)), float(elem.get("cy", 0))
            r = float(elem.get("r", 0))
            x, y, w, h = cx - r, cy - r, 2 * r, 2 * r
        elif tag == "ellipse":
            cx, cy = float(elem.get("cx", 0)), float(elem.get("cy", 0))
            rx, ry = float(elem.get("rx", 0)), float(elem.get("ry", 0))
            x, y, w, h = cx - rx, cy - ry, 2 * rx, 2 * ry
        elif tag == "image":
            x, y = float(elem.get("x", 0)), float(elem.get("y", 0))
            w, h = float(elem.get("width", 0)), float(elem.get("height", 0))
        else:
            return None
    except ValueError:
        return None
    if w <= 0 or h <= 0:
        return None
    return _Box(x * sx, y * sy, w * sx, h * sy)


def scan_svg(svg_path: Path, thresholds: Dict[str, float]) -> List[Finding]:
    """Flag hand-drawn table/chart/node patterns missing a data-pptx-role marker.

    Parses the given SVG source and looks for three unmarked hand-drawn
    patterns: a stacked grid of equally-sized rects (table_like), a
    baseline-aligned run of equal-width rects with varying heights
    (chart_like), and a `<g>` holding a tight box+icon+label cluster
    (node_cluster_like). Any element nested under an ancestor carrying
    `data-pptx-role="table"|"chart"|"group"` is treated as already handled
    and excluded from detection.

    Args:
        svg_path: Path to the SVG file to scan.
        thresholds: Detection thresholds as returned by `load_thresholds`.

    Returns:
        A list of `Finding`s for every unmarked pattern detected in this
        SVG; empty if the SVG is clean.
    """
    tree = etree.parse(str(svg_path))
    root = tree.getroot()
    sx, sy = _viewbox_scale(root)
    slide_name = svg_path.name

    rect_boxes: List[_Box] = []
    for elem in root.iter():
        if _local_tag(elem) != "rect" or _has_pptx_role_ancestor(elem):
            continue
        box = _element_box(elem, sx, sy)
        if box is not None:
            rect_boxes.append(box)
    content = _content_boxes(rect_boxes, thresholds)

    findings: List[Finding] = []
    grid = _detect_grid(content, thresholds)
    if grid is not None:
        rows, columns = grid
        findings.append(Finding(
            slide_name, "table_like",
            f"{rows} stacked rows of {columns} equally-sized rects look like a "
            f"hand-drawn table. {_TABLE_REMEDIATION}",
        ))
    bars = _detect_bars(content, thresholds)
    if bars is not None:
        count, baseline = bars
        findings.append(Finding(
            slide_name, "chart_like",
            f"{count} equal-width rects sharing baseline y={baseline:.0f} with "
            f"varying heights look like a hand-drawn bar chart. "
            f"{_CHART_REMEDIATION}",
        ))
    findings.extend(_scan_svg_node_clusters(slide_name, root, thresholds, sx, sy))
    findings.extend(_scan_missing_text_coordinates(slide_name, root))
    return findings


def _scan_missing_text_coordinates(slide_name: str, root: Any) -> List[Finding]:
    """Flag `<text>` elements missing `x` or `y` on the element itself.

    Every renderer in this skill writes `x`/`y` directly on the `<text>`
    tag, which `svg_to_pptx`'s label-attachment and standalone-textbox code
    (`_compute_text_attachments`, `_collect_text_lines`, `add_textbox`)
    reads to decide which shape a label belongs to and where it lands.
    Coordinates set only on a child `<tspan>` are valid SVG and the
    converter now tolerates them via a fallback (`_text_xy`), but that
    fallback is a second line of defense, not the authoring convention --
    this is the one place that convention is enforced.

    Args:
        slide_name: Name of the SVG file being scanned, used as the
            `Finding.slide` label.
        root: The parsed SVG document's root element.

    Returns:
        A `Finding` for every `<text>` element missing `x` and/or `y` on
        itself; empty if every `<text>` in the document is compliant.
    """
    findings: List[Finding] = []
    for elem in root.iter():
        if _local_tag(elem) != "text":
            continue
        missing = [attr for attr in ("x", "y") if elem.get(attr) is None]
        if not missing:
            continue
        snippet = "".join(elem.itertext()).strip()[:40]
        findings.append(Finding(
            slide_name, "text_missing_coords",
            f"<text> is missing {' and '.join(missing)} on the element "
            f"itself (content: {snippet!r}). {_TEXT_COORDS_REMEDIATION}",
        ))
    return findings


def _scan_svg_node_clusters(slide_name: str, root: Any,
                            thresholds: Dict[str, float],
                            sx: float, sy: float) -> List[Finding]:
    """Flag `<g>` elements that hold an unmarked box+icon+label node cluster."""
    max_area = thresholds["max_cluster_area_ratio"] * REFERENCE_CANVAS_AREA
    findings: List[Finding] = []
    for group in root.iter():
        if _local_tag(group) != "g" or group.get("data-pptx-role") is not None:
            continue
        if _has_pptx_role_ancestor(group):
            continue
        shape_children = [c for c in group if _local_tag(c) in
                          ("rect", "circle", "ellipse", "path", "image")]
        if len(shape_children) < thresholds["min_cluster_shapes"]:
            continue
        boxes = [box for box in (_element_box(c, sx, sy) for c in shape_children)
                 if box is not None]
        content = _content_boxes(boxes, thresholds)
        if not content:
            continue
        span_w = max(b.right for b in content) - min(b.x for b in content)
        span_h = max(b.bottom for b in content) - min(b.y for b in content)
        if span_w * span_h > max_area:
            continue
        if not _has_contained_label(group, content, sx, sy):
            continue
        findings.append(Finding(
            slide_name, "node_cluster_like",
            f"<g> with {len(shape_children)} shapes and a label bound inside "
            f"one of them looks like an unmarked diagram node. "
            f"{_GROUP_REMEDIATION}",
        ))
    return findings


def _has_contained_label(group: Any, boxes: Sequence[_Box],
                         sx: float, sy: float) -> bool:
    """Return whether a `<text>` child sits inside one of the group's shapes.

    A label merely *present* alongside several shapes is not evidence of a
    node; a label anchored *inside* one specific shape is.
    """
    for child in group:
        if _local_tag(child) != "text":
            continue
        try:
            tx = float(child.get("x", 0)) * sx
            ty = float(child.get("y", 0)) * sy
        except ValueError:
            continue
        if any(box.contains_point(tx, ty) for box in boxes):
            return True
    return False


# ── Produced-PPTX scan ───────────────────────────────────────────────────────

def scan_pptx(pptx_path: Path, thresholds: Dict[str, float]) -> List[Finding]:
    """Flag autoshape patterns in a produced PPTX with no native counterpart.

    Opens the given .pptx and, for each slide, normalizes top-level autoshape
    geometry back into the reference canvas and applies the same three
    structural detectors `scan_svg` applies to source SVGs. A grid is only
    reported when the slide has no native table, a bar run only when it has
    no native chart; node clusters are reported per cluster (shapes already
    inside a native Group are not top-level and so never reach the detector).

    Args:
        pptx_path: Path to the produced .pptx file to scan.
        thresholds: Detection thresholds as returned by `load_thresholds`.

    Returns:
        A list of `Finding`s, one per flagged slide/pattern combination;
        empty if every autoshape pattern has an accompanying native object.
    """
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(str(pptx_path))
    scale_x = REFERENCE_CANVAS_W / float(prs.slide_width)
    scale_y = REFERENCE_CANVAS_H / float(prs.slide_height)
    findings: List[Finding] = []

    for slide_index, slide in enumerate(prs.slides, start=1):
        has_native_table = any(getattr(s, "has_table", False) for s in slide.shapes)
        has_native_chart = any(getattr(s, "has_chart", False) for s in slide.shapes)

        boxes: List[_Box] = []
        for shape in slide.shapes:
            if shape.shape_type != MSO_SHAPE_TYPE.AUTO_SHAPE:
                continue
            if None in (shape.left, shape.top, shape.width, shape.height):
                continue
            boxes.append(_Box(
                float(shape.left) * scale_x, float(shape.top) * scale_y,
                float(shape.width) * scale_x, float(shape.height) * scale_y,
                labeled=_shape_has_label(shape),
            ))
        content = _content_boxes(boxes, thresholds)
        slide_label = f"slide{slide_index}"

        grid = _detect_grid(content, thresholds)
        if grid is not None and not has_native_table:
            rows, columns = grid
            findings.append(Finding(
                slide_label, "table_like",
                f"{rows} stacked rows of {columns} equally-sized autoshapes and "
                "no native table object -- converter likely fell through to "
                "shape-flattening.",
            ))
        bars = _detect_bars(content, thresholds)
        if bars is not None and not has_native_chart:
            count, baseline = bars
            findings.append(Finding(
                slide_label, "chart_like",
                f"{count} equal-width autoshapes sharing baseline y={baseline:.0f} "
                "with varying heights and no native chart object -- converter "
                "likely fell through to shape-flattening.",
            ))
        for cluster in _detect_node_clusters(content, thresholds):
            findings.append(Finding(
                slide_label, "node_cluster_like",
                f"{len(cluster)} ungrouped autoshapes form a labeled box with an "
                "icon -- if these are one semantic diagram node, they should be "
                "a native Group.",
            ))

    return findings


def _shape_has_label(shape: Any) -> bool:
    """Return whether an autoshape carries non-empty text of its own."""
    if not getattr(shape, "has_text_frame", False):
        return False
    return bool(shape.text_frame.text.strip())


def main(argv: Sequence[str] | None = None) -> int:
    """CLI entry point for validate_native_objects.py.

    Parses `--svg-dir`/`--pptx` (mutually exclusive, one required) and an
    optional `--thresholds` override, runs the corresponding scan, and
    prints results: `OK: ...` on stdout and exit 0 when no findings, or
    `ERROR <slide> [<kind>]: <message>` lines on stderr and exit 1 when
    findings exist.

    Args:
        argv: Command-line arguments to parse, excluding the program name.
            Defaults to `sys.argv[1:]` via argparse when None.

    Returns:
        Process exit code: 0 if no findings were flagged, 1 otherwise.
    """
    parser = argparse.ArgumentParser(
        description="Detect hand-drawn table/chart/node patterns missing "
                     "native PPTX table/chart/group objects."
    )
    targets = parser.add_mutually_exclusive_group(required=True)
    targets.add_argument("--svg-dir", type=Path,
                         help="Directory of slide*.svg files (pre-conversion scan).")
    targets.add_argument("--pptx", type=Path,
                         help="Produced .pptx file (post-conversion scan).")
    parser.add_argument("--thresholds", type=Path, default=DEFAULT_THRESHOLDS_PATH)
    args = parser.parse_args(argv)

    thresholds = load_thresholds(args.thresholds)

    if args.svg_dir is not None:
        findings: List[Finding] = []
        for svg_path in sorted(args.svg_dir.glob("slide*.svg")):
            findings.extend(scan_svg(svg_path, thresholds))
    else:
        findings = scan_pptx(args.pptx, thresholds)

    if findings:
        for finding in findings:
            print(f"ERROR {finding.slide} [{finding.kind}]: {finding.message}",
                 file=sys.stderr)
        return 1

    print("OK: no unmarked table/chart/node patterns found")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
