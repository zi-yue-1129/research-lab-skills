"""pptx_native.py -- construct real PPTX table, chart, and group objects.

Used exclusively by svg_to_pptx/converter.py's data-pptx-role dispatch
branch, so both the deterministic Python renderer (generate_slides.py) and
agent-authored SVG materialize true native PowerPoint objects instead of
disconnected shapes.
"""
from __future__ import annotations

from typing import Any, Dict, List, Optional, Sequence, Tuple

from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Pt

_DEFAULT_SERIES_COLORS = ("#3b82f6", "#059669", "#d97706", "#dc2626", "#8a2be2")


def _rgb(hex_value: Optional[str], fallback: str = "#000000") -> RGBColor:
    """Convert hex color string to RGBColor.

    Args:
        hex_value: Hex color string (e.g., "#ff0000" or "ff0000")
        fallback: Default color if hex_value is None

    Returns:
        RGBColor instance
    """
    value = (hex_value or fallback).lstrip("#")
    if len(value) == 3:
        value = "".join(ch * 2 for ch in value)
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def _font_family(css_font: str) -> str:
    """Extract font family name from CSS font string.

    Args:
        css_font: CSS font string (e.g., "'Helvetica Neue', Arial, sans-serif")

    Returns:
        First font family name, or "Calibri" as fallback
    """
    first = css_font.split(",")[0].strip().strip("'\"")
    return first or "Calibri"


def add_native_table(
    slide: Any,
    columns: List[str],
    rows: List[List[Any]],
    bbox: Tuple[int, int, int, int],
    style: Dict[str, str],
    highlight_col: Optional[int] = None,
) -> Any:
    """Create a real PPTX table (``GraphicFrame.has_table``) at ``bbox`` (EMU).

    Args:
        slide: PowerPoint slide object
        columns: List of column header names
        rows: List of row data (each row is a list of values)
        bbox: Bounding box as (left, top, width, height) in EMU units
        style: Style dictionary with color and font keys
        highlight_col: Optional column index to apply +/- coloring to

    Returns:
        GraphicFrame containing the table
    """
    left, top, width, height = bbox
    n_cols = len(columns)
    n_rows = len(rows) + 1
    graphic_frame = slide.shapes.add_table(
        n_rows, n_cols, Emu(left), Emu(top), Emu(width), Emu(height)
    )
    table = graphic_frame.table
    table.first_row = False
    table.horz_banding = False

    header_fill = _rgb(style.get("accent"), "#1e3a5f")
    header_text = _rgb(style.get("white"), "#ffffff")
    card_fill = _rgb(style.get("card"), "#f8fafc")
    bg_fill = _rgb(style.get("bg"), "#ffffff")
    body_text = _rgb(style.get("body"), "#374151")
    good_text = _rgb(style.get("good"), "#059669")
    danger_text = _rgb(style.get("danger"), "#dc2626")
    font_name = _font_family(style.get("font", "Calibri"))

    for col_index, column_title in enumerate(columns):
        _set_cell_text(table.cell(0, col_index), str(column_title),
                       header_fill, header_text, font_name, bold=True)

    for row_index, row in enumerate(rows):
        fill = card_fill if row_index % 2 == 0 else bg_fill
        for col_index, value in enumerate(row):
            text_color = body_text
            cell_text = str(value)
            if highlight_col is not None and col_index == highlight_col:
                if "+" in cell_text:
                    text_color = good_text
                elif "-" in cell_text:
                    text_color = danger_text
            _set_cell_text(table.cell(row_index + 1, col_index), cell_text,
                           fill, text_color, font_name)

    return graphic_frame


def _set_cell_text(cell: Any, text: str, fill: RGBColor, text_color: RGBColor,
                   font_name: str, bold: bool = False) -> None:
    """Set text and styling for a table cell.

    Args:
        cell: PowerPoint table cell object
        text: Text content
        fill: Background color
        text_color: Text color
        font_name: Font family name
        bold: Whether text should be bold
    """
    cell.fill.solid()
    cell.fill.fore_color.rgb = fill
    text_frame = cell.text_frame
    text_frame.clear()
    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(13)
    run.font.bold = bold
    run.font.color.rgb = text_color
    if font_name:
        run.font.name = font_name


def add_native_chart(
    slide: Any,
    chart_type: str,
    categories: List[str],
    series: List[Dict[str, Any]],
    bbox: Tuple[int, int, int, int],
    style: Dict[str, str],
    y_max: Optional[float] = None,
    colors: Optional[List[str]] = None,
) -> Any:
    """Create a real PPTX chart (``GraphicFrame.has_chart``) at ``bbox`` (EMU).

    ``chart_type`` is one of "bar", "line", "pie". For "pie", ``series`` must
    contain exactly one entry and ``colors`` (one hex color per category) is
    used instead of per-series coloring.

    Args:
        slide: PowerPoint slide object
        chart_type: Chart type ("bar", "line", or "pie")
        categories: List of category names
        series: List of series data, each with "label", "values", and optional "color"
        bbox: Bounding box as (left, top, width, height) in EMU units
        style: Style dictionary with color and font keys
        y_max: Optional maximum for Y axis (not used for pie charts)
        colors: Optional list of colors for pie chart slices

    Returns:
        GraphicFrame containing the chart
    """
    left, top, width, height = bbox
    chart_data = CategoryChartData()
    chart_data.categories = categories
    for series_index, entry in enumerate(series):
        values = [float(v) for v in entry.get("values", [])]
        chart_data.add_series(entry.get("label", f"Series {series_index + 1}"), values)

    xl_type = {
        "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "line": XL_CHART_TYPE.LINE_MARKERS,
        "pie": XL_CHART_TYPE.PIE,
    }[chart_type]

    graphic_frame = slide.shapes.add_chart(
        xl_type, Emu(left), Emu(top), Emu(width), Emu(height), chart_data
    )
    chart = graphic_frame.chart
    chart.has_title = False
    chart.font.name = _font_family(style.get("font", "Calibri"))
    chart.has_legend = len(series) > 1 or chart_type == "pie"
    if chart.has_legend:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False

    if chart_type == "pie":
        point_colors = colors or [
            _DEFAULT_SERIES_COLORS[i % len(_DEFAULT_SERIES_COLORS)]
            for i in range(len(categories))
        ]
        for point_index, point in enumerate(chart.series[0].points):
            point.format.fill.solid()
            point.format.fill.fore_color.rgb = _rgb(point_colors[point_index])
    else:
        for series_index, (entry, chart_series) in enumerate(zip(series, chart.series)):
            color = entry.get("color") or _DEFAULT_SERIES_COLORS[
                series_index % len(_DEFAULT_SERIES_COLORS)
            ]
            if chart_type == "bar":
                chart_series.format.fill.solid()
                chart_series.format.fill.fore_color.rgb = _rgb(color)
            else:  # line
                chart_series.format.line.color.rgb = _rgb(color)
                chart_series.format.line.width = Pt(2.25)

    if y_max is not None and chart_type != "pie":
        chart.value_axis.maximum_scale = float(y_max)
        chart.value_axis.minimum_scale = 0.0

    return graphic_frame


def add_native_group(slide: Any, shapes: Sequence[Any]) -> Optional[Any]:
    """Group already-placed shapes into one native PPTX Group shape.

    Returns None (and groups nothing) for fewer than 2 shapes -- grouping a
    single shape has no PowerPoint-visible effect and is not a meaningful
    semantic unit for this design.

    Args:
        slide: PowerPoint slide object
        shapes: Sequence of shape objects to group

    Returns:
        GroupShape if len(shapes) >= 2, else None
    """
    if len(shapes) < 2:
        return None
    return slide.shapes.add_group_shape(shapes)
