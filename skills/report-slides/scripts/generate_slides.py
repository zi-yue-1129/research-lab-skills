#!/usr/bin/env python3
"""generate_slides.py — Research Report SVG Slide Generator (Path A)

Usage:
    python scripts/generate_slides.py --data slide_data.json --out ./output/
    python scripts/generate_slides.py --data slide_data.json --out ./output/ --slide 5

Slide types:
    title, bullet_list, bar_chart, table, metric_cards, two_column, timeline, conclusion
"""

import json
import argparse
import os
from pathlib import Path
from typing import Any, Dict, List, Optional, Sequence, Tuple

from design_tokens import DEFAULT_TOKENS_PATH, DesignTokens, TokenError, TypeRole
from fonts import resolve_font_stack, text_width, vertical_metrics
from presentation_gates import ProductionGateError, assert_production_allowed
from presentation_state import find_project_root


# ── Style constants ───────────────────────────────────────────────────────────

S = {
    "w":      1200,
    "h":      675,
    "bg":     "#ffffff",
    "accent": "#1e3a5f",
    "good":   "#059669",
    "warn":   "#d97706",
    "danger": "#dc2626",
    "blue":   "#3b82f6",
    "card":   "#f8fafc",
    "border": "#e2e8f0",
    "body":   "#374151",
    "muted":  "#64748b",
    "white":  "#ffffff",
    "font":   "'Helvetica Neue', Arial, sans-serif",
    "top_bar_h": 6,
}

TYPE: Dict[str, TypeRole] = {}
_TOKENS: Optional[DesignTokens] = None


class SlideCapacityError(ValueError):
    """Raised when content cannot be rendered legibly at presentation scale.

    Shrinking type below a token role, or clipping rows, would hide the problem
    rather than solve it. The caller must split the content across slides or
    reduce it.
    """


def apply_tokens(tokens_path: Optional[Path]) -> None:
    """Load a design-token file into the module-level style state.

    Args:
        tokens_path: Path to a `.tokens.yaml` file, or None for the shipped
            default contract.

    Raises:
        TokenError: If the token file is missing or schema-invalid.
        FontError: If no family in the token font stack is installed.
    """
    global _TOKENS
    _TOKENS = DesignTokens.load(tokens_path or DEFAULT_TOKENS_PATH)
    TYPE.clear()
    for role in _TOKENS.raw["typography"]["roles"]:
        TYPE[role] = _TOKENS.type_role(role)
    S["w"] = _TOKENS.raw["canvas"]["width"]
    S["h"] = _TOKENS.raw["canvas"]["height"]
    S["grid"] = _TOKENS.raw["canvas"]["grid"]
    S["safe"] = dict(_TOKENS.raw["canvas"]["safe_area"])
    S["spacing"] = list(_TOKENS.raw["spacing"]["scale"])
    for role in ("bg", "body", "muted", "card", "primary",
                 "positive", "warn", "danger", "line", "divider"):
        S[role] = _TOKENS.color(role)
    # Legacy key names still referenced by the renderers.
    S["accent"] = _TOKENS.color("primary")
    S["good"] = _TOKENS.color("positive")
    S["border"] = _TOKENS.color("divider")
    S["blue"] = _TOKENS.raw["chart"]["palette"][0]
    S["chart_palette"] = list(_TOKENS.raw["chart"]["palette"])
    S["white"] = "#ffffff"
    S["font"] = _TOKENS.font_stack("sans")
    S["font_resolved"] = resolve_font_stack(S["font"])
    S["top_bar_h"] = 0


def _role(role: str) -> TypeRole:
    """Return one resolved type role.

    Args:
        role: Role key such as `body`.

    Returns:
        The resolved `TypeRole`.

    Raises:
        RuntimeError: If `apply_tokens` has not been called.
        TokenError: If the role is undefined.
    """
    if not TYPE:
        raise RuntimeError(
            "apply_tokens() must be called before rendering; "
            "type roles are not loaded"
        )
    if role not in TYPE:
        raise TokenError(
            f"undefined type role {role!r}; defined roles: {sorted(TYPE)}"
        )
    return TYPE[role]


def t_size(role: str) -> float:
    """Return the font size for a type role.

    Args:
        role: Role key such as `body`.

    Returns:
        The role's font size in SVG units.
    """
    return _role(role).size


def t_weight(role: str) -> int:
    """Return the numeric font weight for a type role.

    Args:
        role: Role key such as `body`.

    Returns:
        The role's CSS numeric font weight.
    """
    return _role(role).weight


def t_lh(role: str) -> float:
    """Return the line-height multiplier for a type role.

    Args:
        role: Role key such as `body`.

    Returns:
        The role's line-height multiplier.
    """
    return _role(role).line_height


# ── Style loading ─────────────────────────────────────────────────────────────

def _parse_frontmatter(path: str) -> dict:
    """Parse YAML frontmatter from a .md file without requiring PyYAML."""
    try:
        with open(path, encoding="utf-8") as f:
            lines = f.readlines()
    except OSError as e:
        print(f"  [style] Cannot read {path}: {e}")
        return {}
    if not lines or lines[0].strip() != "---":
        return {}
    fm: dict = {}
    for line in lines[1:]:
        stripped = line.rstrip("\n")
        if stripped.strip() == "---":
            break
        if ":" not in stripped:
            continue
        key, _, val = stripped.partition(":")
        key = key.strip()
        val = val.strip().strip('"').strip("'")
        if key and val:
            fm[key] = val
    return fm


# Style Markdown keys are historical; the token roles are the contract. This map
# is the only place the two vocabularies meet. `border` resolves to the
# `divider` role because that is what it has always meant: `apply_tokens` sets
# `S["border"]` from `color.roles.divider`, and there is no `border` role.
STYLE_KEY_TO_ROLE: Dict[str, str] = {
    "primary": "primary", "bg": "bg", "body": "body", "muted": "muted",
    "border": "divider", "card": "card", "positive": "positive",
    "warn": "warn", "danger": "danger", "font": "font",
}

# Frontmatter keys that describe the style itself rather than a colour role.
# They are part of the documented schema, so rejecting them as unknown roles
# would make every shipped style file a hard error.
STYLE_METADATA_KEYS = frozenset({"name", "description"})


def effective_tokens(
    tokens_path: Optional[Path], style_path: Optional[str], out_dir: Path
) -> Tuple[Path, str]:
    """Compose tokens with a style override and write the result.

    The composed file is the single artifact the renderer loads, the linter
    reads, and the workflow gate digests. Composing before anything reads the
    tokens is what keeps those three from disagreeing.

    Args:
        tokens_path: Token file, or `None` for the shipped default.
        style_path: Style Markdown file, or `None`.
        out_dir: Directory to write `_effective.tokens.yaml` into.

    Returns:
        The written path and the composed set's digest.

    Raises:
        TokenError: If the style names an unknown role or the composed set
            fails validation.
        ValueError: If the style file has no usable frontmatter.
    """
    tokens = DesignTokens.load(tokens_path or DEFAULT_TOKENS_PATH)
    if style_path is not None:
        frontmatter = _parse_frontmatter(style_path)
        if not frontmatter:
            raise ValueError(
                f"style file {style_path} has no usable YAML frontmatter; "
                f"expected keys such as primary/bg/body "
                f"(see references/styles/STYLES.md)"
            )
        unknown = sorted(
            set(frontmatter) - set(STYLE_KEY_TO_ROLE) - STYLE_METADATA_KEYS)
        if unknown:
            raise TokenError(
                f"style file {style_path} sets {', '.join(unknown)}, which "
                f"name no colour role"
            )
        overrides = {
            STYLE_KEY_TO_ROLE[key]: value
            for key, value in frontmatter.items() if key in STYLE_KEY_TO_ROLE
        }
        if "font" in overrides:
            overrides["font"] = resolve_font_stack(overrides["font"])
        tokens = tokens.with_overrides(overrides)
    path = out_dir / "_effective.tokens.yaml"
    tokens.dump(path)
    return path, tokens.digest


def series_color(series: dict, index: int) -> str:
    """Return the colour for one chart series.

    Args:
        series: One series mapping from slide_data.json.
        index: Zero-based position of the series in the chart.

    Returns:
        The series' own `color` when it names one, otherwise the token chart
        palette entry for its position.

    Every series used to fall back to a single colour, so a two-series chart
    drew both in one ink and its legend claimed to distinguish two things it
    rendered identically. The palette exists in the token file for this.
    """
    explicit = series.get("color")
    if isinstance(explicit, str) and explicit.strip():
        return explicit
    palette = S["chart_palette"]
    return palette[index % len(palette)]


def content_top() -> float:
    """Return the first y coordinate below the frame rule.

    Four renderers used to hard-code their own origin (75, 80, 66), each tuned
    against the former 20-unit title. They now share one derivation so a change
    to the `slide_title` role moves all of them together.

    Returns:
        The top edge of the slide's content area in SVG units.
    """
    return S["safe"]["top"] + t_size("slide_title") + 40


def chart_legend(series: list, left: float, bottom: float) -> list:
    """Render one legend row for a series list.

    Entries advance by the measured width of the label they carry. The bar
    chart's former fixed pitch of 230 units put every second entry on top of
    the first whenever a label ran long.

    Args:
        series: Series mappings from slide_data.json, in chart order.
        left: Left edge of the plot area; the legend starts here.
        bottom: Bottom edge of the plot area.

    Returns:
        SVG markup fragments for the legend, one per swatch and label.
    """
    parts: list = []
    legend_x = left
    legend_y = bottom + t_size("axis") * t_lh("axis") + 16
    swatch = t_size("axis") * 0.75
    for index, one in enumerate(series):
        color = series_color(one, index)
        label = str(one.get("label", ""))
        parts.append(f'<rect x="{legend_x:.1f}" '
                     f'y="{legend_y - swatch:.1f}" '
                     f'width="{swatch:.1f}" height="{swatch:.1f}" '
                     f'fill="{color}"/>')
        parts.append(f'<text x="{legend_x + swatch + 8:.1f}" '
                     f'y="{legend_y:.1f}" '
                     f'font-size="{t_size("axis"):g}" '
                     f'data-style-role="axis" '
                     f'fill="{S["body"]}">{esc(label)}</text>')
        legend_x += swatch + 8 + measured_width(label, "axis") + 32
    return parts


def chart_note(note: str, right: float, bottom: float) -> str:
    """Render the chart note on its own row below the legend.

    The note used to share the legend's baseline and was right anchored, so a
    wide last legend entry ran straight into it.

    Args:
        note: Note text; assumed non-empty.
        right: Right edge of the plot area, which the note is anchored to.
        bottom: Bottom edge of the plot area.

    Returns:
        SVG markup for the note text element.
    """
    note_y = (bottom + t_size("axis") * t_lh("axis") + 16
              + t_size("footnote") * t_lh("footnote") + 12)
    return (f'<text x="{right}" y="{note_y:.1f}" '
            f'font-size="{t_size("footnote"):g}" '
            f'data-style-role="footnote" '
            f'fill="{S["muted"]}" text-anchor="end">{esc(note)}</text>')


# Chart drawing area
_MAX_TICK_DECIMALS = 3


def tick_decimals(y_max: float) -> int:
    """Return the fewest decimals that keep the six y ticks distinct.

    `{val:.0f}` was fixed, so the shipped example deck -- QWK on a 0-1 axis --
    labelled its six grid lines 0, 0, 0, 1, 1, 1. An axis whose labels repeat
    tells the reader nothing about the scale.

    Args:
        y_max: The axis maximum.

    Returns:
        A decimal count between 0 and 3. A degenerate axis (`y_max` of 0) has
        no distinct ticks at any precision and gets the maximum.
    """
    for decimals in range(_MAX_TICK_DECIMALS + 1):
        labels = {f"{y_max * i / 5:.{decimals}f}" for i in range(6)}
        if len(labels) == 6:
            return decimals
    return _MAX_TICK_DECIMALS


def axis_tick_labels(y_max: float, unit: str) -> List[str]:
    """Return the six y-axis tick labels for a bar or line chart.

    The renderers appended a literal `%` to every tick, so a throughput series
    of 340 requests per second was labelled `340%`. A percentage suffix on
    non-percentage data is not a styling choice; it states something false
    about the numbers. The unit is now declared per slide and defaults to
    empty, because a bare number is never wrong.

    Args:
        y_max: The axis maximum.
        unit: The slide's `unit` string, appended verbatim.

    Returns:
        Six labels, from 0 to `y_max`.
    """
    decimals = tick_decimals(y_max)
    return [f"{y_max * i / 5:.{decimals}f}{unit}" for i in range(6)]


def chart_area(widest_y_tick: str = "100%") -> tuple:
    """Compute the plot rectangle from the active tokens.

    The area clears the frame rule at the top, the widest y-tick label on the
    left, and the category-label plus legend plus note rows at the bottom. The
    former fixed `130, 1100, 100, 520` was tuned for 10-12 unit chart text and
    a 20 unit title, both of which changed.

    Args:
        widest_y_tick: The widest label the caller will draw beside the axis.
            The gutter was measured from a hardcoded `100%`, so a chart whose
            unit is wider than that ran its ticks into the left margin.

    Returns:
        `(left, right, top, bottom)` in SVG units.
    """
    safe = S["safe"]
    rule_y = safe["top"] + t_size("slide_title") + 16
    top = rule_y + 24
    left = safe["left"] + measured_width(widest_y_tick, "axis") + 8
    right = S["w"] - safe["right"]
    axis_adv = t_size("axis") * t_lh("axis")
    foot_adv = t_size("footnote") * t_lh("footnote")
    # category labels, then legend row, then note row
    bottom = S["h"] - safe["bottom"] - (axis_adv + axis_adv + foot_adv + 24)
    return left, right, top, bottom


# ── Text helpers ──────────────────────────────────────────────────────────────

def esc(v) -> str:
    return str(v).replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;").replace('"', "&quot;")


def wrap(text: str, max_chars: int = 60) -> list:
    words = str(text).split()
    lines, cur, n = [], [], 0
    for w in words:
        if n + len(w) + 1 > max_chars and cur:
            lines.append(" ".join(cur))
            cur, n = [w], len(w)
        else:
            cur.append(w)
            n += len(w) + 1
    if cur:
        lines.append(" ".join(cur))
    return lines or [""]


def tlines(lines: list, x, y, size, color, anchor="start", weight="normal",
           lh=1.45, *, role: str) -> str:
    """Render a multi-line text element.

    Args:
        lines: Already-wrapped lines, one per rendered line.
        x: Left, centre, or right coordinate, per `anchor`.
        y: Baseline of the first line.
        size: Font size in SVG units.
        color: Fill colour.
        anchor: SVG `text-anchor` value.
        weight: SVG `font-weight` value.
        lh: Line-height multiplier, applied to `size` for each `dy` after the
            first line.
        role: The typography role this text realises. Keyword-only and
            mandatory: the visual linter skips a `<text>` with no
            `data-style-role`, so an optional marker with a default would let a
            caller silently disable the type-floor and colour rules for its
            text.

    Returns:
        SVG markup for one `<text>` element.
    """
    spans = []
    for i, line in enumerate(lines):
        dy = "0" if i == 0 else f"{size * lh:.1f}"
        spans.append(f'<tspan x="{x}" dy="{dy}">{esc(line)}</tspan>')
    return (f'<text x="{x}" y="{y}" font-size="{size}" font-weight="{weight}" '
            f'fill="{color}" text-anchor="{anchor}" '
            f'data-style-role="{role}">{"".join(spans)}</text>')


def measured_width(text: str, role: str) -> float:
    """Measure a string's advance width at a type role's size.

    Args:
        text: The string to measure.
        role: Type role key, such as `body`.

    Returns:
        Advance width in SVG units.
    """
    return text_width(text, S["font_resolved"], t_size(role), t_weight(role))


def wrap_to_width(text: str, max_width: float, role: str) -> list:
    """Wrap text to a pixel budget using real font metrics.

    Character-count wrapping cannot survive a size change: 88 characters at 14
    units and at 21 units occupy different widths. A word wider than the budget
    is kept on its own line rather than dropped.

    Args:
        text: The string to wrap.
        max_width: Available width in SVG units.
        role: Type role key used for measurement.

    Returns:
        Wrapped lines; always at least one entry.
    """
    words = str(text).split()
    if not words:
        return [""]
    lines: list = []
    current: list = []
    for word in words:
        candidate = " ".join(current + [word])
        if current and measured_width(candidate, role) > max_width:
            lines.append(" ".join(current))
            current = [word]
        else:
            current.append(word)
    if current:
        lines.append(" ".join(current))
    return lines


# ── Common slide frame ────────────────────────────────────────────────────────

_FRAME_VARIANTS = ("left", "centered")


def frame(title: str, footer: str = "", *, variant: str = "left") -> str:
    """Render the shared slide frame: title, rule, and footer.

    The former 6px top accent bar and centred 20pt title are gone. The title now
    uses the `slide_title` role inside the token safe area, and the rule sits
    under the title's baseline rather than at a fixed y=54.

    Args:
        title: Slide title text.
        footer: Optional footer text, rendered at the `footnote` role.
        variant: `left` for the standard left-aligned title, `centered` for
            section dividers.

    Returns:
        SVG markup for the frame elements.

    Raises:
        ValueError: If `variant` is not a known frame variant.
    """
    if variant not in _FRAME_VARIANTS:
        raise ValueError(
            f"unknown frame variant {variant!r}; expected one of {_FRAME_VARIANTS}"
        )
    safe = S["safe"]
    size = t_size("slide_title")
    # `size` is used here, not the measured ascent, because three later tasks
    # place their content against `rule_y = safe.top + slide_title + 16`. It is
    # a conservative proxy: DejaVu Sans reports ascent 30 at size 32, so the
    # title's em box starts 2 units *inside* the safe area. Task 6's test pins
    # that, so a font whose ascent exceeds its em size fails loudly rather than
    # clipping silently.
    baseline = safe["top"] + size
    rule_y = baseline + 16
    x = S["w"] / 2 if variant == "centered" else safe["left"]
    anchor = "middle" if variant == "centered" else "start"

    parts = [
        f'<rect width="{S["w"]}" height="{S["h"]}" fill="{S["bg"]}" '
        f'data-bleed="true"/>',
        f'<text x="{x:g}" y="{baseline:g}" font-size="{size:g}" '
        f'font-weight="{t_weight("slide_title")}" fill="{S["accent"]}" '
        f'data-style-role="slide_title" '
        f'text-anchor="{anchor}">{esc(title)}</text>',
        f'<line x1="{safe["left"]}" y1="{rule_y:g}" '
        f'x2="{S["w"] - safe["right"]}" y2="{rule_y:g}" '
        f'stroke="{S["divider"]}" stroke-width="1.5" '
        f'data-bleed="true" data-style-role="divider"/>',
    ]
    if footer:
        fs = t_size("footnote")
        # The baseline is lifted by the measured descent so the footer's
        # descenders end *on* the safe-area boundary rather than three units
        # past it. Placing the baseline on the boundary is the obvious-looking
        # choice and is wrong: plan 2's `safe-area` rule reports it on every
        # slide that carries a footer.
        _, descent = vertical_metrics(
            S["font_resolved"], fs, t_weight("footnote"))
        baseline_y = S["h"] - safe["bottom"] - descent
        parts.append(
            f'<text x="{S["w"] - safe["right"]}" y="{baseline_y:g}" '
            f'font-size="{fs:g}" fill="{S["muted"]}" '
            f'data-style-role="footnote" '
            f'text-anchor="end">{esc(footer)}</text>'
        )
    return "\n  ".join(parts)


def svg(body: str) -> str:
    """Wrap slide body markup in the root SVG element.

    Args:
        body: Slide content markup.

    Returns:
        A complete SVG document string.
    """
    return (f'<svg xmlns="http://www.w3.org/2000/svg" '
            f'viewBox="0 0 {S["w"]} {S["h"]}" '
            f'font-family="{S["font_resolved"]}">\n'
            f'  {body}\n</svg>\n')


# ── Renderers ─────────────────────────────────────────────────────────────────

def render_title(sl: dict, meta: dict) -> str:
    title    = sl.get("title",    meta.get("experiment", "Research Report"))
    subtitle = sl.get("subtitle", "")
    author   = sl.get("author",   "")
    date     = sl.get("date",     meta.get("date", ""))
    footer   = meta.get("footer", "")
    cx       = 600

    # The two 8px accent bars that used to run the full width at y=0 and y=667
    # are gone with the rest of the templated signature removed in Task 6.
    parts = [
        f'<rect width="{S["w"]}" height="{S["h"]}" fill="{S["bg"]}" '
        f'data-bleed="true"/>',
    ]

    title_role = "deck_title"
    title_lines = wrap_to_width(title, S["w"] - 2 * S["safe"]["left"] - 160, title_role)
    title_adv = t_size(title_role) * t_lh(title_role)
    title_y = 255 - (len(title_lines) - 1) * (title_adv / 2)
    parts.append(tlines(title_lines, cx, title_y, t_size(title_role),
                        S["accent"], "middle", str(t_weight(title_role)),
                        t_lh(title_role), role=title_role))

    div_y = title_y + len(title_lines) * title_adv
    parts.append(f'<line x1="200" y1="{div_y:g}" x2="1000" y2="{div_y:g}" '
                 f'stroke="{S["divider"]}" stroke-width="1.5"/>')

    base_y = div_y + 40
    if subtitle:
        sub_role = "takeaway"
        sub_lines = wrap_to_width(subtitle, S["w"] - 2 * S["safe"]["left"] - 120,
                                  sub_role)
        parts.append(tlines(sub_lines, cx, base_y, t_size(sub_role),
                            S["muted"], "middle", str(t_weight(sub_role)),
                            t_lh(sub_role), role=sub_role))
        base_y += len(sub_lines) * t_size(sub_role) * t_lh(sub_role) + 12

    meta_str = "  ·  ".join(filter(None, [author, date]))
    if meta_str:
        parts.append(f'<text x="{cx}" y="{base_y + t_size("caption"):g}" '
                     f'font-size="{t_size("caption"):g}" '
                     f'data-style-role="caption" '
                     f'fill="{S["muted"]}" text-anchor="middle">{esc(meta_str)}</text>')
    if footer:
        # Same lift as `frame()`: the baseline sits a measured descent
        # above the safe-area boundary, not on it.
        fs = t_size("footnote")
        _, descent = vertical_metrics(
            S["font_resolved"], fs, t_weight("footnote"))
        parts.append(f'<text x="{S["w"] - S["safe"]["right"]}" '
                     f'y="{S["h"] - S["safe"]["bottom"] - descent:g}" '
                     f'font-size="{fs:g}" fill="{S["muted"]}" '
                     f'data-style-role="footnote" '
                     f'text-anchor="end">{esc(footer)}</text>')

    return svg("\n  ".join(parts))


def render_bullet_list(sl: dict, meta: dict) -> str:
    title    = sl.get("title", "")
    bullets  = sl.get("bullets", [])
    numbered = sl.get("numbered", False)
    footer   = meta.get("footer", "")

    parts = [frame(title, footer)]
    safe = S["safe"]
    x_dot = safe["left"] + 14
    x_text = safe["left"] + 52
    text_budget = S["w"] - x_text - safe["right"]
    body_adv = t_size("body") * t_lh("body")
    y = t_size("slide_title") + safe["top"] + 56
    for i, item in enumerate(bullets):
        lines = wrap_to_width(str(item), text_budget, "body")
        if numbered:
            r = t_size("footnote") * 0.75
            parts.append(f'<circle cx="{x_dot}" cy="{y - t_size("body") * 0.32:g}" '
                         f'r="{r:g}" fill="{S["accent"]}"/>')
            parts.append(f'<text x="{x_dot}" '
                         f'y="{y - t_size("body") * 0.32 + r * 0.55:g}" '
                         f'font-size="{t_size("footnote"):g}" font-weight="700" '
                         f'data-style-role="footnote" '
                         f'fill="{S["white"]}" text-anchor="middle">{i + 1}</text>')
        else:
            parts.append(f'<circle cx="{x_dot}" cy="{y - t_size("body") * 0.30:g}" '
                         f'r="{t_size("body") * 0.28:g}" fill="{S["accent"]}"/>')
        parts.append(tlines(lines, x_text, y, t_size("body"), S["body"],
                            "start", str(t_weight("body")), t_lh("body"),
                            role="body"))
        y += len(lines) * body_adv + S["spacing"][2]

    return svg("\n  ".join(parts))


def render_bar_chart(sl: dict, meta: dict) -> str:
    title      = sl.get("title", "")
    categories = sl.get("categories", [])
    series     = sl.get("series", [])
    y_max      = float(sl.get("y_max", 100))
    unit       = str(sl.get("unit", ""))
    note       = sl.get("note", "")
    footer     = meta.get("footer", "")
    slide_index = sl.get("index", 0)

    parts = [frame(title, footer)]
    chart_parts = []
    tick_labels = axis_tick_labels(y_max, unit)
    # One place finer than the axis: a bar labelled at the axis's own
    # precision rounded a QWK of 0.853 to 0.9.
    value_decimals = tick_decimals(y_max) + 1
    CL, CR, CT, CB = chart_area(
        max(tick_labels, key=lambda label: measured_width(label, "axis")))
    CW, CH = CR - CL, CB - CT

    for i in range(6):
        val = y_max * i / 5
        y   = CB - (val / y_max) * CH
        chart_parts.append(f'<line x1="{CL}" y1="{y:.1f}" x2="{CR}" y2="{y:.1f}" '
                     f'stroke="{S["divider"]}" stroke-width="1"/>')
        chart_parts.append(f'<text x="{CL - 8}" y="{y + 4:.1f}" '
                     f'font-size="{t_size("axis"):g}" '
                     f'data-style-role="axis" '
                     f'fill="{S["muted"]}" text-anchor="end">{tick_labels[i]}</text>')

    chart_parts.append(f'<line x1="{CL}" y1="{CT}" x2="{CL}" y2="{CB}" '
                 f'stroke="{S["muted"]}" stroke-width="1.5"/>')
    chart_parts.append(f'<line x1="{CL}" y1="{CB}" x2="{CR}" y2="{CB}" '
                 f'stroke="{S["muted"]}" stroke-width="1.5"/>')

    n_cats = len(categories)
    n_ser  = len(series)
    if not n_cats or not n_ser:
        parts.append(_wrap_pptx_role("chart", slide_index, chart_parts,
                                     bbox=(CL, CT, CR - CL, CB - CT),
                                     style_keys=("font",)))
        return svg("\n  ".join(parts))

    cat_slot = CW / n_cats
    group_w  = cat_slot * 0.70
    bar_w    = group_w / n_ser
    pad      = cat_slot * 0.15

    for ci, cat in enumerate(categories):
        gx = CL + ci * cat_slot + pad
        lx = gx + group_w / 2
        chart_parts.append(f'<text x="{lx:.1f}" y="{CB + 20}" '
                     f'font-size="{t_size("axis"):g}" '
                     f'font-weight="{t_weight("axis")}" '
                     f'data-style-role="axis" '
                     f'fill="{S["body"]}" text-anchor="middle">{esc(cat)}</text>')

        for si, ser in enumerate(series):
            vals  = ser.get("values", [])
            if ci >= len(vals):
                continue
            val   = float(vals[ci])
            color = series_color(ser, si)
            bx    = gx + si * bar_w
            bh    = max((val / y_max) * CH, 2)
            by    = CB - bh

            chart_parts.append(f'<rect x="{bx:.1f}" y="{by:.1f}" '
                         f'width="{bar_w - 3:.1f}" height="{bh:.1f}" fill="{color}"/>')
            if bh > 16:
                chart_parts.append(f'<text x="{bx + (bar_w - 3) / 2:.1f}" y="{by - 4:.1f}" '
                             f'font-size="{t_size("footnote"):g}" '
                             f'font-weight="700" '
                             f'data-style-role="footnote" fill="{color}" '
                             f'text-anchor="middle">'
                             f'{val:.{value_decimals}f}{unit}</text>')

    chart_parts.extend(chart_legend(series, CL, CB))

    if note:
        chart_parts.append(chart_note(note, CR, CB))

    parts.append(_wrap_pptx_role("chart", slide_index, chart_parts,
                                 bbox=(CL, CT, CR - CL, CB - CT),
                                 style_keys=("font",)))
    return svg("\n  ".join(parts))


def render_line_chart(sl: dict, meta: dict) -> str:
    """SVG preview for a line chart. Same slide_data.json schema as
    bar_chart (categories/series/y_max/note); svg_to_pptx/converter.py
    replaces this hand-drawn preview with a real native line chart."""
    title      = sl.get("title", "")
    categories = sl.get("categories", [])
    series     = sl.get("series", [])
    y_max      = float(sl.get("y_max", 100))
    unit       = str(sl.get("unit", ""))
    note       = sl.get("note", "")
    footer     = meta.get("footer", "")
    slide_index = sl.get("index", 0)

    parts = [frame(title, footer)]
    chart_parts = []
    tick_labels = axis_tick_labels(y_max, unit)
    # One place finer than the axis: a bar labelled at the axis's own
    # precision rounded a QWK of 0.853 to 0.9.
    value_decimals = tick_decimals(y_max) + 1
    CL, CR, CT, CB = chart_area(
        max(tick_labels, key=lambda label: measured_width(label, "axis")))
    CW, CH = CR - CL, CB - CT

    for i in range(6):
        val = y_max * i / 5
        y   = CB - (val / y_max) * CH
        chart_parts.append(f'<line x1="{CL}" y1="{y:.1f}" x2="{CR}" y2="{y:.1f}" '
                     f'stroke="{S["divider"]}" stroke-width="1"/>')
        chart_parts.append(f'<text x="{CL - 8}" y="{y + 4:.1f}" '
                     f'font-size="{t_size("axis"):g}" '
                     f'data-style-role="axis" '
                     f'fill="{S["muted"]}" text-anchor="end">{tick_labels[i]}</text>')
    chart_parts.append(f'<line x1="{CL}" y1="{CT}" x2="{CL}" y2="{CB}" '
                 f'stroke="{S["muted"]}" stroke-width="1.5"/>')
    chart_parts.append(f'<line x1="{CL}" y1="{CB}" x2="{CR}" y2="{CB}" '
                 f'stroke="{S["muted"]}" stroke-width="1.5"/>')

    n_cats = len(categories)
    if n_cats:
        step = CW / max(n_cats - 1, 1)
        for ci, cat in enumerate(categories):
            x = CL + ci * step
            chart_parts.append(f'<text x="{x:.1f}" y="{CB + 20}" '
                         f'font-size="{t_size("axis"):g}" '
                         f'data-style-role="axis" '
                         f'fill="{S["body"]}" text-anchor="middle">{esc(cat)}</text>')
        for si, ser in enumerate(series):
            vals = ser.get("values", [])
            color = series_color(ser, si)
            points = []
            for ci, val in enumerate(vals[:n_cats]):
                x = CL + ci * step
                y = CB - (float(val) / y_max) * CH
                points.append(f"{x:.1f},{y:.1f}")
                chart_parts.append(f'<circle cx="{x:.1f}" cy="{y:.1f}" r="4" fill="{color}"/>')
            if len(points) > 1:
                chart_parts.append(f'<polyline points="{" ".join(points)}" '
                             f'fill="none" stroke="{color}" stroke-width="2.5"/>')

    # The legend row `chart_area` budgets for was never drawn here, so two
    # distinguishable lines were still unlabelled.
    chart_parts.extend(chart_legend(series, CL, CB))

    if note:
        chart_parts.append(chart_note(note, CR, CB))

    parts.append(_wrap_pptx_role("chart", slide_index, chart_parts,
                                 bbox=(CL, CT, CR - CL, CB - CT),
                                 style_keys=("font",)))
    return svg("\n  ".join(parts))


def render_pie_chart(sl: dict, meta: dict) -> str:
    """SVG preview for a pie chart. Schema: categories/values/colors(optional)/
    note. svg_to_pptx/converter.py replaces this preview with a real native
    pie chart."""
    import math
    title      = sl.get("title", "")
    categories = sl.get("categories", [])
    values     = sl.get("values", [])
    colors     = sl.get("colors") or [S["blue"], S["good"], S["warn"], S["danger"], S["accent"]]
    note       = sl.get("note", "")
    footer     = meta.get("footer", "")
    slide_index = sl.get("index", 0)

    parts = [frame(title, footer)]
    if not categories or not values:
        return svg("\n  ".join(parts))

    CL, CR, CT, CB = chart_area()
    cx, cy, r = 420, 340, 200
    total = sum(values) or 1
    chart_parts = []
    angle = -math.pi / 2
    for i, val in enumerate(values):
        frac = val / total
        end_angle = angle + frac * 2 * math.pi
        x1, y1 = cx + r * math.cos(angle), cy + r * math.sin(angle)
        x2, y2 = cx + r * math.cos(end_angle), cy + r * math.sin(end_angle)
        large_arc = 1 if (end_angle - angle) > math.pi else 0
        color = colors[i % len(colors)]
        chart_parts.append(
            f'<path d="M{cx},{cy} L{x1:.1f},{y1:.1f} '
            f'A{r},{r} 0 {large_arc} 1 {x2:.1f},{y2:.1f} Z" fill="{color}"/>'
        )
        angle = end_angle

    swatch = t_size("axis")
    legend_adv = t_size("axis") * t_lh("axis") + 12
    for i, cat in enumerate(categories):
        ly = 160 + i * legend_adv
        color = colors[i % len(colors)]
        chart_parts.append(f'<rect x="700" y="{ly:.1f}" '
                     f'width="{swatch:g}" height="{swatch:g}" fill="{color}"/>')
        chart_parts.append(f'<text x="{700 + swatch + 8:g}" y="{ly + swatch * 0.8:.1f}" '
                     f'font-size="{t_size("axis"):g}" '
                     f'data-style-role="axis" '
                     f'fill="{S["body"]}">{esc(cat)}</text>')

    if note:
        # Same row as the other two renderers put it in, so a deck mixing
        # chart types keeps its notes on one baseline.
        chart_parts.append(chart_note(note, CR, CB))

    parts.append(_wrap_pptx_role("chart", slide_index, chart_parts,
                                 bbox=(CL, CT, CR - CL, CB - CT),
                                 style_keys=("font",)))
    return svg("\n  ".join(parts))


def _wrap_pptx_role(role: str, slide_index: int, inner_parts: list,
                    bbox: tuple, style_keys: tuple, node_id: str = "",
                    extra_style: Optional[Dict[str, Any]] = None) -> str:
    """Wrap hand-drawn preview markup in the data-pptx-role marker so
    svg_to_pptx/converter.py materializes a real native PPTX object instead
    of flattening these shapes. bbox is (x, y, w, h) in SVG user units.

    `extra_style` carries resolved values that are not colour roles -- the
    table's type sizes -- so the native object is typeset at the same scale as
    the preview beside it rather than at a size of its own.
    """
    resolved = {k: S[k] for k in style_keys if k in S}
    resolved.update(extra_style or {})
    style_json = esc(json.dumps(resolved))
    bx, by, bw, bh = bbox
    attrs = [
        f'data-pptx-role="{role}"',
        f'data-pptx-source="slide_data.json#{slide_index}"',
        f'data-pptx-bbox="{bx:.1f},{by:.1f},{bw:.1f},{bh:.1f}"',
        f'data-pptx-style="{style_json}"',
    ]
    if node_id:
        attrs.append(f'data-node-id="{esc(node_id)}"')
    return (f'<g {" ".join(attrs)}>\n  ' + "\n  ".join(inner_parts) + '\n  </g>')


def render_table(sl: dict, meta: dict) -> str:
    title         = sl.get("title", "")
    columns       = sl.get("columns", [])
    rows          = sl.get("rows", [])
    highlight_col = sl.get("highlight_col")   # 0-indexed; colorizes +/- values
    footer        = meta.get("footer", "")
    slide_index   = sl.get("index", 0)

    parts = [frame(title, footer)]

    n_cols = len(columns)
    if not n_cols:
        return svg("\n  ".join(parts))

    safe = S["safe"]
    tl = safe["left"]
    tr = S["w"] - safe["right"]
    tw = tr - tl
    col_w = tw / n_cols
    top_y = content_top()

    # A row is its text's line box plus symmetric vertical padding. The former
    # `min(50, 450 / (len(rows) + 1))` shrank the row to fit the slide, which
    # at 13-unit text was invisible and at the body role would clip descenders.
    header_h = t_size("node_label") * t_lh("node_label") + 2 * 12
    body_h = t_size("body") * t_lh("body") + 2 * 12
    available = S["h"] - top_y - safe["bottom"]
    capacity = int((available - header_h) // body_h)
    if len(rows) > capacity:
        raise SlideCapacityError(
            f"table has {len(rows)} rows but only {capacity} fit at the body "
            f"role ({t_size('body'):g} units); split the table across slides or "
            f"drop columns rather than shrinking the text"
        )
    row_h = body_h
    table_h = header_h + row_h * len(rows)

    table_parts = [
        f'<rect x="{tl}" y="{top_y}" width="{tw}" '
        f'height="{header_h:.1f}" fill="{S["accent"]}"/>'
    ]
    for ci, col in enumerate(columns):
        cx = tl + ci * col_w + col_w / 2
        table_parts.append(f'<text x="{cx:.1f}" y="{top_y + header_h * 0.66:.1f}" '
                     f'font-size="{t_size("node_label"):g}" '
                     f'font-weight="{t_weight("node_label")}" fill="{S["white"]}" '
                     f'data-style-role="node_label" '
                     f'text-anchor="middle">{esc(col)}</text>')

    for ri, row in enumerate(rows):
        ry = top_y + header_h + ri * row_h
        bg = S["card"] if ri % 2 == 0 else S["bg"]
        table_parts.append(f'<rect x="{tl}" y="{ry:.1f}" width="{tw}" '
                     f'height="{row_h:.1f}" fill="{bg}" '
                     f'stroke="{S["border"]}" stroke-width="0.5"/>')
        for ci, cell in enumerate(row):
            cx    = tl + ci * col_w + col_w / 2
            cy    = ry + row_h * 0.66
            color = S["body"]
            if highlight_col is not None and ci == highlight_col:
                cs = str(cell)
                if "+" in cs:
                    color = S["good"]
                elif "-" in cs:
                    color = S["danger"]
            table_parts.append(f'<text x="{cx:.1f}" y="{cy:.1f}" '
                         f'font-size="{t_size("body"):g}" '
                         f'data-style-role="body" '
                         f'fill="{color}" text-anchor="middle">{esc(cell)}</text>')

    table_parts.append(f'<rect x="{tl}" y="{top_y}" width="{tw}" '
                 f'height="{table_h:.1f}" fill="none" '
                 f'stroke="{S["border"]}" stroke-width="1.5"/>')

    parts.append(_wrap_pptx_role(
        "table", slide_index, table_parts,
        bbox=(tl, top_y, tw, table_h),
        style_keys=("accent", "white", "card", "bg", "body", "good", "danger", "font"),
        extra_style={"header_size": t_size("node_label"),
                     "body_size": t_size("body")},
    ))

    return svg("\n  ".join(parts))


def _metric_card_height() -> float:
    """Return the height one metric card needs for its three type roles.

    Returns:
        Card height in SVG units: the caption, display value, and footnote
        stacked with grid spacing between them, inside the card surface's
        vertical padding.
    """
    surface = _TOKENS.surface("card")
    return (surface["padding"]["y"]
            + t_size("caption")
            + S["spacing"][4]
            + t_size("deck_title")
            + S["spacing"][3]
            + t_size("footnote")
            + surface["padding"]["y"])


def render_metric_cards(sl: dict, meta: dict) -> str:
    title   = sl.get("title", "")
    metrics = sl.get("metrics", [])
    footer  = meta.get("footer", "")

    parts = [frame(title, footer)]
    n = len(metrics)
    if not n:
        return svg("\n  ".join(parts))

    cols = 2 if n == 4 else min(n, 3)
    rows = (n + cols - 1) // cols
    safe = S["safe"]
    band_top = content_top()
    band_h = S["h"] - band_top - safe["bottom"]
    gap = S["spacing"][4]
    cw = (S["w"] - safe["left"] - safe["right"] - (cols - 1) * gap) / cols
    # The card is as tall as its content, not as tall as the slide. Dividing the
    # whole band by the row count left the label at the card's top edge, the
    # value adrift in its middle, and the change stranded at the bottom.
    ch = _metric_card_height()
    grid_h = rows * ch + (rows - 1) * gap
    if grid_h > band_h:
        raise SlideCapacityError(
            f"{n} metric cards need {grid_h:.0f} units of height in {rows} "
            f"rows but only {band_h:.0f} are available; use fewer cards per "
            f"slide"
        )
    top = band_top + (band_h - grid_h) / 2
    surface = _TOKENS.surface("card")

    for i, m in enumerate(metrics):
        col = i % cols
        row = i // cols
        cx = safe["left"] + col * (cw + gap)
        cy = top + row * (ch + gap)
        color = m.get("color", S["blue"])
        label = m.get("label", "")
        value = m.get("value", "")
        change = m.get("change", "")

        parts.append(f'<rect x="{cx:.1f}" y="{cy:.1f}" width="{cw:.1f}" '
                     f'height="{ch:.1f}" rx="{surface["radius"]}" '
                     f'fill="{S["card"]}" stroke="{S["divider"]}" '
                     f'stroke-width="{surface["border_width"]}"/>')
        # Inset by the card radius so the bar's own corners cannot render
        # outside the card outline, which a full-width bar with a smaller
        # radius did on every card.
        parts.append(f'<rect x="{cx + surface["radius"]:.1f}" y="{cy:.1f}" '
                     f'width="{cw - 2 * surface["radius"]:.1f}" '
                     f'height="5" rx="2.5" fill="{color}"/>')
        label_y = cy + surface["padding"]["y"] + t_size("caption")
        parts.append(f'<text x="{cx + cw / 2:.1f}" y="{label_y:.1f}" '
                     f'font-size="{t_size("caption"):g}" '
                     f'data-style-role="caption" '
                     f'fill="{S["muted"]}" text-anchor="middle">{esc(label)}</text>')
        value_y = label_y + S["spacing"][4] + t_size("deck_title")
        parts.append(f'<text x="{cx + cw / 2:.1f}" '
                     f'y="{value_y:.1f}" '
                     f'font-size="{t_size("deck_title"):g}" '
                     f'font-weight="{t_weight("deck_title")}" fill="{color}" '
                     f'data-style-role="deck_title" '
                     f'text-anchor="middle">{esc(value)}</text>')
        if change:
            cc = (S["good"] if "+" in str(change)
                  else (S["danger"] if "-" in str(change) else S["muted"]))
            parts.append(f'<text x="{cx + cw / 2:.1f}" '
                         f'y="{value_y + S["spacing"][3] + t_size("footnote"):.1f}" '
                         f'font-size="{t_size("footnote"):g}" fill="{cc}" '
                         f'data-style-role="footnote" '
                         f'text-anchor="middle">{esc(change)}</text>')

    return svg("\n  ".join(parts))


_PANEL_COUNT = 2


def _panel_geometry() -> tuple:
    """Return the two-column panel rectangle geometry.

    Returns:
        `(left_x, right_x, top_y, width, height)` in SVG units. The former
        literals (60, 640, 66, 500, 562) predate the token safe area and left
        an asymmetric margin against a 48-unit safe edge.
    """
    safe = S["safe"]
    gap = S["spacing"][4]
    top = content_top()
    width = (S["w"] - safe["left"] - safe["right"] - gap) / _PANEL_COUNT
    height = S["h"] - top - safe["bottom"]
    return safe["left"], safe["left"] + width + gap, top, width, height


def _panel_text_width() -> float:
    """Return the width available to prose inside one panel.

    Returns:
        The panel width less the card surface's horizontal padding on both
        sides.
    """
    width = _panel_geometry()[3]
    return width - 2 * _TOKENS.surface("card")["padding"]["x"]


def _panel_lines(text: str, width: float) -> list:
    """Wrap panel prose at the body role.

    Args:
        text: Prose to wrap.
        width: Available width in SVG units.

    Returns:
        Wrapped lines; always at least one entry.
    """
    return wrap_to_width(text, width, "body")


def render_two_column(sl: dict, meta: dict) -> str:
    title  = sl.get("title", "")
    left   = sl.get("left",  {})
    right  = sl.get("right", {})
    footer = meta.get("footer", "")

    parts = [frame(title, footer)]

    surface = _TOKENS.surface("card")
    pad_x = surface["padding"]["x"]
    pad_y = surface["padding"]["y"]
    body_adv = t_size("body") * t_lh("body")

    def panel(p: dict, px: float, py: float, pw: float, ph: float) -> list:
        """Render one column panel.

        Args:
            p: Panel mapping with optional `title` and `content`. `content` is
                either a prose string or a list of bullet items.
            px: Panel left edge.
            py: Panel top edge.
            pw: Panel width.
            ph: Panel height.

        Returns:
            SVG markup fragments for the panel.
        """
        out = [f'<rect x="{px:.1f}" y="{py:.1f}" width="{pw:.1f}" '
               f'height="{ph:.1f}" rx="{surface["radius"]}" fill="{S["card"]}" '
               f'stroke="{S["divider"]}" '
               f'stroke-width="{surface["border_width"]}"/>']
        ty = py + pad_y + t_size("takeaway")
        pt = p.get("title", "")
        if pt:
            out.append(f'<text x="{px + pad_x:.1f}" y="{ty:.1f}" '
                       f'font-size="{t_size("takeaway"):g}" '
                       f'font-weight="{t_weight("takeaway")}" '
                       f'data-style-role="takeaway" '
                       f'fill="{S["accent"]}">{esc(pt)}</text>')
            rule_y = ty + S["spacing"][2]
            out.append(f'<line x1="{px + pad_x:.1f}" y1="{rule_y:.1f}" '
                       f'x2="{px + pw - pad_x:.1f}" y2="{rule_y:.1f}" '
                       f'stroke="{S["divider"]}" stroke-width="1" '
                       f'data-style-role="divider"/>')
            ty = rule_y + S["spacing"][3] + t_size("body")

        content = p.get("content", [])
        if isinstance(content, str):
            # Character-count wrapping (`int(pw / 9)`) was calibrated for
            # 13-unit text; at the body role it overruns the panel by a third.
            out.append(tlines(_panel_lines(content, pw - 2 * pad_x),
                              px + pad_x, ty, t_size("body"), S["body"],
                              "start", str(t_weight("body")), t_lh("body"),
                              role="body"))
        else:
            x_dot = px + pad_x + t_size("body") * 0.28
            x_text = px + pad_x + t_size("body")
            for item in content:
                lines = _panel_lines(str(item), px + pw - pad_x - x_text)
                out.append(f'<circle cx="{x_dot:.1f}" '
                           f'cy="{ty - t_size("body") * 0.30:.1f}" '
                           f'r="{t_size("body") * 0.28:g}" fill="{S["accent"]}"/>')
                out.append(tlines(lines, x_text, ty, t_size("body"), S["body"],
                                  "start", str(t_weight("body")), t_lh("body"),
                                  role="body"))
                ty += len(lines) * body_adv + S["spacing"][2]
        return out

    left_x, right_x, panel_top, panel_w, panel_h = _panel_geometry()
    parts += panel(left, left_x, panel_top, panel_w, panel_h)
    parts += panel(right, right_x, panel_top, panel_w, panel_h)

    return svg("\n  ".join(parts))


def render_timeline(sl: dict, meta: dict) -> str:
    title  = sl.get("title", "")
    events = sl.get("events", [])
    footer = meta.get("footer", "")

    parts = [frame(title, footer)]
    if not events:
        return svg("\n  ".join(parts))

    n = len(events)
    safe = S["safe"]
    top = content_top()
    bottom = S["h"] - safe["bottom"]
    y0 = (top + bottom) / 2
    node_r = 10
    stem = S["spacing"][4]

    label_adv = t_size("node_label") * t_lh("node_label")
    date_adv = t_size("footnote") * t_lh("footnote")
    detail_adv = t_size("caption") * t_lh("caption")

    # Each event owns an equal share of the axis, and its text is wrapped to
    # that share. The former `wrap(label, 18)` / `wrap(detail, 20)` character
    # counts had no relationship to the number of events on the slide.
    axis_left = safe["left"]
    axis_right = S["w"] - safe["right"]
    slot = (axis_right - axis_left) / n
    text_budget = slot - S["spacing"][3]

    blocks = []
    for ev in events:
        label_lines = wrap_to_width(
            str(ev.get("label", "")), text_budget, "node_label")
        date_str = str(ev.get("date", ""))
        detail = str(ev.get("detail", ""))
        detail_lines = (wrap_to_width(detail, text_budget, "caption")
                        if detail else [])
        widths = [measured_width(line, "node_label") for line in label_lines]
        widths += [measured_width(line, "caption") for line in detail_lines]
        if date_str:
            widths.append(measured_width(date_str, "footnote"))
        blocks.append({
            "label_lines": label_lines,
            "date": date_str,
            "detail_lines": detail_lines,
            "height": (len(label_lines) * label_adv
                       + (date_adv if date_str else 0.0)
                       + len(detail_lines) * detail_adv),
            "width": max(widths) if widths else 0.0,
        })

    half_band = (bottom - top) / 2 - node_r - stem
    tallest = max(block["height"] for block in blocks)
    if tallest > half_band:
        raise SlideCapacityError(
            f"{n} timeline events need {tallest:.0f} units of label height but "
            f"only {half_band:.0f} are available on either side of the axis; "
            f"split the timeline across slides or shorten the details"
        )

    max_slot = max(block["width"] for block in blocks) + S["spacing"][3]
    x_left = axis_left + max_slot / 2
    x_right = axis_right - max_slot / 2
    if n > 1:
        pitch = (x_right - x_left) / (n - 1)
        # Half a unit of slack: the pitch equals the slot exactly when every
        # label wraps to its full budget, and floating-point noise there would
        # reject a layout that fits. Half a unit is far below the 8-unit grid,
        # so it cannot hide a real overlap.
        if pitch + 0.5 < max_slot:
            raise SlideCapacityError(
                f"{n} timeline events need {max_slot:.0f} units of horizontal "
                f"pitch but only {pitch:.0f} are available; split the timeline "
                f"across slides or shorten the labels"
            )
    xs = [x_left + i * (x_right - x_left) / max(n - 1, 1) for i in range(n)]

    parts.append(f'<line x1="{axis_left}" y1="{y0:.1f}" x2="{axis_right}" '
                 f'y2="{y0:.1f}" stroke="{S["line"]}" stroke-width="3"/>')

    for i, (ev, x, block) in enumerate(zip(events, xs, blocks)):
        color = ev.get("color", S["accent"])
        above = i % 2 == 0
        event_parts = [f'<circle cx="{x:.1f}" cy="{y0:.1f}" r="{node_r}" '
                       f'fill="{color}" stroke="{S["bg"]}" stroke-width="2"/>']

        if above:
            event_parts.append(
                f'<line x1="{x:.1f}" y1="{y0 - node_r:.1f}" x2="{x:.1f}" '
                f'y2="{y0 - node_r - stem:.1f}" stroke="{color}" '
                f'stroke-width="1.5" stroke-dasharray="3,2"/>')
            # The block is laid out upward from the stem, so a two-line detail
            # grows away from the axis instead of across it.
            ty = y0 - node_r - stem - block["height"] + t_size("node_label")
        else:
            event_parts.append(
                f'<line x1="{x:.1f}" y1="{y0 + node_r:.1f}" x2="{x:.1f}" '
                f'y2="{y0 + node_r + stem:.1f}" stroke="{color}" '
                f'stroke-width="1.5" stroke-dasharray="3,2"/>')
            ty = y0 + node_r + stem + t_size("node_label")

        event_parts.append(tlines(
            block["label_lines"], x, ty, t_size("node_label"), S["accent"],
            "middle", str(t_weight("node_label")), t_lh("node_label"),
            role="node_label"))
        ty += len(block["label_lines"]) * label_adv

        if block["date"]:
            event_parts.append(
                f'<text x="{x:.1f}" y="{ty:.1f}" '
                f'font-size="{t_size("footnote"):g}" '
                f'data-style-role="footnote" fill="{S["muted"]}" '
                f'text-anchor="middle">{esc(block["date"])}</text>')
            ty += date_adv

        if block["detail_lines"]:
            event_parts.append(tlines(
                block["detail_lines"], x, ty, t_size("caption"), S["muted"],
                "middle", str(t_weight("caption")), t_lh("caption"),
                role="caption"))

        parts.append(f'<g data-pptx-role="group" data-node-id="event-{i}">\n    '
                    + "\n    ".join(event_parts) + '\n  </g>')

    return svg("\n  ".join(parts))


def render_conclusion(sl: dict, meta: dict) -> str:
    title       = sl.get("title", "Conclusion & Next Steps")
    conclusions = sl.get("conclusions", [])
    next_steps  = sl.get("next_steps", [])
    footer      = meta.get("footer", "")
    c_head      = sl.get("conclusions_heading", "Conclusions")
    n_head      = sl.get("next_steps_heading",  "Next Steps")

    parts = [frame(title, footer)]

    def block(items, px, py, pw, ph, color, heading, numbered=True) -> list:
        out = []
        out.append(f'<rect x="{px}" y="{py}" width="{pw}" height="{ph}" '
                   f'rx="6" fill="{S["card"]}" stroke="{S["border"]}" stroke-width="1.5"/>')
        out.append(f'<rect x="{px}" y="{py}" width="{pw}" height="5" rx="4" fill="{color}"/>')
        out.append(f'<text x="{px + 24}" y="{py + 24 + t_size("takeaway"):g}" '
                   f'font-size="{t_size("takeaway"):g}" '
                   f'font-weight="{t_weight("takeaway")}" '
                   f'data-style-role="takeaway" '
                   f'fill="{color}">{esc(heading)}</text>')
        rule_y = py + 24 + t_size("takeaway") + 14
        out.append(f'<line x1="{px + 24}" y1="{rule_y:g}" '
                   f'x2="{px + pw - 24}" y2="{rule_y:g}" '
                   f'stroke="{S["divider"]}" stroke-width="1"/>')

        iy = rule_y + 24 + t_size("body")
        body_adv = t_size("body") * t_lh("body")
        for idx, item in enumerate(items):
            lines = wrap_to_width(str(item), pw - 96, "body")
            if numbered:
                r = t_size("footnote") * 0.85
                out.append(f'<circle cx="{px + 32}" '
                           f'cy="{iy - t_size("body") * 0.32:g}" '
                           f'r="{r:g}" fill="{color}"/>')
                out.append(f'<text x="{px + 32}" '
                           f'y="{iy - t_size("body") * 0.32 + r * 0.55:g}" '
                           f'font-size="{t_size("footnote"):g}" font-weight="700" '
                           f'data-style-role="footnote" '
                           f'fill="{S["white"]}" text-anchor="middle">{idx + 1}</text>')
                out.append(tlines(lines, px + 32 + r + 16, iy, t_size("body"),
                                  S["body"], "start", str(t_weight("body")),
                                  t_lh("body"), role="body"))
            else:
                out.append(f'<circle cx="{px + 32}" '
                           f'cy="{iy - t_size("body") * 0.30:g}" '
                           f'r="{t_size("body") * 0.28:g}" fill="{color}"/>')
                out.append(tlines(lines, px + 32 + t_size("body") * 0.28 + 16, iy,
                                  t_size("body"), S["body"], "start",
                                  str(t_weight("body")), t_lh("body"),
                                  role="body"))
            iy += len(lines) * body_adv + S["spacing"][2]
        return out

    panel_top = S["safe"]["top"] + t_size("slide_title") + 40
    panel_h = S["h"] - panel_top - S["safe"]["bottom"] - 16
    panel_w = (S["w"] - 2 * S["safe"]["left"] - 40) / 2
    parts += block(conclusions, S["safe"]["left"], panel_top, panel_w, panel_h,
                   S["accent"], c_head, numbered=False)
    parts += block(next_steps, S["safe"]["left"] + panel_w + 40, panel_top,
                   panel_w, panel_h, S["good"], n_head, numbered=True)

    return svg("\n  ".join(parts))


# ── Dispatch ──────────────────────────────────────────────────────────────────

RENDERERS = {
    "title":         render_title,
    "bullet_list":   render_bullet_list,
    "bar_chart":     render_bar_chart,
    "line_chart":    render_line_chart,
    "pie_chart":     render_pie_chart,
    "table":         render_table,
    "metric_cards":  render_metric_cards,
    "two_column":    render_two_column,
    "timeline":      render_timeline,
    "conclusion":    render_conclusion,
}


def generate_slide(sl: dict, meta: dict) -> str:
    renderer = RENDERERS.get(sl["type"])
    if not renderer:
        raise ValueError(
            f"Unknown type: {sl['type']!r}. Available: {list(RENDERERS)}")
    return renderer(sl, meta)


# ── SVG → PNG conversion (fallback chain) ────────────────────────────────────

def svg_to_png(svg_path: str, png_path: str, width: int = 2400, height: int = 1350) -> bool:
    """Convert one SVG to PNG. Tries cairosvg → inkscape → rsvg-convert."""
    # 1. cairosvg (pip install cairosvg)
    try:
        import cairosvg
        cairosvg.svg2png(url=svg_path, write_to=png_path,
                         output_width=width, output_height=height)
        return True
    except ImportError:
        pass
    except Exception as e:
        print(f"    cairosvg error: {e}")

    import subprocess, shutil

    # 2. inkscape (system package)
    if shutil.which("inkscape"):
        r = subprocess.run(
            ["inkscape", svg_path,
             f"--export-filename={png_path}",
             f"--export-width={width}", f"--export-height={height}"],
            capture_output=True,
        )
        if r.returncode == 0:
            return True

    # 3. rsvg-convert (librsvg)
    if shutil.which("rsvg-convert"):
        r = subprocess.run(
            ["rsvg-convert", svg_path,
             "-w", str(width), "-h", str(height), "-o", png_path],
            capture_output=True,
        )
        if r.returncode == 0:
            return True

    return False


# ── SVG dir → PPTX ───────────────────────────────────────────────────────────

def to_pptx(svg_dir: str, output_path: str) -> None:
    """Convert all slide*.svg in svg_dir (sorted) into a single PPTX file."""
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except ImportError:
        print("Missing dependency: python-pptx")
        print("Install with: pip install python-pptx")
        return

    import glob, tempfile

    svgs = sorted(glob.glob(os.path.join(svg_dir, "slide*.svg")))
    if not svgs:
        print(f"No slide*.svg found in {svg_dir}")
        return

    # Check at least one PNG converter is available before starting
    import shutil
    has_converter = (
        _try_import("cairosvg")
        or shutil.which("inkscape")
        or shutil.which("rsvg-convert")
    )
    if not has_converter:
        print("No SVG→PNG converter found. Install one of:")
        print("  pip install cairosvg")
        print("  sudo apt install inkscape   (or librsvg2-bin for rsvg-convert)")
        return

    prs = Presentation()
    prs.slide_width  = Inches(13.333)   # 16:9, matches 1200px wide
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]        # blank layout

    print(f"Converting {len(svgs)} slides → {output_path}")

    with tempfile.TemporaryDirectory() as tmp:
        for svg_path in svgs:
            name     = os.path.basename(svg_path).replace(".svg", ".png")
            png_path = os.path.join(tmp, name)

            ok = svg_to_png(svg_path, png_path)
            if not ok:
                print(f"  ✗ {os.path.basename(svg_path)} (conversion failed, skipped)")
                continue

            slide = prs.slides.add_slide(blank)
            slide.shapes.add_picture(
                png_path,
                left=0, top=0,
                width=prs.slide_width,
                height=prs.slide_height,
            )
            print(f"  ✓ {os.path.basename(svg_path)}")

    prs.save(output_path)
    size_kb = os.path.getsize(output_path) // 1024
    print(f"\n✅ PPTX saved: {output_path}  ({size_kb} KB, {len(prs.slides)} slides)")
    print("\nNote: each slide in the PPTX is an embedded image; elements cannot be edited individually in PowerPoint.")
    print("To edit a single slide, modify the corresponding SVG and re-run with --to-pptx.")


def _try_import(name: str) -> bool:
    try:
        __import__(name)
        return True
    except ImportError:
        return False


# ── CLI ───────────────────────────────────────────────────────────────────────

def _authorize_production(deck_id: str, project_root: Path | None) -> Path | None:
    """Authorize the CLI before it reads inputs or creates output paths."""
    try:
        root = (
            Path(project_root).resolve()
            if project_root is not None
            else find_project_root(Path.cwd())
        )
        assert_production_allowed(root, deck_id)
    except ProductionGateError as exc:
        print(
            json.dumps(
                {
                    "error": type(exc).__name__,
                    "predicate": exc.predicate,
                    "deck_id": exc.deck_id,
                    "blockers": exc.blockers,
                }
            )
        )
        return None
    except Exception as exc:  # noqa: BLE001 - CLI failures stay machine-readable
        print(
            json.dumps(
                {
                    "error": "ProductionGateError",
                    "predicate": "production_allowed",
                    "deck_id": deck_id,
                    "blockers": [{"reason": "project_root_invalid", "message": str(exc)}],
                }
            )
        )
        return None
    return root


def main(arguments: Sequence[str] | None = None) -> int:
    ap = argparse.ArgumentParser(description="Generate SVG research slides from JSON")
    ap.add_argument("--data",     help="Path to slide_data.json")
    ap.add_argument("--out",      help="Output directory for SVG files")
    ap.add_argument("--deck-id", required=True, help="Approved presentation Deck identifier")
    ap.add_argument(
        "--project-root",
        type=Path,
        default=None,
        help="Project root (defaults to the nearest ancestor of the current directory)",
    )
    ap.add_argument("--slide",    type=int, default=None, help="Render only slide N")
    ap.add_argument("--style",    metavar="FILE",      default=None,
                    help="Style .md file whose colours are composed into the "
                         "token set before rendering (see "
                         "references/styles/STYLES.md in skill bundle)")
    ap.add_argument("--tokens", metavar="FILE", type=Path, default=None,
                    help="Design-token .tokens.yaml file "
                         "(default: references/tokens/default.tokens.yaml)")
    ap.add_argument("--to-pptx",  metavar="SVG_DIR",
                    help="Convert all slide*.svg in SVG_DIR to PPTX (skips SVG generation)")
    ap.add_argument("--pptx-out", metavar="FILE", default=None,
                    help="Output PPTX path (default: SVG_DIR/deck.pptx)")
    args = ap.parse_args(arguments)

    project_root = _authorize_production(args.deck_id, args.project_root)
    if project_root is None:
        return 1

    # ── Mode: convert existing SVGs to PPTX ──
    if args.to_pptx:
        out = args.pptx_out or os.path.join(args.to_pptx, "deck.pptx")
        Path(out).parent.mkdir(parents=True, exist_ok=True)
        to_pptx(args.to_pptx, out)
        return 0

    # ── Mode: generate SVGs from JSON ──
    if not args.data or not args.out:
        ap.error("--data and --out are required when not using --to-pptx")

    # One token set per deck: the style override is composed into the tokens
    # and written out before anything reads them, so the renderer, the linter,
    # and the gate all hold this slide to the same contract.
    os.makedirs(args.out, exist_ok=True)
    tokens_path, tokens_digest = effective_tokens(
        args.tokens, args.style, Path(args.out))
    apply_tokens(tokens_path)
    print(f"  [tokens] {tokens_path} sha256={tokens_digest[:12]}")

    with open(args.data, encoding="utf-8") as f:
        data = json.load(f)

    meta   = data.get("meta", {})
    slides = data.get("slides", [])
    os.makedirs(args.out, exist_ok=True)

    generated = 0
    for sl in slides:
        if args.slide is not None and sl.get("index") != args.slide:
            continue
        content  = generate_slide(sl, meta)
        idx      = sl.get("index", 0)
        filename = f"slide{idx:02d}_{sl['type']}.svg"
        path     = os.path.join(args.out, filename)
        with open(path, "w", encoding="utf-8") as f:
            f.write(content)
        print(f"  ✓ {filename}")
        generated += 1

    print(f"\n{generated} slide(s) written to {args.out}")

    import shutil
    shutil.copy2(args.data, os.path.join(args.out, "slide_data.json"))

    return 0


if __name__ == "__main__":
    raise SystemExit(main())
