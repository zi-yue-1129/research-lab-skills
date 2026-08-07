"""Tests for validate_pptx_structure.py."""
import json
import subprocess
import sys
import zipfile
from pathlib import Path
from typing import Optional

from pptx import Presentation
from pptx.util import Inches

SCRIPT = Path(__file__).resolve().parent.parent / "validate_pptx_structure.py"


def _build_pptx(path: Path, n_slides: int, picture_path: Optional[Path] = None) -> None:
    prs = Presentation()
    blank_layout = prs.slide_layouts[6]
    for i in range(n_slides):
        slide = prs.slides.add_slide(blank_layout)
        if picture_path is not None:
            slide.shapes.add_picture(str(picture_path), 0, 0, width=prs.slide_width, height=prs.slide_height)
        else:
            box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
            box.text_frame.text = f"Slide {i}"
    prs.save(str(path))


def _make_test_png(path: Path) -> Path:
    # A valid, minimal 1x1 RGBA PNG. Note: the brief's literal hex string for
    # this fixture had an odd character count (135 hex digits), which
    # bytes.fromhex() can never parse (it requires an even number of digits
    # to form whole bytes) -- confirmed empirically, not a python-pptx
    # version quirk. Regenerated here as a correct minimal 1x1 PNG (IHDR +
    # IDAT + IEND, each with a valid CRC) that python-pptx's add_picture
    # successfully embeds.
    png_bytes = bytes.fromhex(
        "89504e470d0a1a0a0000000d49484452000000010000000108060000"
        "001f15c4890000000b49444154789c6360000200000500017a5eab3f"
        "0000000049454e44ae426082"
    )
    path.write_bytes(png_bytes)
    return path


def _run(*args: str) -> subprocess.CompletedProcess:
    return subprocess.run([sys.executable, str(SCRIPT), *args], capture_output=True, text=True, check=False)


def test_valid_native_pptx_passes(tmp_path: Path) -> None:
    pptx_path = tmp_path / "deck.pptx"
    _build_pptx(pptx_path, n_slides=3)

    result = _run("--pptx", str(pptx_path), "--expected-slides", "3", "--json")

    assert result.returncode == 0, result.stderr
    data = json.loads(result.stdout)
    assert data["status"] == "passed"
    assert data["slide_count_actual"] == 3
    assert data["relationship_violations"] == []


def test_slide_count_mismatch_fails(tmp_path: Path) -> None:
    pptx_path = tmp_path / "deck.pptx"
    _build_pptx(pptx_path, n_slides=2)

    result = _run("--pptx", str(pptx_path), "--expected-slides", "5", "--json")

    assert result.returncode == 1
    data = json.loads(result.stdout)
    assert data["status"] == "failed"
    assert data["slide_count_actual"] == 2
    assert data["slide_count_expected"] == 5


def test_broken_relationship_target_is_caught(tmp_path: Path) -> None:
    png_path = _make_test_png(tmp_path / "img.png")
    pptx_path = tmp_path / "deck.pptx"
    _build_pptx(pptx_path, n_slides=1, picture_path=png_path)

    corrupted_path = tmp_path / "corrupted.pptx"
    with zipfile.ZipFile(pptx_path) as src, zipfile.ZipFile(corrupted_path, "w") as dst:
        for item in src.infolist():
            if item.filename.startswith("ppt/media/"):
                continue
            dst.writestr(item, src.read(item.filename))

    result = _run("--pptx", str(corrupted_path), "--expected-slides", "1", "--json")

    assert result.returncode == 1
    data = json.loads(result.stdout)
    assert data["status"] == "failed"
    assert len(data["relationship_violations"]) >= 1


def test_editability_mismatch_is_caught(tmp_path: Path) -> None:
    png_path = _make_test_png(tmp_path / "img.png")
    pptx_path = tmp_path / "deck.pptx"
    _build_pptx(pptx_path, n_slides=1, picture_path=png_path)
    declared_path = tmp_path / "declared.json"
    declared_path.write_text(json.dumps({"0": "native"}))

    result = _run("--pptx", str(pptx_path), "--expected-slides", "1", "--declared-editability", str(declared_path), "--json")

    assert result.returncode == 1
    data = json.loads(result.stdout)
    assert data["status"] == "failed"
    assert len(data["editability_mismatches"]) == 1
    assert data["editability_mismatches"][0]["declared"] == "native"
    assert data["editability_mismatches"][0]["observed"] == "raster"


def test_editability_match_passes(tmp_path: Path) -> None:
    pptx_path = tmp_path / "deck.pptx"
    _build_pptx(pptx_path, n_slides=1)  # textbox only -> native
    declared_path = tmp_path / "declared.json"
    declared_path.write_text(json.dumps({"0": "native"}))

    result = _run("--pptx", str(pptx_path), "--expected-slides", "1", "--declared-editability", str(declared_path), "--json")

    assert result.returncode == 0, result.stderr
    data = json.loads(result.stdout)
    assert data["status"] == "passed"
    assert data["editability_mismatches"] == []


def test_hybrid_editability_is_classified_correctly(tmp_path: Path) -> None:
    png_path = _make_test_png(tmp_path / "img.png")
    pptx_path = tmp_path / "deck.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(str(png_path), 0, 0, width=Inches(1), height=Inches(1))
    box = slide.shapes.add_textbox(Inches(2), Inches(2), Inches(4), Inches(1))
    box.text_frame.text = "Hybrid slide"
    prs.save(str(pptx_path))
    declared_path = tmp_path / "declared.json"
    declared_path.write_text(json.dumps({"0": "hybrid"}))

    result = _run("--pptx", str(pptx_path), "--expected-slides", "1", "--declared-editability", str(declared_path), "--json")

    assert result.returncode == 0, result.stderr
    data = json.loads(result.stdout)
    assert data["status"] == "passed"
    assert data["editability_mismatches"] == []


def test_not_a_pptx_file_reports_error(tmp_path: Path) -> None:
    bogus_path = tmp_path / "not-a-pptx.pptx"
    bogus_path.write_text("not a zip file at all")

    result = _run("--pptx", str(bogus_path), "--expected-slides", "1", "--json")

    assert result.returncode == 1
    data = json.loads(result.stdout)
    assert data["status"] == "failed"
    assert "error" in data
