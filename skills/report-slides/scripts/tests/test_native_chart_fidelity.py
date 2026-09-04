"""Tests that a native chart says what its SVG preview said.

The `data-pptx-role="chart"` marker replaces the whole hand-drawn group with a
real PPTX chart, so everything the group carried and the chart does not
declare is lost. On the verification deck that was the value above every bar,
the series name, and the note under the plot -- `Median of five runs.` was in
the SVG and in no shape of the exported slide.
"""
import sys
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Emu

sys.path.insert(0, str(Path(__file__).resolve().parent.parent))

import generate_slides as gs
import pptx_native
from chart_labels import number_format

_SKILL_DIR = Path(__file__).resolve().parents[2]
_DEFAULT_TOKENS = _SKILL_DIR / "references" / "tokens" / "default.tokens.yaml"
_BBOX = (0, 0, Emu(6_096_000), Emu(3_429_000))
_STYLE = {"font": "Inter, 'DejaVu Sans', sans-serif"}
_META = {"footer": "deck 1/1"}


@pytest.fixture(autouse=True)
def _tokens_applied() -> None:
    """Apply the default token file before each test."""
    gs.apply_tokens(_DEFAULT_TOKENS)


def _chart(chart_type: str, **kwargs):
    """Build a native chart and return it.

    Args:
        chart_type: "bar", "line" or "pie".
        **kwargs: Passed through to `add_native_chart`.

    Returns:
        The created PPTX chart.
    """
    prs = Presentation()
    prs.slide_width = Emu(12_192_000)
    prs.slide_height = Emu(6_858_000)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    series = [{"label": "after", "values": [340.0, 210.0]}]
    frame = pptx_native.add_native_chart(
        slide, chart_type, ["p50", "p99"], series, _BBOX, _STYLE, **kwargs)
    return frame.chart


def test_a_bar_chart_shows_the_value_above_each_bar() -> None:
    """The SVG draws the value on every bar; the native chart drew none."""
    assert _chart("bar", y_max=400).plots[0].has_data_labels


def test_the_value_labels_carry_the_unit() -> None:
    """A declared unit reaches the native labels, not only the SVG ones."""
    chart = _chart("bar", y_max=400, unit="%")
    assert chart.plots[0].data_labels.number_format == '0.0"%"'
    assert not chart.plots[0].data_labels.number_format_is_linked


def test_the_value_axis_ticks_carry_the_unit() -> None:
    """The axis is labelled like the SVG axis beside it."""
    chart = _chart("bar", y_max=400, unit="%")
    assert chart.value_axis.tick_labels.number_format == '0"%"'


def test_no_unit_leaves_a_plain_number_format() -> None:
    """Without a unit the labels are bare numbers, as the SVG ticks are."""
    chart = _chart("bar", y_max=400)
    assert chart.plots[0].data_labels.number_format == "0.0"
    assert chart.value_axis.tick_labels.number_format == "0"


def test_a_sub_unit_axis_gets_the_finer_native_format() -> None:
    """A 0-1 axis is formatted at the precision its ticks need."""
    chart = _chart("bar", y_max=1.0)
    assert chart.value_axis.tick_labels.number_format == "0.0"
    assert chart.plots[0].data_labels.number_format == "0.00"


def test_a_percent_sign_is_quoted_so_the_value_is_not_rescaled() -> None:
    """An unquoted `%` makes Excel multiply the value by a hundred.

    A bar of 98.4 would print as 9840%.
    """
    assert number_format("%", 1) == '0.0"%"'


def test_a_single_series_still_gets_its_legend() -> None:
    """The SVG legend names the series, so the native chart must too."""
    assert _chart("bar", y_max=400).has_legend


@pytest.mark.parametrize("chart_type,renderer", [
    ("bar_chart", gs.render_bar_chart),
    ("line_chart", gs.render_line_chart),
])
def test_the_chart_note_is_emitted_outside_the_native_marker(
        chart_type: str, renderer: object) -> None:
    """A note inside the marker group is discarded when the chart replaces it.

    The note is the only place a deck states its sample size or its method, so
    losing it silently drops the one thing on the slide that qualifies the
    numbers.
    """
    slide = {"index": 2, "type": chart_type, "title": "Throughput",
             "categories": ["p50", "p99"],
             "series": [{"label": "after", "values": [340.0, 210.0]}],
             "y_max": 400, "note": "Median of five runs."}
    markup = renderer(slide, _META)
    group_start = markup.index('<g data-pptx-role="chart"')
    group_end = markup.index("</g>", group_start)
    assert "Median of five runs." in markup
    assert "Median of five runs." not in markup[group_start:group_end]


def test_the_pie_note_is_emitted_outside_the_native_marker() -> None:
    """The pie layout carries its note the same way."""
    slide = {"index": 4, "type": "pie_chart", "title": "Split",
             "categories": ["Train", "Dev"], "values": [70, 30],
             "note": "n=1200"}
    markup = gs.render_pie_chart(slide, _META)
    group_start = markup.index('<g data-pptx-role="chart"')
    group_end = markup.index("</g>", group_start)
    assert "n=1200" in markup
    assert "n=1200" not in markup[group_start:group_end]
