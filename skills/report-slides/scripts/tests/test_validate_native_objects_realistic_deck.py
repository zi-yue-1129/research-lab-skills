"""Realistic-deck regression tests for validate_native_objects.py.

The per-heuristic unit tests in `test_validate_native_objects.py` use
hand-written minimal SVG snippets, which cannot show whether the detector
false-positives on the skill's *own* routine output. This module drives the
real `generate_slides.py` renderers for every non-chart slide type plus a
real table and chart, converts them through the real
`svg_to_pptx.converter.convert_file`, and asserts both scans report zero
findings -- while the deliberately-unmarked hand-drawn fixtures still fire.
"""

import json
import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent.parent))

from generate_slides import generate_slide  # noqa: E402
from svg_to_pptx.converter import convert_file  # noqa: E402
from validate_native_objects import (  # noqa: E402
    load_thresholds,
    scan_pptx,
    scan_svg,
)

REALISTIC_DECK = {
    "meta": {"footer": "Native objects regression deck", "date": "2026-08-13"},
    "slides": [
        {"index": 1, "type": "title", "title": "Native PPTX Objects",
         "subtitle": "A realistic multi-slide-type deck",
         "author": "report-slides", "date": "2026-08-13"},
        {"index": 2, "type": "bullet_list", "title": "Agenda",
         "bullets": ["Background and motivation", "Method overview",
                     "Experimental results", "Discussion and limitations"]},
        {"index": 3, "type": "metric_cards", "title": "Headline Metrics",
         "metrics": [
             {"label": "QWK", "value": "0.83", "change": "+0.05"},
             {"label": "Accuracy", "value": "91%", "change": "+2.1"},
             {"label": "Latency", "value": "120ms", "change": "-8ms"},
             {"label": "Cost", "value": "$0.02", "change": "-12%"},
         ]},
        {"index": 4, "type": "table", "title": "Model Comparison",
         "columns": ["Model", "QWK", "Delta"],
         "rows": [["GPT-4", "0.671", "+0.02"],
                  ["ZH-BERT-FT", "0.828", "+0.16"],
                  ["Baseline", "0.512", "-0.11"]],
         "highlight_col": 2},
        {"index": 5, "type": "bar_chart", "title": "Scores by Split",
         "categories": ["Train", "Dev", "Test"],
         "series": [
             {"label": "EN", "color": "#1e3a5f", "values": [88.0, 85.5, 84.2]},
             {"label": "ZH", "color": "#059669", "values": [86.0, 83.1, 82.7]},
         ],
         "y_max": 100},
        {"index": 6, "type": "line_chart", "title": "Training Curve",
         "categories": ["E1", "E2", "E3", "E4"],
         "series": [{"label": "Loss", "values": [80, 60, 45, 38]}],
         "y_max": 100},
        {"index": 7, "type": "pie_chart", "title": "Corpus Split",
         "categories": ["Train", "Dev", "Test"], "values": [70, 15, 15]},
        {"index": 8, "type": "timeline", "title": "Project Timeline",
         "events": [
             {"label": "Kickoff", "date": "2025-01"},
             {"label": "Prototype", "date": "2025-04"},
             {"label": "Evaluation", "date": "2025-08"},
             {"label": "Ship", "date": "2025-12"},
         ]},
        {"index": 9, "type": "two_column", "title": "Strengths vs Risks",
         "left": {"title": "Strengths",
                  "content": ["Fast inference", "Small footprint"]},
         "right": {"title": "Risks",
                   "content": ["Domain shift", "Annotation noise"]}},
        {"index": 10, "type": "conclusion", "title": "Conclusion",
         "conclusions": ["Native objects improve editability",
                         "Preview fidelity is unchanged"],
         "next_steps": ["Ship to main", "Monitor adoption"]},
    ],
}


def _build_deck(tmp_path: Path) -> Path:
    """Render the realistic deck through the real Path A renderers."""
    out_dir = tmp_path / "slides"
    out_dir.mkdir()
    meta = REALISTIC_DECK["meta"]
    for slide in REALISTIC_DECK["slides"]:
        svg_text = generate_slide(slide, meta)
        name = f'slide{slide["index"]:02d}_{slide["type"]}.svg'
        (out_dir / name).write_text(svg_text, encoding="utf-8")
    # generate_slides.main() copies slide_data.json next to the SVGs so the
    # data-pptx-source="slide_data.json#N" markers resolve at conversion time.
    (out_dir / "slide_data.json").write_text(
        json.dumps(REALISTIC_DECK), encoding="utf-8")
    return out_dir


def _write_svg(tmp_path: Path, name: str, body: str) -> Path:
    path = tmp_path / name
    path.write_text(
        f'<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 1200 675">{body}</svg>',
        encoding="utf-8",
    )
    return path


def test_realistic_deck_svg_scan_reports_zero_findings(tmp_path: Path):
    out_dir = _build_deck(tmp_path)
    thresholds = load_thresholds()
    findings = []
    for svg_path in sorted(out_dir.glob("slide*.svg")):
        findings.extend(scan_svg(svg_path, thresholds))
    assert findings == [], [(f.slide, f.kind) for f in findings]


def test_realistic_deck_pptx_scan_reports_zero_findings(tmp_path: Path):
    out_dir = _build_deck(tmp_path)
    pptx_path = tmp_path / "deck.pptx"
    convert_file(str(out_dir), str(pptx_path))
    thresholds = load_thresholds()
    findings = scan_pptx(pptx_path, thresholds)
    assert findings == [], [(f.slide, f.kind) for f in findings]


def test_realistic_deck_pptx_still_contains_the_native_objects(tmp_path: Path):
    """Zero findings must mean 'correctly native', not 'nothing was built'."""
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    out_dir = _build_deck(tmp_path)
    pptx_path = tmp_path / "deck.pptx"
    convert_file(str(out_dir), str(pptx_path))

    prs = Presentation(str(pptx_path))
    # Slide order matches REALISTIC_DECK: index 4 = table, 5/6/7 = bar/line/pie
    # charts, 8 = timeline (one native Group per event).
    assert any(getattr(s, "has_table", False) for s in prs.slides[3].shapes)
    for chart_slide in (4, 5, 6):
        assert any(getattr(s, "has_chart", False)
                   for s in prs.slides[chart_slide].shapes)
    groups = [s for s in prs.slides[7].shapes
              if s.shape_type == MSO_SHAPE_TYPE.GROUP]
    assert len(groups) == 4


def test_unmarked_hand_drawn_table_still_fires(tmp_path: Path):
    grid = "".join(
        f'<rect x="{60 + c * 200}" y="{75 + r * 50}" width="190" height="45" fill="#eee"/>'
        for r in range(3) for c in range(3)
    )
    svg_path = _write_svg(tmp_path, "slide01_table.svg", grid)
    findings = scan_svg(svg_path, load_thresholds())
    assert [f.kind for f in findings] == ["table_like"]


def test_unmarked_hand_drawn_bars_still_fire(tmp_path: Path):
    bars = "".join(
        f'<rect x="{130 + i * 100}" y="{500 - i * 40}" width="60" '
        f'height="{20 + i * 40}" fill="#369"/>'
        for i in range(4)
    )
    svg_path = _write_svg(tmp_path, "slide02_bar.svg", bars)
    findings = scan_svg(svg_path, load_thresholds())
    assert [f.kind for f in findings] == ["chart_like"]


def test_unmarked_node_cluster_still_fires(tmp_path: Path):
    cluster = (
        '<g><rect x="100" y="100" width="200" height="80" fill="#eee"/>'
        '<circle cx="120" cy="120" r="10" fill="#333"/>'
        '<text x="200" y="140">Encoder</text></g>'
    )
    svg_path = _write_svg(tmp_path, "slide03_arch.svg", cluster)
    findings = scan_svg(svg_path, load_thresholds())
    assert [f.kind for f in findings] == ["node_cluster_like"]


def test_unmarked_node_cluster_in_pptx_still_fires(tmp_path: Path):
    """A converted, unmarked box+icon+label node must fire in PPTX mode too."""
    cluster = (
        '<rect x="100" y="100" width="200" height="80" fill="#eeeeee"/>'
        '<circle cx="130" cy="130" r="12" fill="#333333"/>'
        '<text x="200" y="140">Encoder</text>'
    )
    out_dir = tmp_path / "svgs"
    out_dir.mkdir()
    _write_svg(out_dir, "slide01_arch.svg", cluster)
    pptx_path = tmp_path / "deck.pptx"
    convert_file(str(out_dir), str(pptx_path))
    findings = scan_pptx(pptx_path, load_thresholds())
    assert [f.kind for f in findings] == ["node_cluster_like"]
