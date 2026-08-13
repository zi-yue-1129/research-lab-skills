"""End-to-end integration test for the native-object PPTX pipeline.

Exercises the full chain: SVG preview generation (generate_slides.py's
render_table/render_bar_chart/render_pie_chart) -> SVG-to-PPTX conversion
(svg_to_pptx.converter.convert_file) -> real PPTX inspection (python-pptx)
-> the safety-net detector (validate_native_objects.scan_pptx). This proves
the four independently-tested components (Tasks 1, 3, 4, 6) actually agree
with each other when wired together, which none of their per-task unit
tests could verify in isolation.
"""

import json
import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent.parent))

from pptx import Presentation
from pptx.enum.chart import XL_CHART_TYPE

from generate_slides import render_table, render_bar_chart, render_pie_chart
from svg_to_pptx.converter import convert_file
from validate_native_objects import scan_pptx, load_thresholds


SLIDE_DATA = {
    "meta": {"footer": "Integration test"},
    "slides": [
        {"index": 0, "type": "table", "title": "Results",
         "columns": ["Model", "QWK"], "rows": [["GPT-4", "0.671"], ["BERT", "0.836"]]},
        {"index": 1, "type": "bar_chart", "title": "Scores",
         "categories": ["Q1", "Q2"],
         "series": [{"label": "EN", "color": "#1e3a5f", "values": [0.8, 0.9]}],
         "y_max": 1.0},
        {"index": 2, "type": "pie_chart", "title": "Share",
         "categories": ["A", "B", "C"], "values": [50, 30, 20]},
    ],
}


def test_end_to_end_table_and_chart_slides_produce_native_objects(tmp_path: Path):
    out_dir = tmp_path / "svgs"
    out_dir.mkdir()
    (out_dir / "slide_data.json").write_text(json.dumps(SLIDE_DATA), encoding="utf-8")

    renderers = {"table": render_table, "bar_chart": render_bar_chart,
                "pie_chart": render_pie_chart}
    meta = SLIDE_DATA["meta"]
    for sl in SLIDE_DATA["slides"]:
        svg_text = renderers[sl["type"]](sl, meta)
        (out_dir / f'slide{sl["index"]:02d}_{sl["type"]}.svg').write_text(
            svg_text, encoding="utf-8"
        )

    pptx_path = tmp_path / "deck.pptx"
    convert_file(str(out_dir), str(pptx_path))

    prs = Presentation(str(pptx_path))
    assert len(prs.slides) == 3
    has_table = any(getattr(s, "has_table", False) for s in prs.slides[0].shapes)
    assert has_table

    bar_frame = next(
        (s for s in prs.slides[1].shapes if getattr(s, "has_chart", False)), None
    )
    assert bar_frame is not None
    assert bar_frame.chart.chart_type == XL_CHART_TYPE.COLUMN_CLUSTERED

    pie_frame = next(
        (s for s in prs.slides[2].shapes if getattr(s, "has_chart", False)), None
    )
    assert pie_frame is not None
    assert pie_frame.chart.chart_type == XL_CHART_TYPE.PIE

    thresholds = load_thresholds()
    findings = scan_pptx(pptx_path, thresholds)
    assert findings == []
