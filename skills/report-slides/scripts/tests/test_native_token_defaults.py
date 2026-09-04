"""Tests that the native table and chart route obeys the token contract.

`converter._pptx_style` carried its own colour literals and a proprietary font
stack, and `pptx_native` set every table cell at a hardcoded 13pt and coloured
chart series from a five-entry tuple of its own. Nothing in that path had ever
been pointed at the design tokens, so a deck's own palette and type scale
stopped at the edge of its tables and charts.
"""
import sys
from pathlib import Path

import pytest
from lxml import etree
from pptx import Presentation
from pptx.util import Emu, Pt

sys.path.insert(0, str(Path(__file__).resolve().parent.parent))

import pptx_native
from design_tokens import default_tokens
from svg_to_pptx.converter import SvgConverter

_BBOX = (0, 0, Emu(6_096_000), Emu(3_429_000))


def _bare_marker():
    """Return a native-object marker that declares no style of its own."""
    return etree.fromstring(
        '<g data-pptx-role="table" data-pptx-source="slide_data.json#1" '
        'data-pptx-bbox="48,108,1104,400"/>'
    )


def _defaults() -> dict:
    """Return the style a marker with no `data-pptx-style` resolves to."""
    converter = SvgConverter.__new__(SvgConverter)
    converter.svg_path = "<test>"
    return converter._pptx_style(_bare_marker())


@pytest.mark.parametrize("style_key,color_role", [
    ("accent", "primary"),
    ("white", "bg"),
    ("card", "card"),
    ("bg", "bg"),
    ("body", "body"),
    ("good", "positive"),
    ("danger", "danger"),
])
def test_the_native_style_defaults_are_the_token_colours(
        style_key: str, color_role: str) -> None:
    """Each default colour is the token role it stands for.

    The literals had drifted: `good` was `#059669` against the contract's
    `#047857`, and `danger` `#dc2626` against `#b91c1c`.
    """
    assert _defaults()[style_key] == default_tokens().color(color_role)


def test_the_native_style_default_font_is_the_token_stack() -> None:
    """The default font stack is the deck's, not a proprietary guess.

    `'Helvetica Neue'` is installed on neither Linux nor the token stack, so
    the table asked for a face nothing could render.
    """
    font = _defaults()["font"]
    assert font == default_tokens().font_stack("sans")
    assert "Helvetica Neue" not in font


def _table(style: dict):
    """Build a two-row native table and return it.

    Args:
        style: The style dict handed to the native route.

    Returns:
        The created PPTX table.
    """
    prs = Presentation()
    prs.slide_width = Emu(12_192_000)
    prs.slide_height = Emu(6_858_000)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    frame = pptx_native.add_native_table(
        slide, ["Metric", "After"], [["Accuracy", "100%"]], _BBOX, style)
    return frame.table


def _cell_size(table, row: int, col: int):
    """Return the point size of a cell's only run."""
    return table.cell(row, col).text_frame.paragraphs[0].runs[0].font.size


def test_table_cells_are_set_at_the_type_scale() -> None:
    """A cell's type comes from the contract, not from a 13pt literal.

    The SVG table renderer sets its header at the `node_label` role and its
    body rows at `body`. The native route set both at 13pt, a size the token
    scale does not contain, so the two renderings of the same table disagreed.
    """
    tokens = default_tokens()
    table = _table(dict(_defaults()))
    assert _cell_size(table, 0, 0) == Pt(tokens.type_role("node_label").size)
    assert _cell_size(table, 1, 0) == Pt(tokens.type_role("body").size)


def test_a_deck_can_override_the_cell_sizes() -> None:
    """A deck's own resolved sizes travel with its marker."""
    style = dict(_defaults())
    style["header_size"] = 15
    style["body_size"] = 17
    table = _table(style)
    assert _cell_size(table, 0, 0) == Pt(15)
    assert _cell_size(table, 1, 0) == Pt(17)


def _chart_series_fills(series: list):
    """Return each series' solid fill colour from a native bar chart."""
    prs = Presentation()
    prs.slide_width = Emu(12_192_000)
    prs.slide_height = Emu(6_858_000)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    frame = pptx_native.add_native_chart(
        slide, "bar", ["A", "B"], series, _BBOX, dict(_defaults()))
    return [str(s.format.fill.fore_color.rgb) for s in frame.chart.series]


def test_chart_series_fall_back_to_the_token_palette() -> None:
    """Uncoloured series take the contract's chart palette.

    They took a five-entry tuple private to `pptx_native`, so a deck whose
    tokens defined a palette never saw it used.
    """
    series = [{"label": f"s{i}", "values": [1.0, 2.0]} for i in range(3)]
    palette = default_tokens().chart_palette()
    assert _chart_series_fills(series) == [
        color.lstrip("#").upper() for color in palette[:3]]


def test_an_explicit_series_colour_still_wins() -> None:
    """A series that names its colour keeps it."""
    series = [{"label": "s0", "values": [1.0, 2.0], "color": "#123456"}]
    assert _chart_series_fills(series) == ["123456"]
