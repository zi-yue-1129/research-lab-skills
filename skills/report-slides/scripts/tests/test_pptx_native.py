import shutil
import sys
from pathlib import Path

import pytest
sys.path.insert(0, str(Path(__file__).resolve().parent.parent))

from pptx import Presentation
from pptx.util import Emu

import pptx_native
from pptx_native import add_native_table, add_native_chart, add_native_group

STYLE = {
    "accent": "#1e3a5f", "white": "#ffffff", "card": "#f8fafc",
    "bg": "#ffffff", "body": "#374151", "good": "#059669",
    "danger": "#dc2626", "font": "'Helvetica Neue', Arial, sans-serif",
}


def _new_slide():
    prs = Presentation()
    return prs.slides.add_slide(prs.slide_layouts[6])


def test_add_native_table_creates_real_table_with_header_and_rows():
    slide = _new_slide()
    frame = add_native_table(
        slide,
        columns=["Model", "QWK"],
        rows=[["GPT-4", "0.671"], ["ZH-BERT-FT", "0.828"]],
        bbox=(Emu(600000), Emu(750000), Emu(10800000), Emu(4500000)),
        style=STYLE,
    )
    assert frame.has_table
    table = frame.table
    assert len(table.rows) == 3
    assert len(table.columns) == 2
    assert table.cell(0, 0).text_frame.paragraphs[0].runs[0].text == "Model"
    assert table.cell(1, 0).text_frame.paragraphs[0].runs[0].text == "GPT-4"


def test_add_native_table_highlight_col_colors_plus_minus():
    slide = _new_slide()
    frame = add_native_table(
        slide,
        columns=["Model", "Delta"],
        rows=[["A", "+0.05"], ["B", "-0.03"]],
        bbox=(Emu(600000), Emu(750000), Emu(10800000), Emu(4500000)),
        style=STYLE,
        highlight_col=1,
    )
    table = frame.table
    good_run = table.cell(1, 1).text_frame.paragraphs[0].runs[0]
    danger_run = table.cell(2, 1).text_frame.paragraphs[0].runs[0]
    assert str(good_run.font.color.rgb) == "059669"
    assert str(danger_run.font.color.rgb) == "DC2626"


def test_add_native_chart_bar_creates_real_chart():
    slide = _new_slide()
    frame = add_native_chart(
        slide, "bar",
        categories=["Q1", "Q2"],
        series=[{"label": "EN", "color": "#1e3a5f", "values": [0.85, 0.83]}],
        bbox=(Emu(1200000), Emu(900000), Emu(9000000), Emu(4200000)),
        style=STYLE,
        y_max=1.0,
    )
    assert frame.has_chart
    chart = frame.chart
    assert list(chart.plots[0].categories) == ["Q1", "Q2"]
    assert len(chart.series) == 1


def test_add_native_chart_pie_colors_each_point():
    slide = _new_slide()
    frame = add_native_chart(
        slide, "pie",
        categories=["A", "B", "C"],
        series=[{"label": "Share", "values": [50, 30, 20]}],
        bbox=(Emu(1200000), Emu(900000), Emu(9000000), Emu(4200000)),
        style=STYLE,
        colors=["#ff0000", "#00ff00", "#0000ff"],
    )
    assert frame.has_chart
    points = frame.chart.series[0].points
    assert str(points[0].format.fill.fore_color.rgb) == "FF0000"
    assert str(points[2].format.fill.fore_color.rgb) == "0000FF"


def test_add_native_group_groups_two_or_more_shapes():
    slide = _new_slide()
    rect = slide.shapes.add_shape(1, Emu(0), Emu(0), Emu(100000), Emu(100000))
    oval = slide.shapes.add_shape(9, Emu(0), Emu(0), Emu(100000), Emu(100000))
    group = add_native_group(slide, [rect, oval])
    assert group is not None
    assert group.shape_type == 6  # MSO_SHAPE_TYPE.GROUP


def test_add_native_group_returns_none_for_single_shape():
    slide = _new_slide()
    rect = slide.shapes.add_shape(1, Emu(0), Emu(0), Emu(100000), Emu(100000))
    assert add_native_group(slide, [rect]) is None


def test_add_native_chart_raises_on_missing_values():
    """Verify that ValueError is raised when 'values' key is missing from series."""
    slide = _new_slide()
    with pytest.raises(ValueError, match="'values' key is required but missing"):
        add_native_chart(
            slide, "bar",
            categories=["Q1", "Q2"],
            series=[{"label": "EN", "color": "#1e3a5f"}],  # Missing "values"
            bbox=(Emu(1200000), Emu(900000), Emu(9000000), Emu(4200000)),
            style=STYLE,
        )


_NEEDS_FONTCONFIG = pytest.mark.skipif(
    shutil.which("fc-match") is None,
    # fontconfig is an optional external binary; deciding which family in a
    # stack is installed has no meaning without it.
    reason="fontconfig (fc-match) is not installed",
)


@_NEEDS_FONTCONFIG
def test_font_family_skips_an_absent_first_family() -> None:
    """A stack whose first family is missing resolves to an installed one.

    The shipped token stack begins with `Inter`, which is not installed on a
    stock Linux box. Naming it anyway put the deck's tables in one face and
    every other shape in another, and the renderer measured the text with the
    resolved face while the table used something else.
    """
    resolved = pptx_native._font_family(
        "Inter, 'DejaVu Sans', 'Liberation Sans', Arial, sans-serif")
    assert resolved != "Inter"
    assert resolved in ("DejaVu Sans", "Liberation Sans", "Arial")


@_NEEDS_FONTCONFIG
def test_font_family_agrees_with_the_shape_route() -> None:
    """The table route and the shape route must resolve identically.

    `shapes._apply_font` resolves through `resolve_font_stack`. If the table
    route resolved differently, the same deck would carry two faces and no
    test anywhere would notice.
    """
    from fonts import resolve_font_stack

    stack = "Inter, 'DejaVu Sans', 'Liberation Sans', Arial, sans-serif"
    assert pptx_native._font_family(stack) == resolve_font_stack(stack)


@_NEEDS_FONTCONFIG
def test_font_family_keeps_a_multi_word_family_intact() -> None:
    """A quoted multi-word family is never truncated to its first word."""
    assert pptx_native._font_family("'DejaVu Sans', sans-serif") == "DejaVu Sans"
